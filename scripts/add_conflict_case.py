#!/usr/bin/env python3
"""Add one conflict case slide to a deck, cloned from its own What if? slide.

Katherine, 30 August 2026, on the ordering: the Concept Bank supplies
co-presence and the conflict case does the relating, so the case follows the
bank and precedes the Day 3 divider. The case "carries its own what-if".

`vt-fusion-retrofit` §3 step 5 says the slide is cloned from the matching
template slide type and never built by hand. So this copies the deck's own
What if? slide - its layout, geometry, icon, fills and type scale - and
replaces only the text. Nothing about the look is authored here, which is why
a deck whose design drifts still gets a slide that matches itself.

The two writing boxes the template already carries are used for the two things
§2c says must both be individual and written: the coordination question, which
is the occasion, and the counterfactual, which is the evidence.

Usage:  python3 add_conflict_case.py <deck.pptx> <case.json> [--apply]
"""
import argparse
import copy
import json

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt

INK = "111111"

WHATIF_MARK = "What if?"
BANK_MARK = "CONCEPT BANK"


def find(prs, needle, first_line_only=False):
    for i, s in enumerate(prs.slides):
        for sh in s.shapes:
            if not sh.has_text_frame:
                continue
            t = sh.text_frame.text
            if first_line_only:
                t = t.split("\n")[0]
            if needle in t:
                return i
    return None


R_NS = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"


def clone_slide(prs, index):
    """Copy slide `index` onto the end of the deck, relationships and all.

    A copied shape refers to its picture by relationship id, and those ids are
    per-slide. Copying the XML alone gives a slide pointing at rIds that mean
    something else on the new part, which is how a cloned slide ends up with
    the wrong image or none. So every relationship is re-created on the new
    part and every reference in the copied XML is repointed to the new id.
    """
    src = prs.slides[index]
    dest = prs.slides.add_slide(src.slide_layout)
    for shape in list(dest.shapes):
        shape._element.getparent().remove(shape._element)

    remap = {}
    for rid, rel in src.part.rels.items():
        if rel.reltype.endswith("/slideLayout"):
            continue
        if rel.is_external:
            remap[rid] = dest.part.relate_to(rel.target_ref, rel.reltype,
                                             is_external=True)
        else:
            remap[rid] = dest.part.relate_to(rel.target_part, rel.reltype)

    for el in src.shapes._spTree:
        if el.tag.endswith("}nvGrpSpPr") or el.tag.endswith("}grpSpPr"):
            continue
        new_el = copy.deepcopy(el)
        for node in [new_el] + list(new_el.iter()):
            for attr in ("embed", "id", "link", "pict"):
                key = R_NS + attr
                if key in node.attrib and node.attrib[key] in remap:
                    node.attrib[key] = remap[node.attrib[key]]
        dest.shapes._spTree.append(new_el)
    return dest


def move_to(prs, slide, position):
    sldIdLst = prs.slides._sldIdLst
    el = [s for s in sldIdLst][-1]
    sldIdLst.remove(el)
    sldIdLst.insert(position, el)


def set_para(tf, texts, size_pt=None, bold=None, color=None):
    """Replace the paragraphs of `tf`, keeping the first run's formatting.

    `texts` may be a list of strings, or of (string, bold) pairs where the
    slide needs the case and the question to read differently.
    """
    paras = list(tf.paragraphs)
    template = None
    for p in paras:
        if p.runs:
            template = p.runs[0]
            break
    for p in paras[1:]:
        p._p.getparent().remove(p._p)
    p0 = tf.paragraphs[0]
    for r in list(p0.runs):
        r._r.getparent().remove(r._r)
    for i, t in enumerate(texts):
        t, b = t if isinstance(t, tuple) else (t, bold)
        para = p0 if i == 0 else tf.add_paragraph()
        para.space_after = Pt(6)
        run = para.add_run()
        run.text = t
        if template is not None:
            run.font.size = template.font.size
            run.font.bold = template.font.bold
            run.font.name = template.font.name
            try:
                run.font.color.rgb = template.font.color.rgb
            except Exception:
                pass
        if size_pt is not None:
            run.font.size = Pt(size_pt)
        if b is not None:
            run.font.bold = b
        if color is not None and i > 0:
            run.font.color.rgb = RGBColor.from_string(color)


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("deck")
    ap.add_argument("case")
    ap.add_argument("--apply", action="store_true")
    args = ap.parse_args()

    case = json.load(open(args.case))
    prs = Presentation(args.deck)

    wi = find(prs, WHATIF_MARK, first_line_only=True)
    bank = find(prs, BANK_MARK, first_line_only=True)
    if wi is None or bank is None:
        raise SystemExit("no What if? slide (%s) or Concept Bank (%s)" % (wi, bank))
    print("cloning slide %d (What if?), inserting after slide %d (Concept Bank)"
          % (wi + 1, bank + 1))

    new = clone_slide(prs, wi)
    move_to(prs, new, bank + 1)

    boxes = sorted([sh for sh in new.shapes if sh.has_text_frame],
                   key=lambda s: s.top)
    title, question = boxes[0], boxes[1]

    def starts(prefix):
        for b in boxes:
            if b.text_frame.text.strip().startswith(prefix):
                return b
        raise SystemExit("no %r label on the template slide" % prefix)

    # Cycle 02's labels carry their grading marker on a second line, so match
    # on the opening words rather than the whole string. The marker itself is
    # deliberately not copied: it names a question this slide does not ask,
    # and a duplicated id reads as one box in two places.
    label1, label2 = starts("Your first answer"), starts("Your revised answer")
    fills = [sh for sh in new.shapes if not sh.has_text_frame
             or not sh.text_frame.text.strip()]
    box1 = min(fills, key=lambda s: abs((s.top or 0) - (label1.top + 300000)))
    box2 = min(fills, key=lambda s: abs((s.top or 0) - (label2.top + 300000)))

    set_para(title.text_frame, ["Conflict case",
                                "Critical aspects: " + case["aspects"]])

    # Anything else the template carried belongs to the slide it came from.
    # Cycle 02's What if? slide has a "Watch first" video link and a footer
    # reading "Think first. Discuss with a partner. Then write." - the first
    # points at the wrong resource here and the second contradicts section 2c,
    # which requires the what-if to be individual or the diagnosis is empty.
    # Compare the underlying XML elements, not the shape wrappers - python-pptx
    # builds a new wrapper object every time `.shapes` is iterated, so identity
    # on the wrapper is never true and this loop deleted the labels it was
    # supposed to protect.
    keep = {id(x._element) for x in (title, question, label1, label2, box1, box2)}
    for sh in list(new.shapes):
        if id(sh._element) in keep or not sh.has_text_frame:
            continue
        if sh.text_frame.text.strip():
            print("   dropping cloned text: %r"
                  % sh.text_frame.text.strip().split("\n")[0][:52])
            sh._element.getparent().remove(sh._element)

    # The template holds one short counterfactual in the question box. A
    # conflict case carries a situation and a question naming both aspects,
    # roughly three times the text, so the box grows into whatever room sits
    # between the title and the first label, and the type drops to the 14pt
    # body size rather than the template's heading size. Geometry is measured
    # per deck, because slide 1 is not the only page whose layout has drifted.
    set_para(question.text_frame,
             [(case["case"], False), (case["question"], True)], size_pt=14)
    q_top = title.top + title.height + Emu(76200)
    q_room = label1.top - q_top - Emu(60960)
    question.top = q_top
    question.height = min(Emu(1280160), q_room)

    # The template's icon sits top right, below its shorter question box. The
    # taller box would run underneath it, so the text stops at the icon.
    pics = [sh for sh in new.shapes if not sh.has_text_frame
            and sh.left is not None and sh.left > question.left]
    if pics:
        pic = min(pics, key=lambda s: s.left)
        if pic.top < question.top + question.height:
            question.width = min(question.width,
                                 pic.left - question.left - Emu(137160))

    set_para(label1.text_frame, ["Your answer"])

    # The what-if is the evidence that coordination happened - retrofit skill
    # section 2c - so it goes on the slide as a question the student answers,
    # not as a relabelled empty box. It takes the second label's place, which
    # grows upward into the gap the template leaves above it, and the second
    # writing box starts below it and keeps its original foot.
    set_para(label2.text_frame,
             [("What if", True), (case["whatif"], False)], size_pt=12,
             color=INK)
    box2_bottom = box2.top + box2.height
    gap = label2.top - (box1.top + box1.height)
    label2.top = label2.top - max(Emu(0), min(gap - Emu(91440), Emu(228600)))
    label2.height = Emu(640080)
    box2.top = label2.top + label2.height + Emu(60960)
    box2.height = box2_bottom - box2.top

    notes = new.notes_slide.notes_text_frame
    old = prs.slides[wi].notes_slide.notes_text_frame.text
    head, _, tail = old.partition("\n")
    body = tail.split("\n", 1)[1] if "\n" in tail else tail
    notes.text = head + "\nQUESTION: " + case["question"] + "\n" + body

    if args.apply:
        prs.save(args.deck)
        print("written. deck is now %d slides." % len(prs.slides.__iter__.__self__._sldIdLst))
    else:
        print("report only. re-run with --apply to write.")


if __name__ == "__main__":
    main()

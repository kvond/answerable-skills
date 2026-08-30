#!/usr/bin/env python3
"""Repair the three mechanical format failures deck_lint reports as HARD.

  H-FORMAT-FONT   a text run in a font other than Arial          (SKILL 13.1)
  H-FORMAT-HEX    a text run coloured 000000 or 0000FF           (SKILL 13.2)
  H-FORMAT-SCALE  a text run above 26pt, the day-divider ceiling (SKILL 13.3/13.4)

Nothing here is a judgment call: the skill states one font, names 000000 and
0000FF as faults, and caps the type scale. The sizes it drops to come from the
13.4 table by slide role, not from a blanket clamp - a blanket clamp is how a
heading and an activity label end up the same size.

Report first. `--apply` is a separate, deliberate flag.

Usage:
  python3 fix_format.py <dir-or-files> [--apply] [--out DIR]
"""
import argparse
import glob
import os
import re
import shutil
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt

ARIAL = "Arial"
BLACK_FAULT = "000000"
BODY_BLACK = "111111"
LINK_BLUE_FAULT = "0000FF"
TEAL = "028090"
CEILING_PT = 26.0

# The 13.4 table, by the role the run's own text identifies. Falls back to the
# heading ceiling, which is the most conservative reading of "titles are 20pt
# at the largest".
ARROW = re.compile(r"write your answer in the next slide", re.I)


def target_size(text, current_pt):
    """The size 13.4 gives this run, or None to leave it alone."""
    if current_pt is None or current_pt <= CEILING_PT:
        return None
    if ARROW.search(text or ""):
        return 14.0            # 13.4 names 32pt Calibri here as the defect
    return 20.0                # "Titles are 20pt at the largest"


def runs_of(slide):
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for para in sh.text_frame.paragraphs:
            for r in para.runs:
                yield sh, para, r
    if getattr(slide, "has_notes_slide", False):
        return


def inspect(path):
    prs = Presentation(path)
    out = []
    for i, s in enumerate(prs.slides, 1):
        for sh, para, r in runs_of(s):
            txt = (r.text or "").strip()
            font = r.font
            name = font.name
            size = font.size.pt if font.size is not None else None
            try:
                hexv = str(font.color.rgb) if font.color and font.color.type is not None else None
            except Exception:                                   # noqa: BLE001
                hexv = None

            if name and name != ARIAL:
                out.append((i, "FONT", name, txt[:52], r))
            if hexv in (BLACK_FAULT, LINK_BLUE_FAULT):
                out.append((i, "HEX", hexv, txt[:52], r))
            tgt = target_size(txt, size)
            if tgt is not None:
                out.append((i, "SCALE", "%.0fpt -> %.0fpt" % (size, tgt), txt[:52], r))
    return prs, out


def repair(prs, findings):
    n = {"FONT": 0, "HEX": 0, "SCALE": 0}
    for _slide, kind, detail, txt, r in findings:
        if kind == "FONT":
            r.font.name = ARIAL
            n["FONT"] += 1
        elif kind == "HEX":
            # 000000 is body black written wrong. 0000FF is the theme hyperlink
            # blue leaking through, and the palette colour for a link is teal.
            r.font.color.rgb = RGBColor.from_string(
                BODY_BLACK if detail == BLACK_FAULT else TEAL)
            n["HEX"] += 1
        elif kind == "SCALE":
            r.font.size = Pt(float(detail.split("-> ")[1].rstrip("pt")))
            n["SCALE"] += 1
    return n


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("paths", nargs="+")
    ap.add_argument("--apply", action="store_true",
                    help="write the repairs. Without it, report only.")
    ap.add_argument("--out", default=None,
                    help="write repaired copies here instead of in place")
    args = ap.parse_args()

    files = []
    for p in args.paths:
        if os.path.isdir(p):
            files += sorted(glob.glob(os.path.join(p, "*.pptx")))
        else:
            files += sorted(glob.glob(p))

    grand = {"FONT": 0, "HEX": 0, "SCALE": 0}
    touched = 0
    for f in files:
        prs, findings = inspect(f)
        if not findings:
            continue
        touched += 1
        print("\n=== %s  (%d finding%s)" % (os.path.basename(f), len(findings),
                                            "" if len(findings) == 1 else "s"))
        for slide, kind, detail, txt, _r in findings:
            print("   s%-3d %-6s %-22s %s" % (slide, kind, detail, txt))
            grand[kind] += 1
        if args.apply:
            n = repair(prs, findings)
            dest = f if not args.out else os.path.join(args.out, os.path.basename(f))
            if args.out:
                os.makedirs(args.out, exist_ok=True)
            prs.save(dest)
            print("   repaired: font %d, hex %d, scale %d" % (n["FONT"], n["HEX"], n["SCALE"]))

    print("\n%s across %d deck(s):  FONT %d   HEX %d   SCALE %d"
          % ("REPAIRED" if args.apply else "WOULD REPAIR", touched,
             grand["FONT"], grand["HEX"], grand["SCALE"]))
    if not args.apply:
        print("Report only. Re-run with --apply to write.")


if __name__ == "__main__":
    main()

#!/usr/bin/env python3
"""
deck_inventory.py — read-only inventory of the VT Biology cycle decks
=====================================================================
Job 1 of `01_deck_work_sequence.md`. Walks the cycle folders under
`~/deck_work/exports/`, opens every `.pptx`, and writes one CSV row per deck
plus a readable Markdown summary. It touches nothing: no deck is opened for
write, no deck is modified, no file outside the two report paths is created.

Vocabulary follows `vt-bio-skill`: **5 core questions** and **slide types**.
The word "beat" is forbidden as a structural noun (spec, Vocabulary); this
script reports where it appears rather than using it.

WHY THIS FILE EXISTS
--------------------
`audits/deck_inventory_2026-08-29.csv` was produced by a script that was never
committed. This is that pass, reconstructed from the CSV's 43 columns and the
decks themselves, and extended with the columns the 2026-08-29 spec introduces.
Detection rules were re-derived by fitting against the existing CSV; where this
script and that CSV disagree the divergence is documented in the run report,
not smoothed over.

WHAT IS DETECTED, AND HOW
-------------------------
Slide typing is inferred from title text and the §7 marker strings, in the
precedence order in SLIDE_TYPE_ORDER below. A deck whose markers drifted is
typed by what it says, not by what it was meant to be.

  Teacher reference        slide 1, or "TEACHER REFERENCE" in the slide text
  Teacher Prep             "Teacher Prep" on the slide (the Additional
                           Resources slide, NOT the one-slide teacher note)
  Day divider              title matching "Day N of M"
  Bellringer               title "Bellringer"
  Response slide           both writing boxes: "Your first answer" and
                           "Your revised answer"; "(What if)" when the slide
                           also carries the "What if?" marker
  Optional challenge       title "Optional challenge"
  Relates to me            title "Relates to me"
  REMOVED: ...             the three slide types retired 2026-08-12
  Closing checklist        "Think -> Write -> Submit"
  Links slide              "Activity and resource links"
  Image credits            "Image credits"
  Concept Bank             "CONCEPT BANK"
  Teacher note             "TEACHER NOTE" on the slide itself (spec §Teacher
                           note slide) — distinct from Teacher Prep, and
                           distinct from anything living in speaker notes
  What if                  "What if?" marker with no writing boxes
  Pattern Break            "Pattern break"
  Build a Rule             "Finish this sentence as a rule"
  3-Tier Question          all three tier labels on one slide
  Continuity question      "Continuation question" / "Continuity question" /
                           "no answer to look up"
  Contrast Set             "CASE A" and "CASE B"
  Stock-and-flow model     an interactive tank model ("... Tank Model")
  Activity                 title beginning "Activity"
  Critical Aspect question the FIRST slide carrying a given "Critical aspect:"
                           label; later labelled slides that match no other
                           type are "Explanation"
  Content/Explanation      everything else

NEW COLUMNS (2026-08-29 spec)
-----------------------------
  teacher_note_slide          present/absent. Slide text and speaker notes are
                              read separately so the one-slide teacher note is
                              never confused with the Teacher Prep slide or
                              with design intent buried in speaker notes.
  teacher_note_declares       which of the seven required declarations the
                              teacher note makes (spec §Teacher note slide)
  visibility_rung_declared    yes/no, plus the declared rung text
  simultaneity_declared       synchronic / diachronic / both / none
  what_if_present             yes/no
  what_if_individual_written  does the slide or its notes require an individual
                              written answer (spec §6c: group work is a
                              visibility reducer and no evidence of fusion)
  concept_bank_relates        does the bank ask the student to state how two
                              terms connect, or does it only list them
  forbidden_vocab_beat        does the forbidden structural noun appear
  ngss_codes                  NGSS performance expectations found in the deck

USAGE
-----
  python3 deck_inventory.py [exports_dir] [-o OUT_STEM]

  exports_dir  default ~/deck_work/exports
  -o           output stem; default audits/deck_inventory_<today>
               writes <stem>.csv and <stem>.md
  --compare P  cross-check the run against an earlier inventory CSV and print
               a per-deck disagreement report on the shared columns

Dependencies: python-pptx only. Read-only. Never writes a .pptx.
"""
from __future__ import print_function

import argparse
import csv
import datetime
import os
import re
import sys
from collections import OrderedDict

from pptx import Presentation


# ---------------------------------------------------------------------------
# Format tokens (SKILL §13) — must agree with deck_lint.py
# ---------------------------------------------------------------------------

SLIDE_W, SLIDE_H = 9144000, 6858000          # 4:3

PALETTE = {
    "028090", "111111", "666666", "C0392B", "1E8449",
    "F2F6F9", "CCCCCC", "F5F5F5", "EAF6F7", "FFFFFF",
}
# NOTE: EFDF85 ("Working On It") is deliberately NOT in the palette here.
# The inventory reports it as off-token because that is what the reference
# inventory did; deck_lint.py accepts it. The disagreement is real and is
# reported rather than resolved by this script — it is a spec question
# (SKILL §13.2 says body black, the reference build ships EFDF85).

ALLOWED_FONT = "Arial"

TIER_LABELS = ("Getting Started", "Working On It", "Mastery")

REQUIRED_TYPES = OrderedDict([
    ("Teacher Prep",    lambda t: "Teacher Prep" in t),
    ("Concept Bank",    lambda t: "Concept Bank" in t),
    ("Links slide",     lambda t: "Links slide" in t),
    ("Response slides", lambda t: any(x.startswith("Response slide") for x in t)),
    ("Day divider",     lambda t: "Day divider" in t),
])

REMOVED_TYPES = (
    "REMOVED: Then and Now",
    "REMOVED: Turn answer into draft",
    "REMOVED: teacher slide-type list",
)

CORE_FIVE = ("Critical Aspect question", "Contrast Set", "Build a Rule",
             "Pattern Break", "3-Tier Question")


# ---------------------------------------------------------------------------
# Advisory 1 — does the move-1 question itself carry a difference?
# ---------------------------------------------------------------------------
# Three device families, per `08_nine_moves_attribution.md` "Move 1 — the flag".
# These regexes were fitted against the labelled move-1 questions in the
# reference inventory: 47/49 agreement, the two residuals documented in the
# run report. A keyword heuristic, deliberately: the verbatim question text is
# carried in the CSV so the call stays with the teacher.

DEVICE_COMPARISON = [
    r"\bor\b", r"\bbut\b", r"\bboth\b", r"\bsame\b", r"\bdifferent(ly)?\b",
    r"\bthan\b", r"\bmore\b", r"\bless\b", r"\binstead of\b", r"\bunlike\b",
    r"\bwhile\b", r"\bversus\b", r"\bvs\.?\b", r"\bwhy can.?.?t\b",
    r"\bchanging\b", r"\bstill\b",
]
DEVICE_NAMED_ALTERNATIVES = [r"[A-Za-z]+ or [A-Za-z]+"]
DEVICE_CHANGE_CONDITION = [
    r"\bchange[ds]?\b", r"\bwhat if\b",
    r"\bwould\b.{0,60}\bif\b", r"\bif\b.{0,60}\bwould\b",
]
DEVICES = OrderedDict([
    ("change condition", DEVICE_CHANGE_CONDITION),
    ("comparison", DEVICE_COMPARISON),
    ("named alternatives", DEVICE_NAMED_ALTERNATIVES),
])


# ---------------------------------------------------------------------------
# Advisory 2 — coordination structures (spec §"The nine moves")
# ---------------------------------------------------------------------------
# Three objects of different kinds (08 §"correctly distinguished"):
#   stock-and-flow  a REPRESENTATION — rate against amount
#   compensatory pair  a CASE SET — two cases, same outcome, inverted values
#   conflict case   a SINGLE CASE with opposing pulls
#
# Scanned over student-facing slides only. The reference inventory scanned
# every slide including the teacher reference, the Teacher Prep buy-ahead list
# and the Relates-to-me boilerplate, which produced hits that cannot be a
# coordination structure (a lab "balance", a slide that merely names the two
# critical aspects). Those are excluded here.

ADV2_TWO_CASE = re.compile(
    r"CASE A.{0,4000}CASE B|\btwo (cases|examples|organisms|species|populations)\b"
    r"|\bboth\b", re.I | re.S)

ADV2_EXCLUDED_TYPES = {
    "Teacher reference", "Teacher Prep", "Teacher note", "Links slide",
    "Image credits", "Closing checklist", "Relates to me", "Concept Bank",
    "Day divider", "REMOVED: Then and Now", "REMOVED: Turn answer into draft",
    "REMOVED: teacher slide-type list",
}

# Stock-and-flow is a REPRESENTATION, and in this arc it is always a slide: the
# interactive model. It is taken from the slide type rather than from keywords,
# because rate-against-amount wording turns up in bellringers and explanations
# that are not the model.
ADV2_FROM_SLIDE_TYPE = {"Stock-and-flow model": "stock-and-flow model"}

ADV2_PATTERNS = OrderedDict([
    # A compensatory pair is a CASE SET: two cases, same outcome, inverted
    # values on two aspects. The keyword alone is not enough — "trade-off"
    # also names a cost-benefit weighing, which is not a compensatory pair —
    # so a two-case structure must be present on the same slide as well
    # (ADV2_TWO_CASE below).
    ("compensatory pair", [
        r"\bcompensat",
        r"\btrade[- ]offs?\b",
        r"\bmakes up for\b",
        r"\boffsets?\b",
        r"\bboth\b.{0,120}\bsame (outcome|result|end|number|size|level)\b",
        r"\bone .{0,40}\bhigh\b.{0,60}\bthe other .{0,40}\blow\b",
        r"\bfew\b.{0,80}\bmany\b.{0,80}\bboth\b",
    ]),
    ("conflict case", [
        r"\bpulls? (in )?(two|opposite|both) (ways|directions)\b",
        r"\bopposing\b",
        r"\btension between\b",
        r"\bcontradict",
        r"\bparadox",
        r"\bat the same time\b.{0,120}\bbut\b",
        r"\bpoints? one way\b.{0,120}\bthe other\b",
        r"\bcannot be (resolved|settled) (from|by) one\b",
    ]),
])


# ---------------------------------------------------------------------------
# New-column patterns (2026-08-29 spec)
# ---------------------------------------------------------------------------

RE_TEACHER_NOTE_SLIDE = re.compile(r"\bTEACHER NOTE\b", re.I)
# Scoped deliberately. "Teacher prep" also appears as a cross-reference on the
# links slide ("Teacher prep (Slide 2) — ..."), and "TEACHER REFERENCE" appears
# inside the retired teacher slide-type list. Matching either loosely swallows
# the wrong slide.
RE_TEACHER_PREP = re.compile(
    r"Additional Resources\s*[\u2014\u2013-]\s*Teacher Prep", re.I)
RE_STOCK_FLOW = re.compile(r"tank model", re.I)
RE_EXPLORE_MODEL = re.compile(r"^Explore\s*[\u2014\u2013-]\s*the\b.*\bModel\b", re.I)

# The seven required declarations, spec §"Teacher note slide".
TN_DECLARATIONS = OrderedDict([
    ("critical aspect", [r"critical aspect"]),
    ("held invariant", [r"held (invariant|constant)", r"\binvariant\b",
                        r"one dimension", r"differ in one"]),
    ("what breaks if substituted", [r"substitut", r"what breaks",
                                    r"if an example (is|were) (swapped|replaced|changed)"]),
    ("position in sequence", [r"position in the sequence", r"what this assumes",
                              r"depends on it", r"\bassumes\b.{0,60}\bdepends\b"]),
    ("slide-type map", [r"slide[- ]type map", r"slide types (that )?appear",
                        r"left out", r"conditional .{0,30}(left out|omitted)"]),
    ("simultaneity", [r"synchronic", r"diachronic", r"simultaneity"]),
    ("visibility rung", [r"visibilit", r"\brung\b", r"unattributed"]),
])

RE_VISIBILITY = re.compile(
    r"(visibility\s*[:—-]\s*[^\n]{0,160}"
    r"|\brung\s*\d[^\n]{0,120}"
    r"|written and private[^\n]{0,120}"
    r"|read[- ]aloud[^\n]{0,120}unattributed[^\n]{0,60})", re.I)
RE_SYNCHRONIC = re.compile(r"synchronic", re.I)
RE_DIACHRONIC = re.compile(r"diachronic", re.I)

RE_WHATIF = re.compile(r"what if\?", re.I)
RE_INDIVIDUAL_WRITTEN = re.compile(
    r"(individually|on your own|by yourself|your own words|write your (own|first) answer)", re.I)
RE_GROUP = re.compile(r"(with your group|as a group|in groups|group answer)", re.I)

RE_BANK_RELATES = re.compile(
    r"(how (any )?two[^\n]{0,60}(connect|relate|fit together)"
    r"|pick two[^\n]{0,80}(connect|relate)"
    r"|connect (two|any two)[^\n]{0,60}terms"
    r"|state how[^\n]{0,60}(connect|relate))", re.I)

RE_BEAT = re.compile(r"\bbeats?\b|\bbeating\b", re.I)
RE_NGSS = re.compile(r"\b(?:HS|MS)-[A-Z]{2,4}\d?-\d+\b")

RE_URL = re.compile(r"https?://[^\s<>\"'\)\]]+")
RE_CA_LABEL = re.compile(r"Critical aspect\s*:\s*(.+)", re.I)
RE_DAY_DIVIDER = re.compile(r"^\s*Day\s+\d+\s+of\s+\d+", re.I)
RE_CYCLE_KEY = re.compile(r"Cycle\s*0*(\d{1,2})\s*([a-d])?", re.I)


# ---------------------------------------------------------------------------
# pptx reading helpers (idiom shared with deck_lint.py / build_concept_banks.py)
# ---------------------------------------------------------------------------

def walk_shapes(shapes):
    """Yield every shape, descending into groups."""
    for sh in shapes:
        yield sh
        try:
            if sh.shape_type == 6:          # MSO_SHAPE_TYPE.GROUP
                for inner in walk_shapes(sh.shapes):
                    yield inner
        except Exception:
            pass


def shape_text(slide):
    """Slide-surface text only. Speaker notes are read separately, on purpose:
    the one-slide teacher note must never be confused with speaker notes."""
    out = []
    for sh in walk_shapes(slide.shapes):
        try:
            if sh.has_text_frame and sh.text_frame.text.strip():
                out.append(sh.text_frame.text)
        except Exception:
            pass
    return "\n".join(out)


def notes_text(slide):
    try:
        if slide.has_notes_slide:
            t = slide.notes_slide.notes_text_frame.text
            return t if t else ""
    except Exception:
        pass
    return ""


def first_line(text):
    for line in text.split("\n"):
        if line.strip():
            return line.strip()
    return ""


def run_hex(run):
    try:
        c = run.font.color
        if c and c.type is not None and c.rgb is not None:
            return str(c.rgb).upper()
    except Exception:
        pass
    return None


def collect_fonts_colors(prs):
    fonts, colors = set(), set()
    for slide in prs.slides:
        for sh in walk_shapes(slide.shapes):
            try:
                if sh.has_text_frame:
                    for para in sh.text_frame.paragraphs:
                        for r in para.runs:
                            if r.font.name:
                                fonts.add(r.font.name)
                            h = run_hex(r)
                            if h:
                                colors.add(h)
            except Exception:
                pass
            try:
                if sh.fill.type == 1:
                    colors.add(str(sh.fill.fore_color.rgb).upper())
            except Exception:
                pass
            try:
                if sh.line.color is not None and sh.line.color.type is not None:
                    colors.add(str(sh.line.color.rgb).upper())
            except Exception:
                pass
    return fonts, colors


def slide_urls(slide, text):
    urls = []
    def add(u):
        u = u.strip()
        if u and u not in urls:
            urls.append(u)
    for sh in walk_shapes(slide.shapes):
        try:
            if sh.has_text_frame:
                for para in sh.text_frame.paragraphs:
                    for r in para.runs:
                        try:
                            if r.hyperlink is not None and r.hyperlink.address:
                                add(r.hyperlink.address)
                        except Exception:
                            pass
        except Exception:
            pass
        try:
            if sh.click_action is not None and sh.click_action.hyperlink.address:
                add(sh.click_action.hyperlink.address)
        except Exception:
            pass
    for m in RE_URL.finditer(text):
        add(m.group(0).rstrip(".,;"))
    return urls


# ---------------------------------------------------------------------------
# Slide typing
# ---------------------------------------------------------------------------

def type_slide(idx, text, notes, seen_aspects):
    """Return (slide_type, critical_aspect_or_None).

    Precedence matters and is the order below. It was derived from the
    reference inventory: e.g. the Optional challenge slide that links a tank
    model is an Optional challenge, not a Stock-and-flow model; an Activity
    slide that links one IS typed by the model.
    """
    low = text.lower()
    title = first_line(text)
    ca = None
    m = RE_CA_LABEL.search(text)
    if m:
        ca = m.group(1).strip()
        ca = re.split(r"\s{2,}|\s+/\s+", ca)[0].strip()

    def has(s):
        return s.lower() in low

    # The one-slide teacher note is tested BEFORE Teacher Prep, and both before
    # Teacher reference: the Teacher Prep slide also opens with the words
    # "TEACHER REFERENCE", so testing the generic banner first swallows it.
    if RE_TEACHER_NOTE_SLIDE.search(text) and not RE_TEACHER_PREP.search(text):
        return "Teacher note", ca
    if RE_TEACHER_PREP.search(text) or title.strip().lower() == "teacher prep":
        return "Teacher Prep", ca
    if idx == 1 or "teacher reference" in first_line(text).lower():
        return "Teacher reference", ca
    if RE_DAY_DIVIDER.match(title):
        return "Day divider", ca
    if title.strip().lower() == "bellringer":
        return "Bellringer", ca
    if has("Your first answer") and has("Your revised answer"):
        return ("Response slide (What if)" if RE_WHATIF.search(text)
                else "Response slide"), ca
    if title.lower().startswith("optional challenge"):
        return "Optional challenge", ca
    if title.lower().startswith("relates to me"):
        return "Relates to me", ca
    if title.lower().startswith("then and now"):
        return "REMOVED: Then and Now", ca
    if has("Turn your answer into a draft"):
        return "REMOVED: Turn answer into draft", ca
    if "teacher navigation" in low and ("legend" in low or "what each slide is doing" in low):
        return "REMOVED: teacher slide-type list", ca
    if has("Think") and ("write" in low and "submit" in low and "→" in text):
        return "Closing checklist", ca
    if has("Activity and resource links"):
        return "Links slide", ca
    if title.lower().startswith("image credits"):
        return "Image credits", ca
    if "CONCEPT BANK" in text:
        return "Concept Bank", ca
    if RE_WHATIF.search(text):
        return "What if", ca
    if has("Pattern break"):
        return "Pattern Break", ca
    if has("Finish this sentence as a rule"):
        return "Build a Rule", ca
    if all(lbl in text for lbl in TIER_LABELS):
        return "3-Tier Question", ca
    if (has("Continuation question") or has("Continuity question")
            or has("no answer to look up")):
        return "Continuity question", ca
    if "CASE A" in text and "CASE B" in text:
        return "Contrast Set", ca
    if RE_STOCK_FLOW.search(text) or RE_EXPLORE_MODEL.match(title):
        return "Stock-and-flow model", ca
    if (title.lower().startswith("activity")
            or re.match(r"^Day\s+\d+\s+activity\b", title, re.I)
            or re.search(r"\bLab$", title)
            or has("Activity (Critical Aspect")):
        return "Activity", ca
    if ca:
        return ("Critical Aspect question" if ca not in seen_aspects
                else "Explanation"), ca
    return "Content/Explanation", ca


def move1_devices(question_text):
    found = []
    for name, pats in DEVICES.items():
        for p in pats:
            if re.search(p, question_text, re.I):
                found.append(name)
                break
    return sorted(found)


def ca_question_text(text):
    """The question on a Critical Aspect question slide: the longest line that
    is not the aspect label and not the standing 'Think first' instruction."""
    best = ""
    for line in text.split("\n"):
        s = line.strip()
        if not s or RE_CA_LABEL.match(s):
            continue
        if s.lower().startswith("think first"):
            continue
        if s.startswith("[[") or s.startswith("Critical aspect"):
            continue
        if len(s) > len(best):
            best = s
    return best


# ---------------------------------------------------------------------------
# Per-deck inventory
# ---------------------------------------------------------------------------

def inventory_deck(path, folder):
    base = os.path.basename(path)
    prs = Presentation(path)
    slides = list(prs.slides)

    m = RE_CYCLE_KEY.search(base)
    cycle = (m.group(1) + (m.group(2).lower() if m.group(2) else "")) if m else ""
    deck_set = "with-Concept-Bank" if "with Concept Bank" in base else "LIVE export"

    texts, notes_all, types, cas = [], [], [], []
    seen_aspects = []
    for i, s in enumerate(slides, start=1):
        t = shape_text(s)
        n = notes_text(s)
        texts.append(t)
        notes_all.append(n)
        st, ca = type_slide(i, t, n, seen_aspects)
        types.append(st)
        cas.append(ca)
        if ca and ca not in seen_aspects:
            seen_aspects.append(ca)

    all_text = "\n".join(texts)
    all_notes = "\n".join(notes_all)
    everything = all_text + "\n" + all_notes

    # -- critical aspects and the 5 core questions -------------------------
    aspects = seen_aspects
    by_aspect = OrderedDict()
    for a in aspects:
        got = []
        for i, st in enumerate(types):
            if cas[i] == a and st in CORE_FIVE and st not in got:
                got.append(st)
        by_aspect[a] = [q for q in CORE_FIVE if q in got]
    core_desc = " ;; ".join(
        "%s: %d/5 [%s]" % (a, len(v), ", ".join(v)) for a, v in by_aspect.items())
    core_complete = bool(aspects) and all(len(v) == 5 for v in by_aspect.values())

    # -- counts ------------------------------------------------------------
    n_type = lambda name: sum(1 for t in types if t == name)
    ca_q_idx = [i for i, t in enumerate(types) if t == "Critical Aspect question"]
    resp_idx = [i for i, t in enumerate(types) if t.startswith("Response slide")]
    resp_with_prompt = [i for i in resp_idx if "REVISION PROMPT" in notes_all[i]]
    resp_missing = [str(i + 1) for i in resp_idx if i not in resp_with_prompt]

    # -- links -------------------------------------------------------------
    per_slide_urls = [slide_urls(s, texts[i]) for i, s in enumerate(slides)]
    links_idx = [i for i, t in enumerate(types) if t == "Links slide"]
    links_slide_urls = len({u for i in links_idx for u in per_slide_urls[i]}) if links_idx else 0
    all_urls = []
    for us in per_slide_urls:
        for u in us:
            if u not in all_urls:
                all_urls.append(u)          # deck order, not alphabetical
    copy_urls = [u for u in all_urls if "/copy" in u]
    non_copy = [u for u in all_urls if u not in copy_urls]
    gslides_non_copy = [u for u in non_copy if "docs.google.com/presentation" in u]

    # -- format tokens -----------------------------------------------------
    fonts, colors = collect_fonts_colors(prs)
    fonts_off = sorted(f for f in fonts if f != ALLOWED_FONT)
    colors_off = sorted(c for c in colors if c not in PALETTE)
    dims_ok = (prs.slide_width == SLIDE_W and prs.slide_height == SLIDE_H)

    # -- required / removed ------------------------------------------------
    tset = types
    missing = [name for name, test in REQUIRED_TYPES.items() if not test(tset)]
    removed = [r for r in REMOVED_TYPES if r in tset]

    # -- advisory 1 --------------------------------------------------------
    verdicts, verbatims = [], []
    for i in ca_q_idx:
        q = ca_question_text(texts[i])
        dev = move1_devices(q)
        verdicts.append("S%d: %s" % (i + 1, ", ".join(dev) if dev else "NO CONTRAST DEVICE"))
        verbatims.append('S%d: "%s"' % (i + 1, q))
    adv1_pass = bool(ca_q_idx) and all("NO CONTRAST DEVICE" not in v for v in verdicts)

    # -- advisory 2 --------------------------------------------------------
    adv2 = OrderedDict()
    for i, st in enumerate(types):
        if st in ADV2_EXCLUDED_TYPES:
            continue
        if st in ADV2_FROM_SLIDE_TYPE:
            adv2.setdefault(ADV2_FROM_SLIDE_TYPE[st], i + 1)
        hay = texts[i] + "\n" + notes_all[i]
        for name, pats in ADV2_PATTERNS.items():
            if name in adv2:
                continue
            if name == "compensatory pair" and not ADV2_TWO_CASE.search(hay):
                continue
            for p in pats:
                if re.search(p, hay, re.I):
                    adv2[name] = i + 1
                    break
    order = ["stock-and-flow model"] + [n for n in ADV2_PATTERNS
                                       if n != "stock-and-flow model"]
    adv2_names = [n for n in order if n in adv2]
    adv2_structure = "; ".join(adv2_names) if adv2_names else "none"
    adv2_detail = "; ".join("%s (slide %d)" % (n, adv2[n]) for n in adv2_names)

    # -- NEW: teacher note slide and its seven declarations ----------------
    tn_idx = [i for i, t in enumerate(types) if t == "Teacher note"]
    teacher_note_slide = "present" if tn_idx else "absent"
    tn_declares = []
    if tn_idx:
        hay = "\n".join(texts[i] + "\n" + notes_all[i] for i in tn_idx)
        for name, pats in TN_DECLARATIONS.items():
            if any(re.search(p, hay, re.I) for p in pats):
                tn_declares.append(name)

    # -- NEW: visibility rung ---------------------------------------------
    vm = RE_VISIBILITY.search(everything)
    visibility_declared = "yes" if vm else "no"
    visibility_text = re.sub(r"\s+", " ", vm.group(0)).strip()[:200] if vm else ""

    # -- NEW: simultaneity -------------------------------------------------
    syn, dia = bool(RE_SYNCHRONIC.search(everything)), bool(RE_DIACHRONIC.search(everything))
    simultaneity = ("both" if syn and dia else "synchronic" if syn
                    else "diachronic" if dia else "none")

    # -- NEW: what-if ------------------------------------------------------
    wi_idx = [i for i, t in enumerate(types)
              if t in ("What if", "Response slide (What if)")]
    what_if_present = "yes" if wi_idx else "no"
    if wi_idx:
        hay = "\n".join(texts[i] + "\n" + notes_all[i] for i in wi_idx)
        boxes = any("Your first answer" in texts[i] and "Your revised answer" in texts[i]
                    for i in wi_idx)
        if boxes or (RE_INDIVIDUAL_WRITTEN.search(hay) and not RE_GROUP.search(hay)):
            what_if_individual = "yes"
        elif RE_GROUP.search(hay):
            what_if_individual = "no (group)"
        else:
            what_if_individual = "unclear"
    else:
        what_if_individual = "n/a"

    # -- NEW: Concept Bank, relating vs listing ----------------------------
    cb_idx = [i for i, t in enumerate(types) if t == "Concept Bank"]
    if not cb_idx:
        cb_relates = "absent"
    else:
        hay = "\n".join(texts[i] + "\n" + notes_all[i] for i in cb_idx)
        cb_relates = "relates" if RE_BANK_RELATES.search(hay) else "lists only"

    # -- NEW: forbidden vocabulary ----------------------------------------
    beat_hits = []
    for i in range(len(slides)):
        if RE_BEAT.search(texts[i]) or RE_BEAT.search(notes_all[i]):
            beat_hits.append(str(i + 1))
    beat = ("slides " + ", ".join(beat_hits)) if beat_hits else "clean"

    # -- NEW: NGSS ---------------------------------------------------------
    ngss = sorted(set(RE_NGSS.findall(everything)))

    st = os.stat(path)
    return OrderedDict([
        ("file", base),
        ("folder", folder),
        ("set", deck_set),
        ("cycle", cycle),
        ("slides", len(slides)),
        ("mtime", datetime.datetime.fromtimestamp(st.st_mtime).strftime("%Y-%m-%d %H:%M:%S")),
        ("bytes", st.st_size),
        ("critical_aspects", " | ".join(aspects)),
        ("n_aspects", len(aspects)),
        ("core_questions_by_aspect", core_desc),
        ("core_q_complete_all_aspects", core_complete),
        ("ca_question_count", len(ca_q_idx)),
        ("continuity_question", n_type("Continuity question")),
        ("stock_and_flow", n_type("Stock-and-flow model")),
        ("concept_bank", n_type("Concept Bank")),
        ("teacher_prep", n_type("Teacher Prep")),
        ("day_dividers", n_type("Day divider")),
        ("response_slides", len(resp_idx)),
        ("response_slides_with_prompt", len(resp_with_prompt)),
        ("all_response_notes_ok", bool(resp_idx) and len(resp_idx) == len(resp_with_prompt)),
        ("response_missing_prompt_slides", "; ".join(resp_missing)),
        ("links_slide", len(links_idx)),
        ("links_slide_urls", links_slide_urls),
        ("links_total_unique", len(all_urls)),
        ("links_copy", len(copy_urls)),
        ("links_non_copy", len(non_copy)),
        ("links_gslides_non_copy", len(gslides_non_copy)),
        ("links_non_copy_sample", " ; ".join(non_copy[:4])),
        ("dims", "%dx%d" % (prs.slide_width, prs.slide_height)),
        ("dims_ok", dims_ok),
        ("fonts_ok", not fonts_off),
        ("fonts_found", "; ".join(sorted(fonts))),
        ("fonts_offending", "; ".join(fonts_off)),
        ("colors_ok", not colors_off),
        ("colors_offending", "; ".join(colors_off)),
        ("format_ok", dims_ok and not fonts_off and not colors_off),
        ("removed_slide_types", "; ".join(removed)),
        ("missing_required", "; ".join(missing)),
        ("advisory1_pass", adv1_pass),
        ("advisory1_verdict", " ;; ".join(verdicts)),
        ("ca_question_text_verbatim", " ;; ".join(verbatims)),
        ("advisory2_structure", adv2_structure),
        ("advisory2_detail", adv2_detail),
        ("slide_types_ordered", " > ".join(types)),
        # ---- new columns, 2026-08-29 spec --------------------------------
        ("teacher_note_slide", teacher_note_slide),
        ("teacher_note_declares", "; ".join(tn_declares)),
        ("teacher_note_declares_count", "%d/7" % len(tn_declares)),
        ("visibility_rung_declared", visibility_declared),
        ("visibility_rung_text", visibility_text),
        ("simultaneity_declared", simultaneity),
        ("what_if_present", what_if_present),
        ("what_if_individual_written", what_if_individual),
        ("concept_bank_relates", cb_relates),
        ("forbidden_vocab_beat", beat),
        ("ngss_codes", "; ".join(ngss)),
    ])


# ---------------------------------------------------------------------------
# Walk, report
# ---------------------------------------------------------------------------

def _is_deck(name):
    if not name.lower().endswith(".pptx"):
        return False
    # backups and Office lock files are not decks
    return not (name.startswith("_") or name.startswith("~$"))


def walk_exports(root):
    """Decks under root.

    The exports tree is root/<cycle folder>/<deck>.pptx, but a deck can also sit
    loose at the root — a one-off rebuild, a scratch build, a single file handed
    in for checking. Those used to be invisible: the walk only descended into
    subfolders, so `deck_inventory.py ~/deck_work/12b_rebuild` inventoried the
    _src/ copies and silently skipped the actual deck above them. Root-level
    decks are now included, folder-labelled by the directory's own name.
    """
    out = []
    if os.path.isfile(root) and _is_deck(os.path.basename(root)):
        return [(root, os.path.basename(os.path.dirname(root)) or ".")]

    for name in sorted(os.listdir(root)):           # loose decks at the root
        if _is_deck(name) and os.path.isfile(os.path.join(root, name)):
            out.append((os.path.join(root, name), os.path.basename(root.rstrip("/")) or "."))

    for folder in sorted(os.listdir(root)):         # the normal exports shape
        fdir = os.path.join(root, folder)
        if not os.path.isdir(fdir):
            continue
        for name in sorted(os.listdir(fdir)):
            if _is_deck(name):
                out.append((os.path.join(fdir, name), folder))
    return out


def cycle_sort_key(c):
    m = re.match(r"(\d+)([a-d]?)", c or "")
    return (int(m.group(1)), m.group(2)) if m else (999, "")


def deck_label(row):
    n = re.sub(r"^▶ LIVE — ", "", row["file"])
    n = re.sub(r"\.pptx$", "", n)
    n = re.sub(r"^Cycle\s*", "", n)
    n = re.sub(r"\s*\(VT deck[^)]*\)|\s*\(VT\)", "", n)
    n = re.sub(r"\s*— with Concept Bank.*$", "", n)
    return n.strip()


def write_markdown(rows, path, root, csv_path):
    live = [r for r in rows if r["set"] == "LIVE export"]
    cb = [r for r in rows if r["set"] == "with-Concept-Bank"]
    L = []
    A = L.append
    A("# VT Biology cycle deck inventory")
    A("")
    A("Read-only pass over `%s`. Generated %s."
      % (root, datetime.datetime.now().strftime("%Y-%m-%d %H:%M")))
    A("Vocabulary follows `vt-bio-skill`: **5 core questions**, **slide types**.")
    A("Full data: `%s` (%d columns, %d rows)."
      % (os.path.basename(csv_path), len(rows[0]) if rows else 0, len(rows)))
    A("")
    A("This is an inventory. It reports; it does not judge. An absence in any")
    A("column below is a measurement, not a verdict.")
    A("")
    A("| Set | Decks |")
    A("|---|---|")
    A("| `LIVE export` | %d |" % len(live))
    A("| `with-Concept-Bank` | %d |" % len(cb))
    A("")
    A("---")
    A("")
    A("## Per deck — the LIVE set")
    A("")
    A("| Deck | Slides | Aspects | 5 core qs | Cont. | S&F | Bank | Prep | Note | Resp (notes ok) | Links (/copy) | Format | Adv 1 | Adv 2 |")
    A("|---|---|---|---|---|---|---|---|---|---|---|---|---|---|")
    for r in sorted(live, key=lambda r: cycle_sort_key(r["cycle"])):
        A("| %s | %s | %s | %s | %s | %s | %s | %s | %s | %s (%s) | %s (%s) | %s | %s | %s |" % (
            deck_label(r)[:44], r["slides"], r["n_aspects"],
            "yes" if r["core_q_complete_all_aspects"] else "NO",
            r["continuity_question"] or "—", r["stock_and_flow"] or "—",
            r["concept_bank"] or "—", r["teacher_prep"] or "—",
            "yes" if r["teacher_note_slide"] == "present" else "—",
            r["response_slides"], "yes" if r["all_response_notes_ok"] else "NO",
            r["links_total_unique"], r["links_copy"],
            "yes" if r["format_ok"] else "NO",
            "yes" if r["advisory1_pass"] else "NO",
            r["advisory2_structure"]))
    A("")
    A("`Cont.` = Continuity question · `S&F` = Stock-and-flow model · `Bank` = Concept Bank ·")
    A("`Prep` = Teacher Prep slide · `Note` = the one-slide teacher note · `Resp` = response")
    A("slides, with whether every one carries the revision prompt in its speaker notes.")
    A("")
    A("---")
    A("")
    A("## The 2026-08-29 columns")
    A("")
    A("| Deck | Teacher note | Declares | Visibility rung | Simultaneity | What-if | individual+written | Concept Bank | \"beat\" | NGSS |")
    A("|---|---|---|---|---|---|---|---|---|---|")
    for r in sorted(live, key=lambda r: cycle_sort_key(r["cycle"])):
        A("| %s | %s | %s | %s | %s | %s | %s | %s | %s | %s |" % (
            deck_label(r)[:40], r["teacher_note_slide"],
            r["teacher_note_declares_count"], r["visibility_rung_declared"],
            r["simultaneity_declared"], r["what_if_present"],
            r["what_if_individual_written"], r["concept_bank_relates"],
            r["forbidden_vocab_beat"], r["ngss_codes"] or "—"))
    A("")
    # summaries
    A("### Missing a required slide type")
    A("")
    A("| Slide type | Decks missing it |")
    A("|---|---|")
    for name in REQUIRED_TYPES:
        miss = [deck_label(r) for r in live if name in r["missing_required"].split("; ")]
        A("| %s | %d of %d |" % (name, len(miss), len(live)))
    A("")
    A("### Teacher note slide (spec §Teacher note slide)")
    A("")
    have = [r for r in live if r["teacher_note_slide"] == "present"]
    A("%d of %d LIVE decks carry a one-slide teacher note. The seven declarations "
      "are checked only where one exists." % (len(have), len(live)))
    A("")
    A("### Visibility and simultaneity")
    A("")
    A("Declared visibility rung: %d of %d. Declared simultaneity: %d of %d."
      % (sum(1 for r in live if r["visibility_rung_declared"] == "yes"), len(live),
         sum(1 for r in live if r["simultaneity_declared"] != "none"), len(live)))
    A("")
    A("### Concept Bank — listing versus relating")
    A("")
    A("Co-presence is the precondition, not the achievement. A bank that only lists")
    A("terms has left the relating to chance.")
    A("")
    for state in ("relates", "lists only", "absent"):
        n = [deck_label(r) for r in rows if r["concept_bank_relates"] == state]
        A("- **%s**: %d deck(s)" % (state, len(n)))
    A("")
    A("### Forbidden vocabulary")
    A("")
    bad = [(deck_label(r), r["forbidden_vocab_beat"]) for r in rows
           if r["forbidden_vocab_beat"] != "clean"]
    if bad:
        A("The word is forbidden everywhere (spec, Vocabulary). Found in %d file(s):" % len(bad))
        A("")
        for n, where in bad:
            A("- %s — %s" % (n, where))
    else:
        A("Clean across every deck.")
    A("")
    A("### NGSS codes found")
    A("")
    A("| Deck | Codes |")
    A("|---|---|")
    for r in sorted(live, key=lambda r: cycle_sort_key(r["cycle"])):
        A("| %s | %s |" % (deck_label(r)[:44], r["ngss_codes"] or "*none in deck*"))
    A("")
    A("### Advisory 1 — move-1 questions with no contrast device")
    A("")
    A("A question that names the aspect and waits for the Contrast Set is")
    A("structurally \"how are polar bears adapted?\". Verbatim text below — these are")
    A("borderline calls and they are the teacher's to make, not this script's.")
    A("")
    for r in sorted(live, key=lambda r: cycle_sort_key(r["cycle"])):
        verds = [v.strip() for v in r["advisory1_verdict"].split(";;") if v.strip()]
        verbs = [v.strip() for v in r["ca_question_text_verbatim"].split(";;") if v.strip()]
        for v, q in zip(verds, verbs):
            if "NO CONTRAST DEVICE" in v:
                A("**%s, %s**" % (deck_label(r), v.split(":")[0]))
                A("")
                A("> %s" % q.split(":", 1)[1].strip().strip('"'))
                A("")
    A("### Advisory 2 — coordination structures")
    A("")
    A("Recorded, not judged. No opinion is offered on whether a cycle ought to")
    A("carry one. An absence is not automatically a hole.")
    A("")
    A("| Structure found | Decks |")
    A("|---|---|")
    groups = {}
    for r in live:
        groups.setdefault(r["advisory2_structure"], []).append(deck_label(r))
    for k in sorted(groups):
        A("| %s | %s |" % (k, ", ".join(groups[k])))
    A("")
    A("---")
    A("")
    A("## Caveats")
    A("")
    A("- Slide typing is inferred from title text and the §7 marker strings, not from")
    A("  a declared slide-type map. A deck whose markers drifted is typed by what it says.")
    A("- Advisory 1 and advisory 2 are keyword heuristics. Every advisory-1 question is")
    A("  quoted in full so the call stays with the teacher.")
    A("- Colors are read from explicit RGB on runs, fills and lines, descending into")
    A("  groups. Theme-inherited colors are invisible to this pass and are not counted.")
    A("- Link counts are unique URLs across the whole deck, from run hyperlinks, shape")
    A("  click actions and printed URLs.")
    A("- This script never opens a deck for write.")
    with open(path, "w") as fh:
        fh.write("\n".join(L) + "\n")


def compare_to(rows, old_csv):
    with open(old_csv) as fh:
        old = list(csv.DictReader(fh))
    old_by = {r["file"]: r for r in old}
    shared = [c for c in old[0].keys() if c in rows[0]] if old else []
    skip = {"mtime", "bytes"}
    print("\n=== cross-check against %s ===" % old_csv)
    print("shared columns: %d" % len(shared))
    only_new = [r["file"] for r in rows if r["file"] not in old_by]
    only_old = [f for f in old_by if f not in {r["file"] for r in rows}]
    if only_new:
        print("in new run only: %s" % ", ".join(only_new))
    if only_old:
        print("in old only:     %s" % ", ".join(only_old))
    total = 0
    for r in rows:
        o = old_by.get(r["file"])
        if not o:
            continue
        diffs = []
        for c in shared:
            if c in skip:
                continue
            a, b = str(r[c]), str(o[c])
            if a != b:
                diffs.append((c, b, a))
        if diffs:
            total += len(diffs)
            print("\n--- %s" % r["file"])
            for c, b, a in diffs:
                print("    %-30s old=%s" % (c, b[:180]))
                print("    %-30s new=%s" % ("", a[:180]))
    print("\ntotal differing cells: %d" % total)


def main(argv=None):
    ap = argparse.ArgumentParser(description=__doc__.split("\n")[1])
    ap.add_argument("exports", nargs="?",
                    default=os.path.expanduser("~/deck_work/exports"))
    ap.add_argument("-o", "--out", default=None,
                    help="output stem (default audits/deck_inventory_<today>)")
    ap.add_argument("--compare", default=None,
                    help="cross-check against an earlier inventory CSV")
    args = ap.parse_args(argv)

    root = os.path.expanduser(args.exports)
    if not (os.path.isdir(root) or os.path.isfile(root)):
        print("no such directory or file: %s" % root, file=sys.stderr)
        return 2

    stem = args.out
    if not stem:
        here = os.path.dirname(os.path.abspath(__file__))
        audits = os.path.join(os.path.dirname(here), "audits")
        stem = os.path.join(audits,
                            "deck_inventory_%s" % datetime.date.today().isoformat())
    outdir = os.path.dirname(os.path.abspath(stem))
    if outdir and not os.path.isdir(outdir):
        os.makedirs(outdir)

    decks = walk_exports(root)
    rows = []
    for path, folder in decks:
        try:
            rows.append(inventory_deck(path, folder))
        except Exception as exc:                      # noqa: BLE001
            print("ERROR %s: %s: %s" % (os.path.basename(path),
                                        type(exc).__name__, exc), file=sys.stderr)
    if not rows:
        print("no decks found under %s" % root, file=sys.stderr)
        return 1

    rows.sort(key=lambda r: (cycle_sort_key(r["cycle"]), r["set"], r["file"]))
    csv_path, md_path = stem + ".csv", stem + ".md"
    with open(csv_path, "w") as fh:
        w = csv.DictWriter(fh, fieldnames=list(rows[0].keys()))
        w.writeheader()
        for r in rows:
            w.writerow(r)
    write_markdown(rows, md_path, root, csv_path)
    print("%d decks -> %s" % (len(rows), csv_path))
    print("%d decks -> %s" % (len(rows), md_path))

    if args.compare:
        compare_to(rows, os.path.expanduser(args.compare))
    return 0


if __name__ == "__main__":
    sys.exit(main())

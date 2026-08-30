#!/usr/bin/env python3
"""Apply the 2026-08-29 change set to every built Biology VT deck.

Four operations, in order:

  1. delete the three slide types Katherine removed
  2. strip [[NOTES:...]] markers left on Concept Bank cells (violates 3b)
  3. capitalise the first letter of every Concept Bank term, matching the
     KEY TERMS convention on the Bellringer
  4. build the slide index - a teacher-facing list of the VT questions in
     order, with what kind of question each one is and what it asks

The index is built by cloning the finished Cycle 02 index slide, so it inherits
its geometry and type scale rather than being restyled deck by deck.

Nothing is overwritten: each deck is written to a new file.

Usage:
  python3 vt_standard_sweep.py <exports_dir> <template.pptx> <out_dir> [--only KEY]
"""
import copy
import os
import re
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt

WRITING_BOX = "F2F6F9"
TEAL = "028090"
REMOVED_PREFIXES = ("Then and Now", "Think → Write", "Turn your answer into a draft")
DAY_DIVIDER = re.compile(r"^Day \d of \d")

# Slide-type markers. Order matters: the first match wins, and the positional
# rule below only sees slides that matched none of these.
TYPED = [
    ("3-Tier Question", lambda t: all(k in t for k in
                                      ("Getting Started", "Working On It", "Mastery"))),
    ("Pattern Break", lambda t: "Pattern break" in t),
    ("Build a Rule", lambda t: "Finish this sentence as a rule" in t),
    ("What if", lambda t: "What if?" in t),
    ("Continuity question", lambda t: "Keep going" in t),
    # Moves 8 and 9, confirmed 2026-08-29. No deck carries one yet - they arrive
    # with the fusion retrofit, which finds them in the content rather than
    # inserting them by rule. Detection is here first so the slide index names
    # them on arrival instead of falling through to the positional rule and
    # calling a compensatory pair a Contrast Set.
    ("Compensatory pair", lambda t: "COMPENSATORY PAIR" in t.upper()
                                    or ("CASE A" in t.upper() and "CASE B" in t.upper())),
    ("Conflict case", lambda t: "CONFLICT CASE" in t.upper()),
]


# ---------------------------------------------------------------- helpers ----
def shapes_text(slide):
    return " \n".join(sh.text_frame.text for sh in slide.shapes if sh.has_text_frame)


def first_line(slide):
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            return sh.text_frame.text.strip().split("\n")[0]
    return ""


def has_writing_box(slide):
    for sh in slide.shapes:
        try:
            if sh.fill.type == 1 and str(sh.fill.fore_color.rgb) == WRITING_BOX:
                return True
        except Exception:                                   # noqa: BLE001
            pass
    return False


def aspect_label(slide):
    """The `Critical aspect: <name>` group label on this slide, or None."""
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for line in sh.text_frame.text.split("\n"):
            m = re.search(r"Critical aspect:\s*(.+?)\s*$", line.strip())
            if m:
                return m.group(1)
    return None


def question_text(slide):
    """The student-facing question on the slide, verbatim.

    The longest paragraph that is not the aspect label, the kicker, or a
    tier label. Katherine's wording is never altered here.
    """
    best = ""
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for para in sh.text_frame.paragraphs:
            t = "".join(r.text for r in para.runs).strip()
            if not t or t.startswith("[["):
                continue
            if t.startswith(("Critical aspect:", "Pattern break", "What if?",
                             "Getting Started", "Working On It", "Mastery",
                             "Keep going", "Explore")):
                continue
            if t.isupper() or len(t) < 25:
                continue
            if len(t) > len(best):
                best = t
    return re.sub(r"\s+", " ", best)


# ------------------------------------------------------------ operation 1 ----
def delete_removed_types(prs):
    lst = prs.slides._sldIdLst
    items = list(lst)
    drop = []
    for i, s in enumerate(prs.slides):
        f = first_line(s)
        if DAY_DIVIDER.match(f):            # 2.5 - never delete a day divider
            continue
        if f.startswith(REMOVED_PREFIXES):
            drop.append(i)
    for i in sorted(drop, reverse=True):
        el = items[i]
        rid = el.get("{http://schemas.openxmlformats.org/officeDocument/"
                     "2006/relationships}id")
        prs.part.drop_rel(rid)
        lst.remove(el)
    return len(drop)


# --------------------------------------------------------- operations 2,3 ----
def clean_concept_bank(prs):
    bank = None
    for s in prs.slides:
        if "Define these in your own words" in shapes_text(s):
            bank = s
            break
    if bank is None:
        return 0, 0
    stripped, capped = 0, 0
    for sh in bank.shapes:
        if not sh.has_text_frame:
            continue
        paras = sh.text_frame.paragraphs
        # 2. marker paragraphs cloned in from a response slide
        for pa in list(paras):
            if pa.runs and pa.runs[0].text.strip().startswith("[[NOTES:"):
                pa._p.getparent().remove(pa._p)
                stripped += 1
        # 3. capitalisation
        paras = sh.text_frame.paragraphs
        if paras and paras[0].runs:
            r = paras[0].runs[0]
            t = r.text
            if t and t[0].islower() and len(t.strip()) < 28:
                r.text = t[0].upper() + t[1:]
                capped += 1
    return stripped, capped


# ------------------------------------------------------------ operation 4 ----
def vt_entries(prs):
    """[(slide_number, type, question_text, aspect_name)] in deck order."""
    slides = list(prs.slides)
    rows = []
    for i, s in enumerate(slides, 1):
        f = first_line(s)
        if DAY_DIVIDER.match(f) or "Define these in your own words" in shapes_text(s):
            continue
        t = shapes_text(s)
        kind = next((k for k, test in TYPED if test(t)), None)
        if kind is None and f.startswith("Explore"):
            kind = "Stock-and-flow model"
        if kind is None:
            if not has_writing_box(s) or aspect_label(s) is None:
                continue
            kind = "?"                       # resolved positionally below
        rows.append([i, kind, question_text(s), aspect_label(s)])

    # positional rule: within one aspect block, the first untyped slide with a
    # writing box is the Critical Aspect question, the second is the Contrast Set
    seen = {}
    for r in rows:
        if r[1] != "?":
            continue
        n = seen.get(r[3], 0)
        r[1] = ["Critical Aspect question", "Contrast Set"][n] if n < 2 else "UNRESOLVED"
        seen[r[3]] = n + 1
    return rows


def build_index(prs, template_slide, entries):
    """Clone the Cycle 02 index slide and refill its two columns."""
    layout = template_slide.slide_layout
    try:
        new = prs.slides.add_slide(prs.slide_layouts[len(prs.slide_layouts) - 1])
    except Exception:                                       # noqa: BLE001
        new = prs.slides.add_slide(layout)
    for sh in list(new.shapes):
        sh._element.getparent().remove(sh._element)
    for sh in template_slide.shapes:
        new.shapes._spTree.append(copy.deepcopy(sh._element))

    cols = sorted((sh for sh in new.shapes
                   if sh.has_text_frame and len(sh.text_frame.paragraphs) > 4),
                  key=lambda sh: sh.left)
    if len(cols) != 2:
        return None

    # group the entries: aspect 1, aspect 2, then anything spanning both
    # A What if slide carries the label of whichever aspect it grew out of, but
    # it is asked after both aspects and answers across them, so it is listed
    # as spanning rather than inside one aspect.
    SPAN = "ACROSS BOTH ASPECTS"

    # A stock-and-flow model, and a continuity question asked before the first
    # aspect opens, carry no `Critical aspect:` label of their own. They belong
    # to the block they sit in, so they inherit the nearest labelled aspect -
    # looking back first, then forward for anything ahead of the first block.
    # Only the What if slide is genuinely listed as spanning.
    labelled = [a for _, _, _, a in entries if a]
    filled, last = [], None
    for n, kind, text, aspect in entries:
        if aspect:
            last = aspect
        filled.append((n, kind, text, aspect or last or (labelled[0] if labelled else None)))
    entries = filled

    order, groups = [], {}
    for n, kind, text, aspect in entries:
        key = SPAN if kind == "What if" else (aspect or SPAN)
        if key not in groups:
            groups[key] = []
            order.append(key)
        groups[key].append((n, kind, text))
    aspects = [k for k in order if k != SPAN]
    blocks = [("CRITICAL ASPECT %d · %s" % (i + 1, a), groups[a])
              for i, a in enumerate(aspects)]
    if SPAN in groups:
        blocks.append((SPAN, groups[SPAN]))

    # Balance the two columns by how much text each holds, choosing the block
    # boundary with the smallest imbalance. Splitting on block count alone put
    # two five-question aspects in one column and a single What if in the other.
    def weight(rows):
        return sum(len(t) for _, _, t in rows)

    if len(blocks) > 1:
        left = min(range(1, len(blocks)),
                   key=lambda k: abs(sum(weight(r) for _, r in blocks[:k]) -
                                     sum(weight(r) for _, r in blocks[k:])))
    else:
        left = 1

    for col, chunk in zip(cols, (blocks[:left], blocks[left:])):
        tf = col.text_frame
        for pa in list(tf.paragraphs)[1:]:
            pa._p.getparent().remove(pa._p)
        first = tf.paragraphs[0]
        for r in list(first.runs):
            r._r.getparent().remove(r._r)
        started = False
        for heading, rows in chunk:
            pa = first if not started else tf.add_paragraph()
            started = True
            r = pa.add_run()
            r.text = heading
            r.font.size, r.font.bold = Pt(9), True
            r.font.color.rgb = RGBColor.from_string("777777")
            for n, kind, text in rows:
                p2 = tf.add_paragraph()
                a = p2.add_run()
                a.text = "%d · %s" % (n, kind)
                a.font.size, a.font.bold = Pt(9), True
                a.font.color.rgb = RGBColor.from_string(TEAL)
                p3 = tf.add_paragraph()
                b = p3.add_run()
                b.text = text
                b.font.size = Pt(9)
                b.font.color.rgb = RGBColor.from_string("111111")
    return new


def move_before_links(prs, slide):
    lst = prs.slides._sldIdLst
    items = list(lst)
    target = None
    for i, s in enumerate(prs.slides):
        if first_line(s).startswith("Activity and resource links"):
            target = i
            break
    if target is None:
        return
    el = items[-1]
    lst.remove(el)
    lst.insert(target, el)


# ------------------------------------------------------------------ main ----
def sweep(path, template_slide, out_dir):
    prs = Presentation(path)
    name = os.path.basename(path)
    res = {"deck": name, "before": len(prs.slides._sldIdLst)}

    res["deleted"] = delete_removed_types(prs)
    if res["deleted"]:
        # drop_rel unlinks the slide but leaves its part in the package, so a
        # slide added afterwards claims a partname that is still occupied and
        # the saved .pptx carries duplicate zip entries. Round-tripping through
        # a temporary file drops the orphans and frees the names.
        import tempfile
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as fh:
            tmp = fh.name
        prs.save(tmp)
        prs = Presentation(tmp)
        os.unlink(tmp)
    res["markers"], res["capped"] = clean_concept_bank(prs)

    entries = vt_entries(prs)
    res["questions"] = len(entries)
    res["unresolved"] = sum(1 for e in entries if e[1] == "UNRESOLVED")
    idx = build_index(prs, template_slide, entries)
    res["index"] = "yes" if idx is not None else "FAILED"
    # rough overflow guard: the column box holds about 46 lines at 9pt, and a
    # line of this width takes about 62 characters
    if idx is not None:
        worst = 0
        for sh in idx.shapes:
            if not sh.has_text_frame or len(sh.text_frame.paragraphs) <= 4:
                continue
            lines = sum(max(1, -(-len(p.text) // 62)) for p in sh.text_frame.paragraphs)
            worst = max(worst, lines)
        res["lines"] = worst
        res["fits"] = "yes" if worst <= 46 else "CHECK"
    if idx is not None:
        move_before_links(prs, idx)

    res["after"] = len(prs.slides._sldIdLst)
    key = re.search(r"Cycle\s*(\d{2}[a-d]?)", name).group(1)
    out = os.path.join(out_dir, "IMPORT — Cycle %s FINAL.pptx" % key)
    prs.save(out)
    res["out"] = os.path.basename(out)
    return res


def main():
    exports, template_path, out_dir = sys.argv[1], sys.argv[2], sys.argv[3]
    only = None
    if "--only" in sys.argv:
        only = sys.argv[sys.argv.index("--only") + 1]
    os.makedirs(out_dir, exist_ok=True)

    tprs = Presentation(template_path)
    template_slide = next(s for s in tprs.slides
                          if "TEACHER NAVIGATION" in shapes_text(s))

    paths = []
    for d in sorted(os.listdir(exports)):
        full = os.path.join(exports, d)
        if not os.path.isdir(full):
            continue
        for f in sorted(os.listdir(full)):
            if f.endswith("with Concept Bank.pptx"):
                if only and only not in f:
                    continue
                paths.append(os.path.join(full, f))

    cols = ["deck", "before", "after", "deleted", "markers", "capped",
            "questions", "unresolved", "index", "lines", "fits"]
    out_rows = []
    for p in paths:
        try:
            out_rows.append(sweep(p, template_slide, out_dir))
        except Exception as e:                              # noqa: BLE001
            out_rows.append({"deck": os.path.basename(p), "index": "ERROR: %s" % e})
    w = {c: max(len(c), max(len(str(r.get(c, ""))) for r in out_rows)) for c in cols}
    w["deck"] = min(w["deck"], 52)
    print(" ".join(c.ljust(w[c]) for c in cols))
    for r in out_rows:
        print(" ".join(str(r.get(c, ""))[:w[c]].ljust(w[c]) for c in cols))


if __name__ == "__main__":
    main()

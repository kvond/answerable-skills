#!/usr/bin/env python3
"""On the Day 3 divider, replace the four-checkbox block with links to the
slides that actually need revising.

Katherine, 2026-08-30: the checklist told the student to go and find her
response slides. The links take her there.

The checklist is deleted and, in the box it occupied, one line per revision
slide: its number, what it asks, and a live jump to it.

Internal slide links are not something python-pptx writes, so the relationship
and the `ppaction://hlinksldjump` action are added to the XML directly. That is
the same mechanism PowerPoint's own "Link to this slide" uses.

Usage:  python3 revision_links.py <dir> [--apply]
"""
import argparse
import copy
import glob
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.util import Pt

TEAL = "028090"
INK = "111111"
CHECKLIST_MARK = "Reopen this deck"
HEADING = "Go back to your answers and revise them — right on the slides"


def texts(slide):
    return [sh.text_frame.text.strip() for sh in slide.shapes
            if sh.has_text_frame and sh.text_frame.text.strip()]


def first_line(slide):
    t = texts(slide)
    return t[0].split("\n")[0] if t else ""


def revision_slides(prs):
    """Slides whose notes carry the revision prompt - the ones to link to."""
    out = []
    for i, s in enumerate(prs.slides, 1):
        if s.has_notes_slide and "REVISION PROMPT" in s.notes_slide.notes_text_frame.text:
            label = first_line(s)
            # the response slides say "Your answer — ..."; trim to the aspect
            label = label.replace("Your answer — ", "")
            out.append((i, label))
    return out


def link_run_to_slide(source_part, run, target_slide_part):
    """Give `run` an internal jump to `target_slide_part`."""
    rId = source_part.relate_to(target_slide_part, RT_SLIDE)
    rPr = run._r.get_or_add_rPr()
    for old in rPr.findall(qn("a:hlinkClick")):
        rPr.remove(old)
    hl = rPr.makeelement(qn("a:hlinkClick"), {})
    hl.set(qn("r:id"), rId)
    hl.set("action", "ppaction://hlinksldjump")
    # hlinkClick must be the last child of rPr per the schema
    rPr.append(hl)


RT_SLIDE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide"


def rebuild(prs, divider_idx, box, targets):
    slide_part = prs.slides[divider_idx]._element  # noqa: F841 - clarity
    src_part = prs.slides[divider_idx].part
    tf = box.text_frame
    for para in list(tf.paragraphs)[1:]:
        para._p.getparent().remove(para._p)
    p0 = tf.paragraphs[0]
    for r in list(p0.runs):
        r._r.getparent().remove(r._r)

    started = False
    for n, label in targets:
        para = p0 if not started else tf.add_paragraph()
        started = True
        para.space_after = Pt(6)
        a = para.add_run()
        a.text = "Slide %d — " % n
        a.font.size, a.font.bold = Pt(14), True
        a.font.color.rgb = RGBColor.from_string(INK)
        a.font.name = "Arial"
        b = para.add_run()
        b.text = label
        b.font.size = Pt(14)
        b.font.color.rgb = RGBColor.from_string(TEAL)
        b.font.underline = True
        b.font.name = "Arial"
        link_run_to_slide(src_part, b, prs.slides[n - 1].part)
    return len(targets)


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("path")
    ap.add_argument("--apply", action="store_true")
    args = ap.parse_args()
    files = sorted(glob.glob(os.path.join(args.path, "*.pptx")))

    total = skipped = 0
    for f in files:
        prs = Presentation(f)
        slides = list(prs.slides)
        div = box = None
        for i, s in enumerate(slides):
            for sh in s.shapes:
                if sh.has_text_frame and CHECKLIST_MARK in sh.text_frame.text:
                    div, box = i, sh
                    break
            if div is not None:
                break
        if div is None:
            continue

        targets = revision_slides(prs)
        name = os.path.basename(f)[9:-11]
        if not targets:
            print("  %-10s slide %-3d ** no slide carries a revision prompt — left alone **"
                  % (name, div + 1))
            skipped += 1
            continue
        print("  %-10s slide %d → links to %s"
              % (name, div + 1, ", ".join("s%d" % n for n, _ in targets)))
        for n, label in targets:
            print("               %-4s %s" % ("s%d" % n, label[:66]))
        if args.apply:
            rebuild(prs, div, box, targets)
            prs.save(f)
        total += 1

    print("\n%s: %d divider(s), %d left alone"
          % ("APPLIED" if args.apply else "WOULD REBUILD", total, skipped))
    if not args.apply:
        print("Report only. Re-run with --apply to write.")


if __name__ == "__main__":
    main()

#!/usr/bin/env python3
"""bio_v3_extract.py — deterministic v3 extractor for the on-slide formative pipeline.

NEW CODE (2026-08-10, Cowork session): written for the v3 Google-Slides collection route.
Implements the behaviors the deployed extract_v2.py established (documented in the project
MANIFEST): response slides detected by BOTH answer labels on the template (case-insensitive
PREFIX match, never equality); answers attributed geometrically (a label claims the nearest
answer shape below it); the Then-and-Now slide never enters analysis; printed slide number
emitted beside the deck position (student rail = printed, teacher rail = position).

Usage: bio_v3_extract.py <teacher_template.pptx> <submissions_dir> <out_dir> --lesson NAME --roster roster.txt
Outputs: <out_dir>/extracted_v3.json, <out_dir>/completion_report.csv
No judgment here: completion = both areas carry any student text. Word counts and
similarity signals are emitted for the judgment layer; nothing is scored or flagged here.
"""
import sys, os, json, re, csv, argparse, unicodedata
from difflib import SequenceMatcher
from pptx import Presentation

FIRST = "your first answer"
REV = "your revised answer"

def norm(s):
    s = unicodedata.normalize("NFKC", s or "")
    return re.sub(r"\s+", " ", s).strip()

def slide_shapes(slide):
    return [sh for sh in slide.shapes if sh.has_text_frame]

def find_response_slides(prs):
    """Return [(deck_pos_1based, printed_number, question_printed, question_text)]."""
    out = []
    slides = list(prs.slides)
    for i, s in enumerate(slides):
        texts = [norm(sh.text_frame.text).lower() for sh in slide_shapes(s)]
        if any(t.startswith(FIRST) for t in texts) and any(t.startswith(REV) for t in texts):
            printed = None
            for sh in slide_shapes(s):
                t = norm(sh.text_frame.text)
                if re.fullmatch(r"\d{1,2}[ab]?", t) and sh.width.inches < 1.2:
                    printed = t
            # Mastery question from the nearest earlier slide carrying the base printed number
            q_printed = printed[:-1] if printed and printed.endswith("b") else None
            q_text = None
            if q_printed is not None:
                for j in range(i - 1, -1, -1):
                    nums = [norm(sh.text_frame.text) for sh in slide_shapes(slides[j])
                            if re.fullmatch(r"\d{1,2}", norm(sh.text_frame.text)) and sh.width.inches < 1.2]
                    if q_printed in nums:
                        joined = [norm(sh.text_frame.text) for sh in slide_shapes(slides[j])]
                        try:
                            k = next(x for x, t in enumerate(joined) if t == "Mastery")
                            q_text = joined[k + 1] if k + 1 < len(joined) else None
                        except StopIteration:
                            pass
                        break
            out.append((i + 1, printed, q_printed, q_text))
    return out

def template_texts(prs):
    """All normalized template text blocks, for new-text and echo comparison."""
    blocks = set()
    for s in prs.slides:
        for sh in slide_shapes(s):
            t = norm(sh.text_frame.text)
            if t:
                blocks.add(t)
    return blocks

def extract_areas(slide, tmpl_slide_blocks):
    """Return (first_answer, revised_answer) using label geometry.

    Excludes only text that exists on the SAME slide of the template (furniture:
    instructions, headers). Text a student typed that copies another slide's
    words still counts as their (copied) answer — the integrity signals catch it.
    """
    shapes = slide_shapes(slide)
    def label(pref):
        c = [sh for sh in shapes if norm(sh.text_frame.text).lower().startswith(pref)]
        return c[0] if c else None
    lab_f, lab_r = label(FIRST), label(REV)
    if not (lab_f and lab_r):
        return None, None
    def band_text(top, bottom):
        out = []
        for sh in shapes:
            if sh in (lab_f, lab_r):
                continue
            if top < sh.top and (bottom is None or sh.top < bottom):
                t = norm(sh.text_frame.text)
                if t and t not in tmpl_slide_blocks:   # exclude same-slide template furniture only
                    out.append(t)
        return " ".join(out).strip()
    first = band_text(lab_f.top, lab_r.top)
    revised = band_text(lab_r.top, None)
    return first, revised

def sim(a, b):
    return SequenceMatcher(None, a.lower(), b.lower()).ratio()

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("template"); ap.add_argument("subs_dir"); ap.add_argument("out_dir")
    ap.add_argument("--lesson", default=""); ap.add_argument("--roster", default=None)
    args = ap.parse_args()
    os.makedirs(args.out_dir, exist_ok=True)

    tmpl = Presentation(args.template)
    resp = find_response_slides(tmpl)
    if not resp:
        sys.exit("HALT: no response slides on template (need both answer labels).")
    tmpl_blocks = template_texts(tmpl)

    roster = []
    if args.roster:
        roster = [l.strip() for l in open(args.roster) if l.strip()]

    students = {}
    for fn in sorted(os.listdir(args.subs_dir)):
        if not fn.lower().endswith(".pptx"):
            continue
        name = re.split(r"\s+[—-]\s+", fn[:-5])[0].strip()
        prs = Presentation(os.path.join(args.subs_dir, fn))
        slides = list(prs.slides)
        entry = {"file": fn, "slides": {}}
        for pos, printed, q_printed, q_text in resp:
            if pos - 1 >= len(slides):
                entry["slides"][pos] = {"error": "slide missing"}
                continue
            tmpl_slide_blocks = {norm(sh.text_frame.text)
                                 for sh in slide_shapes(list(tmpl.slides)[pos - 1])}
            first, revised = extract_areas(slides[pos - 1], tmpl_slide_blocks)
            echo = max((sim(first, t) for t in tmpl_blocks if len(t) > 20), default=0.0) if first else 0.0
            substr = bool(first) and len(first) > 25 and any(
                norm(first).lower() in t.lower() or t.lower() in norm(first).lower()
                for t in tmpl_blocks if len(t) > 25)
            entry["slides"][pos] = {
                "printed": printed, "question_printed": q_printed, "question": q_text,
                "first_answer": first or "", "revised_answer": revised or "",
                "first_words": len((first or "").split()), "revised_words": len((revised or "").split()),
                "first_echo_similarity": round(echo, 3), "first_echo_substring": substr}
        entry["complete"] = all(v.get("first_answer") and v.get("revised_answer")
                                for v in entry["slides"].values() if "error" not in v)
        students[name] = entry

    # pairwise copy signal on first answers
    names = list(students)
    for i, a in enumerate(names):
        peers = {}
        for b in names:
            if b == a: continue
            vals = []
            for pos in students[a]["slides"]:
                fa = students[a]["slides"][pos].get("first_answer", "")
                fb = students[b]["slides"].get(pos, {}).get("first_answer", "")
                if fa and fb: vals.append(sim(fa, fb))
            if vals: peers[b] = round(max(vals), 3)
        students[a]["max_peer_similarity"] = max(peers.values(), default=0.0)

    out = {"lesson": args.lesson, "template": os.path.basename(args.template),
           "response_slides": [{"deck_pos": p, "printed": pr, "q_printed": qp, "question": qt}
                                for p, pr, qp, qt in resp],
           "students": students}
    with open(os.path.join(args.out_dir, "extracted_v3.json"), "w") as f:
        json.dump(out, f, indent=1)

    with open(os.path.join(args.out_dir, "completion_report.csv"), "w", newline="") as f:
        w = csv.writer(f)
        w.writerow(["Student", f"Completion.{args.lesson or 'Lesson'}", "Status"])
        everyone = roster or names
        for n in everyone:
            if n in students:
                st = students[n]
                w.writerow([n, 1 if st["complete"] else 0, "Complete" if st["complete"] else "Incomplete"])
            else:
                w.writerow([n, "", "No submission"])
    print(f"OK {len(students)} decks · response slides {[(p, pr) for p, pr, _, _ in resp]}")

if __name__ == "__main__":
    main()

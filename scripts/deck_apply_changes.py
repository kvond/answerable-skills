#!/usr/bin/env python3
"""deck_apply_changes.py — apply the 2026-08-29 design changes to the Biology VT decks.

The counterpart to deck_lint.py. The linter reports; this applies — but only what
is safe to apply mechanically. Anything that is a content decision is FLAGGED for
Katherine and the slide is left untouched.

Four changes, and only four:

  1. INSERT the teacher note slide (one slide, not speaker notes, deletable in one
     action) carrying the seven declarations from the 2026-08-29 spec. Everything
     derivable from the deck is filled in. The three that require her judgment —
     what is held invariant, what breaks if an example is substituted, and the
     cycle's current visibility rung — are emitted as `NEEDS KATHERINE:` questions
     in alert red C0392B, so an unfilled placeholder is obvious on the slide and
     greppable by the linter. Those three are never invented.

  2. APPEND the relating prompt to a Concept Bank that only lists terms.
     Co-presence is the precondition, not the achievement (09 §7): a page that
     lists the terms has left the relating to chance.

  3. FLAG, never rewrite, a move-1 Critical Aspect question that carries no
     difference. Rewriting it is a content decision. The verbatim text and a
     template reframing go into the report for her to accept or reject.

  4. INSERT the visibility-ladder slide (added 2026-08-29), immediately AFTER
     the teacher note. The note declares the cycle's rung; this slide explains
     what the five rungs are, so the declaration means something to a teacher
     who has never seen the ladder. Teacher-facing and never projected: it
     carries "do not project" and is excluded from every diagnostic count.
     Its copy is fixed rather than derived from the deck, so nothing on it
     needs Katherine's judgment and it carries no placeholder.

Conventions follow scripts/build_concept_banks.py: the input file is never
touched, the result is written to a new file alongside it, and the new slide is
built from the deck's own layout with the §13 design tokens.

    --dry-run is the DEFAULT. Writing requires an explicit --write.

Run with the deck_work venv interpreter, which is the one that has python-pptx:

    ~/deck_work/.venv/bin/python scripts/deck_apply_changes.py --help
"""

import argparse
import csv
import datetime
import glob
import os
import re
import sys
import unicodedata

from pptx import Presentation
from pptx.opc.packuri import PackURI
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor

# --------------------------------------------------------------------------
# §13 design tokens. Do not invent a colour, a size, or a position.
# --------------------------------------------------------------------------

SLIDE_W = 9144000
SLIDE_H = 6858000
FONT = "Arial"

TEAL = "028090"          # kickers, section labels
BODY = "111111"          # all body text. Not 000000
MUTED = "666666"         # teacher-reference kicker, small labels
ALERT = "C0392B"         # the placeholder colour — deliberately loud
WHITE = "FFFFFF"

PALETTE = {
    TEAL, BODY, MUTED, ALERT, "1E8449", "F2F6F9", "CCCCCC",
    "F5F5F5", "EAF6F7", WHITE,
}

# Geometry, copied from the Teacher Prep slide (§13.3 / §2.0), which is the
# slide this one sits beside.
G_KICKER = (502920, 109728, 8321040, 274320)
G_HEADING = (548640, 548640, 8046720, 457200)
G_SUBHEAD = (548640, 1042000, 8046720, 281264)
G_COL_L = (548640, 1400000, 3886200, 5100000)
G_COL_R = (4709160, 1400000, 3886200, 5100000)

# The relating prompt sits in the house bottom-line slot, the same y the
# "Think first" footer uses (§13.3). The Concept Bank grid ends at 5669280 and
# its links block runs to 6309360, so this is the only free band on the slide.
G_RELATE = (548640, 6400800, 8046720, 320040)

# The visibility-ladder slide (added 2026-08-29). It sits immediately after the
# teacher note and reuses that slide's kicker, heading and subhead geometry, so
# the two teacher-facing slides read as one pair rather than as two designs.
# Its body is a SINGLE column: a ladder read across two columns stops being a
# ladder, and the rungs have to run top to bottom in order.
# The body runs to 6675120 — 0.2in above the slide edge — because this slide
# carries no "Think first" footer, exactly as the Concept Bank carries none.
# The subhead runs to two lines here, where the teacher note's runs to one.
G_VL_SUBHEAD = (548640, 1042000, 8046720, 457200)
# y and h are both §13.3 tokens: 1554480 is the Concept Bank's first row, and
# 5120640 is its last. The body ends at 6675120, 0.2in above the slide edge.
G_VL_BODY = (548640, 1554480, 8046720, 5120640)

# 11pt, not 12. MEASURED, not guessed — against the real Arial advance widths
# on this machine, wrapping greedily the way PowerPoint does:
#
#   text width  8046720 - 2 x 91440 inset  = 619.2pt
#   text height 5120640 - 2 x 45720 inset  = 395.9pt
#   the eight blocks below                 = 27 lines
#   27 x 11pt x 1.2 + 7 x 3pt space-after  = 377.4pt   (18.5pt spare)
#
# 1.2 is the conservative line factor; Arial's own hhea metrics give 1.150,
# which leaves 33pt spare. At 12pt the same copy is 31 lines and runs off the
# bottom of the slide. 11pt is an existing §13.4 token (the kicker, the
# word-bank items, the Concept Bank source links) and it is the same fallback
# the teacher note already takes.
SZ_VL_BODY = 11
SZ_VL_SPACE_AFTER = 3

# What _est_lines() scores the shipped copy at. It is a tripwire, not an
# oracle: the estimator assumes an 0.5-em average character and ignores the
# text-frame insets, so it reads 29 lines where real Arial reads 27, and its
# own capacity number (28) is meaningless against that. What it is good for is
# detecting an EDIT — if VL_BLOCKS grows past this, the copy changed and the
# slide has to be measured against Arial again before it ships.
VL_EST_LINES = 29

SZ_KICKER = 11
SZ_HEADING = 16
SZ_SUBHEAD = 13
SZ_LABEL = 10
SZ_BODY = 12
# The teacher note is the densest teacher-facing role in the deck. 12pt is
# tried first; 11pt (an existing token, cf. the kicker and word-bank items)
# is the fallback when the seven declarations will not fit at 12.
SZ_NOTE_STEPS = (12, 11)

# --------------------------------------------------------------------------
# Marker strings. §3c: matched by machine, not read by a person. ASCII, verbatim.
# --------------------------------------------------------------------------

# The kicker deliberately contains "do not project", which is already in
# NON_DIAGNOSTIC_MARKERS in extract_and_grade.py and NON_DIAGNOSTIC in
# deck_lint.py (§7.1). Both tuples are tested BEFORE the four diagnostic
# classifiers, so the teacher note cannot be miscounted as a Pattern Break
# just because its slide-type map names one. No script edit is needed.
NOTE_KICKER = "TEACHER REFERENCE — do not project to students"
NOTE_TITLE = "Teacher note — the design of this cycle"
NOTE_MARKER = "TEACHER NOTE"          # ASCII, for the linter
PLACEHOLDER = "NEEDS KATHERINE:"      # ASCII, for the linter

RELATING_PROMPT = (
    "Pick any two terms on this page. Write one sentence saying how they connect."
)

CONCEPT_BANK_HEADING = "Define these in your own words"
CONCEPT_BANK_KICKER = "CONCEPT BANK"

# Strings that must NEVER appear on the teacher note, because they are unique
# slide-type markers (§7) and a second occurrence makes a teacher prompt count
# the slide twice. The slide-type map therefore says "What if" with no question
# mark and never says "Getting Started".
FORBIDDEN_ON_NOTE = ("what if?", "getting started", "finish this sentence as a rule",
                     "critical aspect:", "pattern break")

# The teacher note has to NAME the slide types, and two of those names are
# themselves unique marker strings (§7). Naming them plainly would make the
# teacher prompt count the note as a second Pattern Break slide — the exact
# fault §7 warns about. Both scripts match on a literal space, so a hyphen
# reads identically to a person and is invisible to the matchers.
NOTE_SAFE_NAME = {
    "Pattern Break": "Pattern-Break",
}

# --------------------------------------------------------------------------
# The visibility-ladder slide (added 2026-08-29)
# --------------------------------------------------------------------------
# The teacher note DECLARES this cycle's rung. This slide EXPLAINS the five
# rungs, so the declaration means something to a teacher who has never seen the
# ladder. It is teacher-facing and never projected, so its kicker carries "do
# not project" — already in NON_DIAGNOSTIC in deck_lint.py and in
# NON_DIAGNOSTIC_MARKERS in the grader (§7.1), both tested before the four
# diagnostic classifiers. The slide therefore cannot inflate a diagnostic count.
#
# THE MARKER-STRING TRAP, AND IT IS REAL. deck_lint.py and the teacher prompts
# match the slide-type markers as case-insensitive substrings ANYWHERE in the
# deck. Writing "Pattern Break" in the ceilings line below would make this deck
# count two Pattern Breaks where it has one. The exclusion tuple saves the
# linter's own count and does not save a raw grep, which is what §11 check 3
# actually runs. So every move name on this slide goes through NOTE_SAFE_NAME
# or is written in a form that has no marker in it — a hyphen reads identically
# to a person and is invisible to the matchers.

VL_MARKER = "VISIBILITY LADDER"        # ASCII, uppercase, for the classifiers
VL_KICKER = ("TEACHER REFERENCE — do not project to students   ·   "
             + VL_MARKER)
VL_TITLE = "The visibility ladder — what happens at each rung"

VL_SUBHEAD = (
    "Slide %d declares the rung this cycle is built for. Every gate below is "
    "something you can see from the front of the room in one period, not a week "
    "on a calendar. Raise visibility across the year; reset it low inside every "
    "cycle."
)
VL_SUBHEAD_NO_NOTE = (
    "The teacher note declares the rung a cycle is built for. Every gate below "
    "is something you can see from the front of the room in one period, not a "
    "week on a calendar. Raise visibility across the year; reset it low inside "
    "every cycle."
)

# (bold teal label, body). Source: 10 §3 for the rungs and their gates, §3c for
# group work, 09 §6a for the compression. Rung 3 gets the longest block on the
# slide on purpose — it is the one that does the work and the one no PD teaches.
VL_BLOCKS = [
    ("Rung 1 — written, seen only by you.",
     " You take the first answer and score it for completion; the revision is "
     "where thinking counts. Nobody else reads it. Opens rung 2 when most "
     "students write a first answer at all and some are wrong in interesting "
     "ways rather than blank or copied. A blank is a refusal to commit even "
     "privately, so rung 1 is not secure yet."),
    ("Rung 2 — written, shown to one assigned partner.",
     " You assign the partner. Chosen partners reproduce the social order "
     "already in the room, and the students most at risk choose nobody. Opens "
     "rung 3 when pairs talk about the answer rather than about something else, "
     "and some students change or extend an answer after talking."),
    ("Rung 3 — you read a wrong but productive answer aloud, unattributed, and "
     "use it.",
     " This is the rung that does the work, and the one no PD teaches. Take a "
     "genuinely useful wrong answer off the stack — not a straw man, they can "
     "tell immediately. Read it with no name attached. Then use it: build on "
     "it, ask what it assumes, put it against a case. Do not correct it and "
     "move on. Never signal whose it is, including by looking at her. Four or "
     "five instances change what a wrong answer means in this room. Opens rung "
     "4 when students volunteer answers you did not ask for, or one of them "
     "says a version of “that was mine” without distress."),
    ("Rung 4 — a student owns an answer aloud, by invitation.",
     " Arrange it privately where you can — “your answer about X is the one I "
     "want to start with; will you say it?” She can do it because she watched "
     "rung 3 and knows what will be done with what she says. Opens rung 5 when "
     "more than a few students speak, disagreement happens between students "
     "rather than through you, and a wrong answer aloud does not end that "
     "student’s participation for the period."),
    ("Rung 5 — public simultaneous commitment.",
     " The vote: pair talk, sequential hands, two students justify, re-vote. "
     "Last, not first. It asks for commitment and exposure at once, and it "
     "works only where everything above has already made the exposure cheap."),
    ("Descent is normal.",
     " A new class, a new semester, a hard week, a student humiliated somewhere "
     "else that day — any of it resets the room. Drop a rung, and do not "
     "announce that you have."),
    ("Group work is a visibility reducer, not a rung.",
     " A student speaks to three people instead of thirty and the product "
     "carries the group’s name, so a group can carry harder content than the "
     "room’s whole-class rung allows. But a group product says nothing about "
     "who did the relating — a good occasion, and no evidence. That is why the "
     "What-if is written and individual."),
    ("Ceilings by move.",
     " The Critical Aspect question and the Contrast Set tolerate the most — "
     "read aloud unattributed, or a vote. Build-a-Rule: a partner. "
     "Pattern-Break and the 3-Tier Question stay written and private — that is "
     "where her own rule breaks. Work at or below the ceiling."),
]

# Strings that must NEVER appear on the visibility-ladder slide. The first five
# are the unique slide-type markers of §7 — a second occurrence makes a teacher
# prompt count this slide as a Pattern Break or a 3-Tier. The rest would make a
# classifier type this slide as some other slide the deck is REQUIRED to carry,
# and a deck that loses its teacher note or its Concept Bank to a misfiled
# teacher slide fails hard for a reason nobody would find.
FORBIDDEN_ON_LADDER = FORBIDDEN_ON_NOTE + (
    "working on it", "mastery", "define these in your own words", "concept bank",
    "teacher note", "design note", "note to the teacher", "teacher navigation",
    "essential claim", "additional resources", "teacher prep",
    "continuation question", "relates to me", "optional challenge",
    "image credits", "activity and resource links", "case a", "case b",
    "your first answer", "bellringer", "tank model",
)

# The word "beat" is forbidden everywhere, including in this file (spec, Vocabulary).

# --------------------------------------------------------------------------
# The nine moves and their visibility ceilings (10 §3b).
# --------------------------------------------------------------------------

def note_name(move):
    """The move's name as it may appear on the teacher note (see NOTE_SAFE_NAME)."""
    n = MOVE_NAME[move]
    return NOTE_SAFE_NAME.get(n, n)


MOVES = [
    ("1", "Critical Aspect question", "read aloud unattributed"),
    ("2", "Contrast Set", "read aloud unattributed, or vote"),
    ("3", "Build a Rule", "partner"),
    ("4", "Pattern Break", "written and private"),
    ("5", "3-Tier Question", "written and private"),
    ("6", "Continuity question", "partner, or aloud"),
    ("7", "Stock-and-flow model", "partner, or group"),
    ("8", "Compensatory pair", "group"),
    ("9", "Conflict case", "group"),
    ("W", "What-if", "written and individual"),
]
CEILING = dict((m[0], m[2]) for m in MOVES)
MOVE_NAME = dict((m[0], m[1]) for m in MOVES)

CORE_MOVES = ("1", "2", "3", "4", "5")
CONDITIONAL_MOVES = ("6", "7", "8", "9")

# inventory slide-type label -> move number
SLIDE_TYPE_TO_MOVE = {
    "critical aspect question": "1",
    "contrast set": "2",
    "build a rule": "3",
    "pattern break": "4",
    "3-tier question": "5",
    "continuity question": "6",
    "stock-and-flow model": "7",
    "response slide (what if)": "W",
}
# advisory2_structure label -> move number
STRUCTURE_TO_MOVE = {
    "stock-and-flow model": "7",
    "compensatory pair": "8",
    "conflict case": "9",
}

NO_DIFFERENCE_VERDICT = "NO CONTRAST DEVICE"

# --------------------------------------------------------------------------
# small helpers
# --------------------------------------------------------------------------


def norm(s):
    """Fold the typographic characters §3c warns about, for matching only."""
    s = unicodedata.normalize("NFKC", s or "")
    for ch in "‐‑‒–—―−":
        s = s.replace(ch, "-")
    s = s.replace("‘", "'").replace("’", "'")
    s = s.replace("“", '"').replace("”", '"')
    s = s.replace("…", "...").replace(" ", " ")
    return re.sub(r"\s+", " ", s).strip()


def slide_text(slide):
    return "\n".join(
        sh.text_frame.text for sh in slide.shapes
        if sh.has_text_frame and sh.text_frame.text
    )


def split_multi(value, sep=" ;; "):
    return [p.strip() for p in (value or "").split(sep) if p.strip()]


def cycle_sort_key(key):
    m = re.match(r"^(\d+)([a-z]?)$", key or "")
    if not m:
        return (999, "")
    return (int(m.group(1)), m.group(2))


# --------------------------------------------------------------------------
# format tokens (§13) — the gate
# --------------------------------------------------------------------------


def check_format_tokens(prs):
    """Return (blocking, strict, detail). blocking = cannot place a house slide."""
    blocking, strict = [], []
    if (prs.slide_width, prs.slide_height) != (SLIDE_W, SLIDE_H):
        blocking.append(
            "slide size %dx%d, not 4:3 %dx%d"
            % (prs.slide_width, prs.slide_height, SLIDE_W, SLIDE_H)
        )

    fonts, colors, runs = set(), set(), 0
    for slide in prs.slides:
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            for para in sh.text_frame.paragraphs:
                for r in para.runs:
                    runs += 1
                    if r.font.name:
                        fonts.add(r.font.name)
                    try:
                        if r.font.color and r.font.color.type is not None:
                            colors.add(str(r.font.color.rgb))
                    except Exception:
                        pass

    bad_fonts = sorted(f for f in fonts if f != FONT)
    bad_colors = sorted(c for c in colors if c not in PALETTE)

    if runs and FONT not in fonts:
        blocking.append("Arial appears on no run — this is not a house deck")

    if bad_fonts:
        strict.append("non-Arial fonts: " + "; ".join(bad_fonts))
    if bad_colors:
        strict.append("off-palette text colors: " + "; ".join(bad_colors))

    return blocking, strict, {"fonts": sorted(fonts), "colors": sorted(colors)}


# --------------------------------------------------------------------------
# move-1 difference test (08 "Move 1 — the flag")
# --------------------------------------------------------------------------

RE_COMPARISON = re.compile(
    r"\b(more|less|greater|fewer|higher|lower|faster|slower|bigger|smaller|"
    r"stronger|weaker|better|worse|differ|different|difference|compare|than|"
    r"versus|vs\.?)\b", re.I)
RE_ALTERNATIVES = re.compile(
    r"\b(which (one|of|is|are)|either\b.*\bor\b|\bA or B\b|choose|pick)\b", re.I)
RE_CHANGE = re.compile(
    r"\b(what happens (if|when)|if you|once you|after you|when you (change|heat|"
    r"remove|add)|what changes|change[sd]? (the|its)|is removed|stops? working)\b",
    re.I)


def carries_difference(text):
    t = norm(text)
    found = []
    if RE_COMPARISON.search(t):
        found.append("comparison")
    if RE_ALTERNATIVES.search(t):
        found.append("named alternatives")
    if RE_CHANGE.search(t):
        found.append("change condition")
    return found


def suggested_reframing(aspect):
    """A TEMPLATE, never a finished question. The content is hers to supply."""
    a = (aspect or "the critical aspect").strip()
    return (
        'Rebuild move 1 as a choice, so the question carries the difference: '
        '"Which of these two — <case A> or <case B> — <the thing that differs>? '
        'Why did you choose?" '
        'Take <case A> and <case B> from this cycle\'s own Contrast Set. '
        'Aspect at stake: %s. '
        '(08: an aspect IS a difference. A question that names the thing and '
        'waits for the Contrast Set to supply the contrast is the polar bear '
        'question and produces silence. Accept, edit, or reject.)' % a
    )


# --------------------------------------------------------------------------
# inventory
# --------------------------------------------------------------------------


def load_inventory(path):
    with open(path, newline="", encoding="utf-8") as fh:
        return list(csv.DictReader(fh))


def arc_positions(rows):
    """cycle key -> (index, total, previous label, next label) across the arc."""
    seen = {}
    for r in rows:
        key = r["cycle"]
        # prefer the plainest deck name for the arc label
        label = re.sub(r"\.pptx$", "", r["file"])
        label = re.sub(r"^▶ LIVE — ", "", label)
        label = re.sub(r" — with Concept Bank$", "", label)
        if key not in seen or len(label) < len(seen[key]):
            seen[key] = label
    keys = sorted(seen, key=cycle_sort_key)
    out = {}
    for i, k in enumerate(keys):
        out[k] = {
            "index": i + 1,
            "total": len(keys),
            "prev": seen[keys[i - 1]] if i > 0 else None,
            "next": seen[keys[i + 1]] if i + 1 < len(keys) else None,
        }
    return out


def moves_present(row):
    """Which of the nine (+ what-if) this deck carries, from the inventory."""
    present = set()
    for label in [s.strip() for s in (row.get("slide_types_ordered") or "").split(">")]:
        mv = SLIDE_TYPE_TO_MOVE.get(label.lower())
        if mv:
            present.add(mv)
    if (row.get("continuity_question") or "0") not in ("0", ""):
        present.add("6")
    if (row.get("stock_and_flow") or "0") not in ("0", ""):
        present.add("7")
    for s in [x.strip().lower() for x in (row.get("advisory2_structure") or "").split(";")]:
        mv = STRUCTURE_TO_MOVE.get(s)
        if mv:
            present.add(mv)
    return present


def simultaneity(row, present):
    """Which simultaneity the cycle works on (09 §2)."""
    dia = (row.get("concept_bank") or "0") not in ("0", "")
    syn = bool(present & set(("7", "8", "9")))
    parts = []
    if syn:
        names = ", ".join(note_name(m) for m in ("7", "8", "9") if m in present)
        parts.append("Synchronic (fusion) — carried by: %s." % names)
    if dia:
        parts.append("Diachronic — carried by the Concept Bank.")
    if not parts:
        parts.append(
            "None declared. No Concept Bank and no coordination structure found, "
            "so this cycle currently works on separation only."
        )
    return parts


def move1_findings(row):
    """Per-slide move-1 verdicts, inventory verdict cross-checked against ours."""
    verdicts = {}
    for item in split_multi(row.get("advisory1_verdict")):
        m = re.match(r"^(S\d+)\s*:\s*(.*)$", item)
        if m:
            verdicts[m.group(1)] = m.group(2).strip()
    texts = {}
    for item in split_multi(row.get("ca_question_text_verbatim")):
        m = re.match(r'^(S\d+)\s*:\s*"?(.*?)"?$', item, re.S)
        if m:
            texts[m.group(1)] = m.group(2).strip()

    aspects = [a.strip() for a in (row.get("critical_aspects") or "").split("|") if a.strip()]
    out = []
    for i, sid in enumerate(sorted(verdicts, key=lambda s: int(s[1:]))):
        inv = verdicts[sid]
        text = texts.get(sid, "")
        ours = carries_difference(text)
        fails = (inv.upper() == NO_DIFFERENCE_VERDICT)
        aspect = aspects[i] if i < len(aspects) else (aspects[0] if aspects else "")
        out.append({
            "slide": sid,
            "inventory_verdict": inv,
            "our_verdict": ", ".join(ours) if ours else NO_DIFFERENCE_VERDICT,
            "agree": bool(ours) != fails,
            "fails": fails,
            "text": text,
            "aspect": aspect,
            "reframing": suggested_reframing(aspect),
        })
    if not verdicts and (row.get("ca_question_count") or "0") in ("0", ""):
        out.append({
            "slide": "-",
            "inventory_verdict": "",
            "our_verdict": "",
            "agree": True,
            "fails": True,
            "text": "",
            "aspect": aspects[0] if aspects else "",
            "reframing": "No move-1 slide found in this deck. The cycle has no "
                         "Critical Aspect question to carry a difference.",
            "missing": True,
        })
    return out


# --------------------------------------------------------------------------
# teacher note copy — the seven declarations
# --------------------------------------------------------------------------


def build_note_copy(row, arc):
    """Return (left_blocks, right_blocks, placeholder_count).

    A block is (label, [(text, kind), ...]) where kind is 'body' or 'ask'.
    """
    aspects = [a.strip() for a in (row.get("critical_aspects") or "").split("|") if a.strip()]
    present = moves_present(row)
    key = row["cycle"]
    pos = arc.get(key, {})

    # 1 — the critical aspect, in plain words
    if aspects:
        aspect_lines = [("%d. %s" % (i + 1, a), "body") for i, a in enumerate(aspects)]
    else:
        aspect_lines = [("No critical aspect found in this deck.", "body")]

    # 2 — what is held invariant. HERS.
    invariant = [
        (PLACEHOLDER + " what is held the same across every example here? "
         "Name it in one line.", "ask"),
        ("The examples in this cycle differ in one dimension on purpose.", "body"),
    ]

    # 3 — what breaks if an example is substituted. HERS.
    substitute = [
        (PLACEHOLDER + " swap one example for a better one - what breaks? "
         "Name the thing that stops working.", "ask"),
    ]

    # 4 — position in the sequence
    seq = []
    if pos:
        seq.append(("Cycle %s — %d of %d in the arc."
                    % (key, pos["index"], pos["total"]), "body"))
        seq.append(("Assumes: %s" % (pos["prev"] or "nothing before it."), "body"))
        seq.append(("Depended on by: %s" % (pos["next"] or "nothing after it."), "body"))
    core = row.get("core_questions_by_aspect") or ""
    counts = re.findall(r":\s*(\d/\d)\s*\[", core)
    if counts:
        seq.append(("5 core questions, per aspect - %s."
                    % ", ".join(counts), "body"))

    # 5 — the slide-type map
    have_core = [m for m in CORE_MOVES if m in present]
    miss_core = [m for m in CORE_MOVES if m not in present]
    have_cond = [m for m in CONDITIONAL_MOVES if m in present]
    miss_cond = [m for m in CONDITIONAL_MOVES if m not in present]
    smap = []
    smap.append(("5 core questions present: %s."
                 % (", ".join(note_name(m) for m in have_core) or "none"), "body"))
    if miss_core:
        smap.append(("Core questions MISSING: %s."
                     % ", ".join(note_name(m) for m in miss_core), "body"))
    smap.append(("Conditional structures present: %s."
                 % (", ".join(note_name(m) for m in have_cond) or "none"), "body"))
    smap.append(("Left out: %s — reason for each absence not recorded."
                 % (", ".join(note_name(m) for m in miss_cond) or "none"), "body"))
    smap.append(("Companion what-if slide - %s."
                 % ("present" if "W" in present else "absent"), "body"))

    # 6 — which simultaneity the cycle works on
    sim = [(t, "body") for t in simultaneity(row, present)]

    # 7 — the visibility rung. THE RUNG IS HERS; the ceilings are derived.
    vis = [
        (PLACEHOLDER + " which rung now? 1 teacher only / 2 assigned partner / "
         "3 you read a wrong-but-useful answer aloud, unattributed, and use it / "
         "4 owned aloud by invitation / 5 the vote.", "ask"),
    ]
    vis.append(("Reset low inside the cycle; raise it across the year.", "body"))

    ceilings = [("Maximum visibility per move, for the moves this deck "
                 "carries:", "body")]
    for num, _name, ceil in MOVES:
        if num in present:
            # never begin a line "what if" / "what-if": extract_and_grade.py
            # classifies on that at line start (§3c).
            name = ("Companion what-if" if num == "W"
                    else "%s %s" % (num, note_name(num)))
            ceilings.append(("%s - %s" % (name, ceil), "body"))

    left = [
        ("1. CRITICAL ASPECT", aspect_lines),
        ("2. HELD INVARIANT", invariant),
        ("3. IF AN EXAMPLE IS SUBSTITUTED", substitute),
        ("4. POSITION IN THE SEQUENCE", seq),
        ("7. VISIBILITY - THE RUNG", vis),
    ]
    right = [
        ("5. SLIDE-TYPE MAP", smap),
        ("6. SIMULTANEITY", sim),
        ("7. VISIBILITY - THE CEILINGS", ceilings),
    ]
    n_ask = sum(1 for _l, items in left + right for _t, k in items if k == "ask")
    return left, right, n_ask


# --------------------------------------------------------------------------
# rendering
# --------------------------------------------------------------------------


def _textbox(slide, geom, wrap=True):
    x, y, w, h = geom
    box = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    return box


def _run(para, text, size, bold, color):
    r = para.add_run()
    r.text = text
    f = r.font
    f.name = FONT
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = RGBColor.from_string(color)
    return r


def _est_lines(text, size, width_emu):
    """Rough wrap estimate: Arial averages ~0.5 em per character."""
    if not text:
        return 1
    chars_per_line = max(10, int((width_emu / 914400.0) * 72.0 / (size * 0.50)))
    return max(1, -(-len(text) // chars_per_line))


def _capacity(height_emu, size):
    return int(height_emu / (Pt(size).emu * 1.30))


def est_column(geom, blocks, size, spacer):
    """Estimated line count for a column, computed before any shape is made."""
    w = geom[2]
    used = 0
    first = True
    for label, items in blocks:
        if not first and spacer:
            used += 1
        first = False
        used += _est_lines(label, SZ_LABEL, w)
        for text, _kind in items:
            used += _est_lines(text, size, w)
    return used


def fit_columns(left, right, geom_l, geom_r):
    """Pick the largest legal body size at which both columns fit.

    Never shrinks below the smallest §13 body token. If nothing fits, the
    caller is told, and the change log carries the warning — the slide is not
    silently allowed to overflow.
    """
    for size in SZ_NOTE_STEPS:
        for spacer in (True, False):
            cap = _capacity(geom_l[3], size)
            ul = est_column(geom_l, left, size, spacer)
            ur = est_column(geom_r, right, size, spacer)
            if ul <= cap and ur <= cap:
                return size, spacer, ul, ur, cap, True
    size = SZ_NOTE_STEPS[-1]
    cap = _capacity(geom_l[3], size)
    ul = est_column(geom_l, left, size, False)
    ur = est_column(geom_r, right, size, False)
    return size, False, ul, ur, cap, False


def render_column(slide, geom, blocks, size, spacer):
    """Render label/body blocks into one column."""
    box = _textbox(slide, geom)
    tf = box.text_frame
    first = True
    for label, items in blocks:
        if not first and spacer:
            p = tf.add_paragraph()
            _run(p, "", size, False, BODY)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        _run(p, label, SZ_LABEL, True, TEAL)
        for text, kind in items:
            p = tf.add_paragraph()
            if kind == "ask":
                _run(p, text, size, True, ALERT)
            else:
                _run(p, text, size, False, BODY)
    return box


def build_teacher_note(prs, row, arc, layout, subhead_text):
    """Append the teacher note slide. Returns (slide, n_placeholders, warnings)."""
    left, right, n_ask = build_note_copy(row, arc)

    slide = blank_house_slide(prs, layout)

    b = _textbox(slide, G_KICKER)
    _run(b.text_frame.paragraphs[0], NOTE_KICKER, SZ_KICKER, True, MUTED)

    b = _textbox(slide, G_HEADING)
    _run(b.text_frame.paragraphs[0], NOTE_TITLE, SZ_HEADING, True, BODY)

    b = _textbox(slide, G_SUBHEAD)
    _run(b.text_frame.paragraphs[0], subhead_text, SZ_SUBHEAD, False, MUTED)

    warnings = []
    size, spacer, used_l, used_r, cap, ok = fit_columns(left, right, G_COL_L, G_COL_R)
    render_column(slide, G_COL_L, left, size, spacer)
    render_column(slide, G_COL_R, right, size, spacer)
    if size != SZ_NOTE_STEPS[0]:
        warnings.append(
            "teacher note body set to %dpt (from %dpt) so the seven declarations "
            "fit on one slide" % (size, SZ_NOTE_STEPS[0]))
    if not ok:
        warnings.append(
            "teacher note may overflow: columns estimated at %d and %d lines "
            "against a capacity of %d at %dpt — render it and look"
            % (used_l, used_r, cap, size))

    # §7 guard: a unique slide-type marker must never appear twice in a deck,
    # and no paragraph may begin "what if" (the grader's line-start fallback).
    txt = norm(slide_text(slide)).lower()
    for bad in FORBIDDEN_ON_NOTE:
        if bad in txt:
            warnings.append("teacher note carries the marker string %r — refuse" % bad)
    for line in slide_text(slide).split("\n"):
        if re.match(r"^\s*what[- ]if\b", norm(line), re.I):
            warnings.append(
                "teacher note has a line starting %r — refuse" % line[:40])
    return slide, n_ask, warnings


RE_SLIDE_PARTNAME = re.compile(r"^/ppt/slides/slide(\d+)\.xml$")


def _fix_slide_partname(prs, slide):
    """Give a newly added slide a partname nothing else in the package holds.

    python-pptx names the new part `slide<len(sldIdLst) + 1>.xml`. A deck that
    has had a slide dropped from the id list keeps the orphaned part in the
    package, so that count is LOWER than the highest index in use and the new
    part silently takes a name another part already has. The .pptx then saves
    with two zip entries of the same name — PowerPoint reads whichever it
    reaches first, which is the old slide, and the new one is simply gone.
    Cycle 12b is exactly that deck: two reference slides were dropped from the
    id list during the build and their parts stayed behind.

    Returns the new partname if it renamed, otherwise None.
    """
    target = str(slide.part.partname)
    used = set()
    seen = 0
    for part in prs.part.package.iter_parts():
        name = str(part.partname)
        if name == target:
            seen += 1
        m = RE_SLIDE_PARTNAME.match(name)
        if m:
            used.add(int(m.group(1)))
    if seen < 2:
        return None
    n = 1
    while n in used:
        n += 1
    free = PackURI("/ppt/slides/slide%d.xml" % n)
    slide.part.partname = free
    return str(free)


def blank_house_slide(prs, layout):
    """A slide on the deck's own layout with every placeholder stripped."""
    slide = prs.slides.add_slide(layout)
    _fix_slide_partname(prs, slide)
    for sh in list(slide.shapes):
        sh._element.getparent().remove(sh._element)
    return slide


def ladder_guard(slide):
    """Every reason this slide must not ship, as a list of warnings."""
    warnings = []
    txt = norm(slide_text(slide)).lower()
    for bad in FORBIDDEN_ON_LADDER:
        if bad in txt:
            warnings.append(
                "visibility ladder carries the marker string %r — refuse" % bad)
    for line in slide_text(slide).split("\n"):
        if re.match(r"^\s*what[- ]if\b", norm(line), re.I):
            warnings.append(
                "visibility ladder has a line starting %r — refuse" % line[:40])
    if "do not project" not in txt:
        warnings.append(
            "visibility ladder has lost its 'do not project' marker, so it would "
            "be counted as a diagnostic slide — refuse")
    return warnings


def build_visibility_ladder(prs, layout, note_slide_number=None):
    """Append the visibility-ladder slide. Returns (slide, warnings).

    One slide, teacher-facing, never projected. It explains the five rungs the
    teacher note's seventh declaration names. Deletable in one action, like the
    note — a teacher who already knows the ladder does not need it twice.
    """
    slide = blank_house_slide(prs, layout)

    b = _textbox(slide, G_KICKER)
    _run(b.text_frame.paragraphs[0], VL_KICKER, SZ_KICKER, True, MUTED)

    b = _textbox(slide, G_HEADING)
    _run(b.text_frame.paragraphs[0], VL_TITLE, SZ_HEADING, True, BODY)

    b = _textbox(slide, G_VL_SUBHEAD)
    sub = (VL_SUBHEAD % note_slide_number if note_slide_number
           else VL_SUBHEAD_NO_NOTE)
    _run(b.text_frame.paragraphs[0], sub, SZ_VL_BODY, False, MUTED)

    box = _textbox(slide, G_VL_BODY)
    tf = box.text_frame
    for i, (label, body) in enumerate(VL_BLOCKS):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if i < len(VL_BLOCKS) - 1:
            p.space_after = Pt(SZ_VL_SPACE_AFTER)
        _run(p, label, SZ_VL_BODY, True, TEAL)
        _run(p, body, SZ_VL_BODY, False, BODY)

    warnings = ladder_guard(slide)

    # Did somebody edit the copy? The slide is fixed, not per-deck, so it was
    # measured against real Arial metrics once (see SZ_VL_BODY) rather than
    # resized per deck. est_column() cannot be used here: it assumes a label on
    # its own line, and this slide runs the bold label and the body in ONE
    # paragraph, which is what keeps eight blocks on one slide.
    used = sum(_est_lines(l + b_, SZ_VL_BODY, G_VL_BODY[2])
               for l, b_ in VL_BLOCKS)
    if used > VL_EST_LINES:
        warnings.append(
            "the visibility-ladder copy has grown: %d estimated lines against "
            "the %d the measured copy scores. The slide fits at %dpt with "
            "18.5pt to spare; measure it against Arial again before it ships"
            % (used, VL_EST_LINES, SZ_VL_BODY))
    return slide, warnings


def move_slide_to(prs, index):
    """build_concept_banks.py's move: pull the appended slide into position."""
    lst = prs.slides._sldIdLst
    items = list(lst)
    el = items[-1]
    lst.remove(el)
    lst.insert(index, el)


# --------------------------------------------------------------------------
# detection — idempotency
# --------------------------------------------------------------------------


def find_teacher_note(prs):
    for i, s in enumerate(prs.slides):
        t = norm(slide_text(s)).lower()
        if "teacher note" in t and "design of this cycle" in t:
            return i
        if NOTE_MARKER.lower() in t and PLACEHOLDER.lower() in t:
            return i
    return None


def find_visibility_ladder(prs):
    """Index of the visibility-ladder slide, or None. Idempotency for change 4."""
    for i, s_ in enumerate(prs.slides):
        t = norm(slide_text(s_))
        # Both sides go through norm(): VL_TITLE carries an em-dash and norm()
        # folds it to a hyphen, so an un-normalised comparison never matches
        # and the insert stops being idempotent.
        if VL_MARKER in t and norm(VL_TITLE).lower() in t.lower():
            return i
    return None


def find_concept_bank(prs):
    for i, s in enumerate(prs.slides):
        t = norm(slide_text(s))
        if CONCEPT_BANK_HEADING.lower() in t.lower() or CONCEPT_BANK_KICKER in t:
            return i
    return None


RE_ALREADY_RELATES = re.compile(
    r"(how (they|two|these) connect|write one sentence saying how|"
    r"pick any two terms|how two of (them|the terms) connect|"
    r"say how .{0,30}\bconnect)", re.I)


def concept_bank_state(prs, idx):
    """'has_prompt' | 'terms_only'"""
    t = norm(slide_text(list(prs.slides)[idx]))
    if RE_ALREADY_RELATES.search(t):
        return "has_prompt"
    return "terms_only"


def teacher_front_matter_end(prs):
    """Index AFTER the leading run of teacher-reference slides."""
    last = -1
    for i, s in enumerate(prs.slides):
        t = norm(slide_text(s)).lower()
        if "teacher reference" in t or "teacher prep" in t or "do not project" in t:
            if i == last + 1:
                last = i
            else:
                break
        else:
            break
    return last + 1 if last >= 0 else 1


# --------------------------------------------------------------------------
# per-deck driver
# --------------------------------------------------------------------------

OUT_SUFFIX = " — with 2026-08-29 changes.pptx"
FORBIDDEN_WRITE_ROOTS = ("/Library/CloudStorage",)


def output_path(src, out_dir):
    base = os.path.basename(src)
    stem = re.sub(r"\.pptx$", "", base)
    name = stem + OUT_SUFFIX
    return os.path.join(out_dir or os.path.dirname(src), name)


def guard_write_path(path):
    real = os.path.realpath(path)
    for root in FORBIDDEN_WRITE_ROOTS:
        if root in real:
            raise RuntimeError("refusing to write inside %s: %s" % (root, real))
    return real


def process_deck(row, exports, arc, args):
    src = os.path.join(exports, row["folder"], row["file"])
    rec = {
        "cycle": row["cycle"],
        "file": row["file"],
        "path": src,
        "status": "planned",
        "applied": [],
        "skipped": [],
        "flagged": [],
        "warnings": [],
        "placeholders": 0,
        "output": "",
    }

    if not os.path.exists(src):
        rec["status"] = "missing"
        rec["skipped"].append("file not found on disk")
        return rec

    try:
        prs = Presentation(src)
    except Exception as exc:
        rec["status"] = "error"
        rec["skipped"].append("could not open: %s: %s" % (type(exc).__name__, exc))
        return rec

    blocking, strict, _detail = check_format_tokens(prs)
    if blocking:
        rec["status"] = "REFUSED"
        rec["skipped"] += ["format tokens (blocking): " + b for b in blocking]
        return rec
    if strict:
        if args.format_gate == "strict":
            rec["status"] = "REFUSED"
            rec["skipped"] += ["format tokens (§13): " + s for s in strict]
            rec["skipped"].append(
                "refused rather than patched. Re-run with --format-gate geometry "
                "to apply anyway, which records the debt instead of hiding it.")
            return rec
        rec["warnings"] += ["pre-existing format debt, NOT patched: " + s for s in strict]

    # ---- change 3: flag move 1 (never rewrite) -------------------------
    for f in move1_findings(row):
        if f["fails"]:
            if f.get("missing"):
                rec["flagged"].append({
                    "kind": "move-1 missing",
                    "slide": "-",
                    "verbatim": "",
                    "reframing": f["reframing"],
                    "note": "",
                })
            else:
                rec["flagged"].append({
                    "kind": "move-1 carries no difference",
                    "slide": f["slide"],
                    "verbatim": f["text"],
                    "reframing": f["reframing"],
                    "note": "inventory: %s / independent check: %s"
                            % (f["inventory_verdict"], f["our_verdict"]),
                })
        elif not f["agree"]:
            rec["warnings"].append(
                "move-1 %s: inventory says %r, independent check says %r"
                % (f["slide"], f["inventory_verdict"], f["our_verdict"]))

    changed = False

    # ---- change 2: relating prompt on the Concept Bank -----------------
    cb = find_concept_bank(prs)
    if cb is None:
        rec["skipped"].append(
            "relating prompt: no Concept Bank slide in this deck "
            "(build_concept_banks.py has to run first)")
    elif concept_bank_state(prs, cb) == "has_prompt":
        rec["skipped"].append(
            "relating prompt: already present on Concept Bank slide %d" % (cb + 1))
    else:
        s = list(prs.slides)[cb]
        box = _textbox(s, G_RELATE)
        _run(box.text_frame.paragraphs[0], RELATING_PROMPT, SZ_BODY, True, BODY)
        rec["applied"].append(
            "relating prompt appended to Concept Bank slide %d "
            "(co-presence is the precondition, not the achievement)" % (cb + 1))
        changed = True

    # ---- change 1: teacher note slide ----------------------------------
    existing = find_teacher_note(prs)
    if existing is not None:
        rec["skipped"].append(
            "teacher note: already present at slide %d" % (existing + 1))
    else:
        insert_at = teacher_front_matter_end(prs)
        layout = list(prs.slides)[min(insert_at, len(prs.slides._sldIdLst) - 1)].slide_layout
        # prefer the layout the Teacher Prep slide uses
        for s in prs.slides:
            t = norm(slide_text(s)).lower()
            if "teacher prep" in t:
                layout = s.slide_layout
                break
        label = re.sub(r"\.pptx$", "", row["file"])
        label = re.sub(r"^▶ LIVE — ", "", label)
        subhead = ("%s · one slide, not speaker notes. Delete it before you "
                   "project the deck." % label)
        _slide, n_ask, warns = build_teacher_note(prs, row, arc, layout, subhead)
        if any("refuse" in w for w in warns):
            rec["status"] = "REFUSED"
            rec["skipped"] += warns
            return rec
        move_slide_to(prs, insert_at)
        rec["applied"].append(
            "teacher note slide inserted at position %d, carrying the seven "
            "declarations; %d left for Katherine" % (insert_at + 1, n_ask))
        rec["placeholders"] = n_ask
        rec["warnings"] += warns
        changed = True

    # ---- change 4: the visibility-ladder slide -------------------------
    # It goes immediately AFTER the teacher note, never at slide 2. The note
    # declares the rung; this slide explains what a rung is. Read the other way
    # round it is a glossary in front of a term nobody has met yet — and it
    # would push the Teacher Prep slide off slide 2, which §2.0 fixes, and
    # break every "slide N" cross-reference the front block already carries.
    existing_vl = find_visibility_ladder(prs)
    if existing_vl is not None:
        rec["skipped"].append(
            "visibility ladder: already present at slide %d" % (existing_vl + 1))
    else:
        note_idx = find_teacher_note(prs)
        insert_at = (note_idx + 1) if note_idx is not None \
            else teacher_front_matter_end(prs)
        layout = list(prs.slides)[min(insert_at, len(prs.slides._sldIdLst) - 1)].slide_layout
        for s_ in prs.slides:
            if "teacher prep" in norm(slide_text(s_)).lower():
                layout = s_.slide_layout
                break
        _slide, warns = build_visibility_ladder(
            prs, layout, (note_idx + 1) if note_idx is not None else None)
        if any("refuse" in w for w in warns):
            rec["status"] = "REFUSED"
            rec["skipped"] += warns
            return rec
        move_slide_to(prs, insert_at)
        rec["applied"].append(
            "visibility-ladder slide inserted at position %d, immediately after "
            "the teacher note. Teacher-facing, marked 'do not project', and "
            "excluded from every diagnostic count" % (insert_at + 1))
        rec["warnings"] += warns
        if note_idx is None:
            rec["warnings"].append(
                "visibility ladder placed in the front block with no teacher "
                "note above it to explain — check the order by eye")
        else:
            rec["warnings"].append(
                "every 'slide N' cross-reference below slide %d in this deck is "
                "now off by one. deck_apply_changes does not renumber prose — "
                "check the teacher note, the image credits and the speaker notes"
                % (insert_at + 1))
        changed = True

    if not changed:
        rec["status"] = "no change needed"
        return rec

    out = output_path(src, args.out_dir)
    rec["output"] = out
    if os.path.realpath(out) == os.path.realpath(src):
        rec["status"] = "REFUSED"
        rec["skipped"].append("output path equals input path")
        return rec

    if not args.write:
        rec["status"] = "would change (dry run)"
        return rec

    guard_write_path(out)
    if os.path.exists(out) and not args.force:
        rec["status"] = "REFUSED"
        rec["skipped"].append("output already exists: %s (use --force)" % out)
        return rec
    d = os.path.dirname(out)
    if d and not os.path.isdir(d):
        os.makedirs(d)
    prs.save(out)
    rec["status"] = "WRITTEN"
    return rec


# --------------------------------------------------------------------------
# reporting
# --------------------------------------------------------------------------


def write_report(records, audits, args, stamp):
    mode = "write" if args.write else "dryrun"
    md = os.path.join(audits, "deck_apply_changes_%s_%s.md" % (stamp, mode))
    csvp = os.path.join(audits, "deck_apply_changes_%s_%s.csv" % (stamp, mode))

    counts = {}
    for r in records:
        counts[r["status"]] = counts.get(r["status"], 0) + 1

    lines = []
    lines.append("# deck_apply_changes — %s run, %s" % (mode, stamp))
    lines.append("")
    lines.append("Applies the 2026-08-29 changes. Input decks are never modified; "
                 "every change is written to a new file alongside the original.")
    lines.append("")
    lines.append("- decks considered: %d" % len(records))
    for k in sorted(counts):
        lines.append("- %s: %d" % (k, counts[k]))
    lines.append("- format gate: %s" % args.format_gate)
    lines.append("")
    lines.append("## What this tool applies, flags, and refuses")
    lines.append("")
    lines.append("APPLIES  the teacher note slide, the visibility-ladder slide "
                 "immediately after it, and the relating prompt on a Concept "
                 "Bank that only lists terms.")
    lines.append("FLAGS    a move-1 Critical Aspect question carrying no "
                 "difference. The slide is not touched.")
    lines.append("REFUSES  a deck that fails the §13 format tokens, or whose "
                 "canvas is not 4:3.")
    lines.append("NEVER    invents what is held invariant, what breaks if an "
                 "example is substituted, or the visibility rung.")
    lines.append("")

    for r in records:
        lines.append("## Cycle %s — %s" % (r["cycle"], r["file"]))
        lines.append("")
        lines.append("- status: **%s**" % r["status"])
        lines.append("- input: `%s`" % r["path"])
        if r["output"]:
            lines.append("- output: `%s`" % r["output"])
        for a in r["applied"]:
            lines.append("- CHANGED: %s" % a)
        for s in r["skipped"]:
            lines.append("- SKIPPED: %s" % s)
        for w in r["warnings"]:
            lines.append("- WARNING: %s" % w)
        if r["placeholders"]:
            lines.append("- FOR KATHERINE: %d placeholder(s) on the teacher note, "
                         "in %s. The linter can find them on the string `%s`."
                         % (r["placeholders"], ALERT, PLACEHOLDER))
        for f in r["flagged"]:
            lines.append("- FLAGGED (no change made): %s, slide %s"
                         % (f["kind"], f["slide"]))
            if f["note"]:
                lines.append("    - %s" % f["note"])
            if f["verbatim"]:
                lines.append("    - verbatim: %s" % f["verbatim"])
            lines.append("    - suggested reframing (accept or reject): %s"
                         % f["reframing"])
        lines.append("")

    with open(md, "w", encoding="utf-8") as fh:
        fh.write("\n".join(lines) + "\n")

    with open(csvp, "w", newline="", encoding="utf-8") as fh:
        w = csv.writer(fh)
        w.writerow(["cycle", "file", "status", "applied", "skipped", "warnings",
                    "placeholders", "flagged_slides", "output"])
        for r in records:
            w.writerow([
                r["cycle"], r["file"], r["status"],
                " ;; ".join(r["applied"]),
                " ;; ".join(r["skipped"]),
                " ;; ".join(r["warnings"]),
                r["placeholders"],
                " ;; ".join("%s %s" % (f["slide"], f["kind"]) for f in r["flagged"]),
                r["output"],
            ])
    return md, csvp


# --------------------------------------------------------------------------


def main(argv=None):
    home = os.path.expanduser("~")
    ap = argparse.ArgumentParser(
        description=__doc__.split("\n")[0],
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="--dry-run is the default. Nothing is written to a deck without --write.")
    ap.add_argument("--inventory",
                    default=os.path.join(home, "code/answerable-skills/audits/"
                                               "deck_inventory_2026-08-29.csv"))
    ap.add_argument("--exports", default=os.path.join(home, "deck_work/exports"))
    ap.add_argument("--audits",
                    default=os.path.join(home, "code/answerable-skills/audits"))
    ap.add_argument("--out-dir", default=None,
                    help="where the new .pptx files go. Default: alongside the input.")
    ap.add_argument("--decks", nargs="*", default=None,
                    help="cycle keys to process, e.g. 10 07a. Default: all.")
    ap.add_argument("--match", default=None,
                    help="substring of the filename to process.")
    ap.add_argument("--write", action="store_true",
                    help="actually write. Without it this is a dry run.")
    ap.add_argument("--dry-run", action="store_true", default=True,
                    help=argparse.SUPPRESS)
    ap.add_argument("--force", action="store_true",
                    help="overwrite an existing OUTPUT file. Never an input.")
    ap.add_argument("--format-gate", choices=("strict", "geometry"), default="strict",
                    help="strict (default): refuse any deck failing the §13 format "
                         "tokens. geometry: refuse only a deck whose canvas or font "
                         "floor makes insertion unsafe, and record the rest as debt.")
    ap.add_argument("--no-report", action="store_true")
    args = ap.parse_args(argv)

    rows = load_inventory(args.inventory)
    arc = arc_positions(rows)

    on_disk = set(os.path.realpath(p)
                  for p in glob.glob(os.path.join(args.exports, "*", "*.pptx")))
    known = set()

    records = []
    for row in rows:
        if args.decks and row["cycle"] not in args.decks:
            continue
        if args.match and args.match not in row["file"]:
            continue
        if OUT_SUFFIX.strip(" .pptx") in row["file"]:
            continue
        known.add(os.path.realpath(
            os.path.join(args.exports, row["folder"], row["file"])))
        rec = process_deck(row, args.exports, arc, args)
        records.append(rec)
        print("%-4s %-22s %s" % (rec["cycle"], rec["status"], rec["file"]))
        for a in rec["applied"]:
            print("       + %s" % a)
        for s in rec["skipped"]:
            print("       - %s" % s)
        for w in rec["warnings"]:
            print("       ! %s" % w)
        for f in rec["flagged"]:
            print("       FLAG %s %s" % (f["slide"], f["kind"]))

    stray = sorted(p for p in on_disk - known
                   if OUT_SUFFIX not in p and "with 2026-08-29" not in p)
    if stray and not args.decks and not args.match:
        print("\n%d .pptx on disk are not in the inventory and were not "
              "considered:" % len(stray))
        for p in stray:
            print("   %s" % p)

    counts = {}
    for r in records:
        counts[r["status"]] = counts.get(r["status"], 0) + 1
    print("\n%d decks considered. " % len(records) +
          ", ".join("%s=%d" % (k, counts[k]) for k in sorted(counts)))
    print("teacher notes planned/inserted: %d" %
          sum(1 for r in records if any("teacher note slide inserted" in a
                                        for a in r["applied"])))
    print("visibility-ladder slides planned/inserted: %d" %
          sum(1 for r in records if any("visibility-ladder slide inserted" in a
                                        for a in r["applied"])))
    print("relating prompts planned/appended: %d" %
          sum(1 for r in records if any("relating prompt appended" in a
                                        for a in r["applied"])))
    print("decks with a move-1 flag: %d" % sum(1 for r in records if r["flagged"]))
    if not args.write:
        print("\nDRY RUN — no deck was written. Add --write to apply.")

    if not args.no_report:
        stamp = datetime.date.today().isoformat()
        md, csvp = write_report(records, args.audits, args, stamp)
        print("report: %s" % md)
        print("report: %s" % csvp)
    return 0


if __name__ == "__main__":
    sys.exit(main())

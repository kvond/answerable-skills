#!/usr/bin/env python3
"""
deck_lint.py — VT deck validator, two tiers (rewritten 2026-08-29)
==================================================================
Run as the LAST gate of every VT deck build.

WHY TWO TIERS
-------------
The 2026-08-29 spec (`deck_work_order_of_operations.md` job 2, `coordination_judgment_and_the_package.md`
§12) splits this script's findings in two:

  HARD      the deck does not ship.  Exit code 1.
  ADVISORY  reported, does not fail the deck.  Exit code 0 on its own.

"A linter that fails every deck gets ignored inside a week." Most decks will not
carry a coordination structure or a teacher note the first time this runs, and
that is not a defect yet — it is work not done. Those live in ADVISORY.

HARD checks
-----------
  H-FORMAT-SIZE     slide size is 4:3, 9144000 x 6858000 EMU            (SKILL §13.1)
  H-FORMAT-FONT     every text run is Arial                             (SKILL §13.1)
  H-FORMAT-HEX      no forbidden palette hex on a text run:
                      000000  body black is 111111, never 000000        (SKILL §13.2)
                      0000FF  theme hyperlink blue — the hlinkClr fault (SKILL §3b)
  H-FORMAT-SCALE    no text run above 26pt, the day-divider headline
                    ceiling; 32pt Calibri is named a defect             (SKILL §13.3/§13.4)
  H-TIER-VERBATIM   "Getting Started" / "Working On It" / "Mastery"
                    each appear verbatim somewhere in the deck
  H-TIER-BLACK      a tier heading rendered 000000 instead of its hex
  H-MARKER-BROKEN   a marker string broken by typographic punctuation —
                    the ellipsis fault                                  (SKILL §3c)
  H-LINKS-SLIDE     links slide present, carrying well-formed URLs      (SKILL §8)
  H-TEACHER-PREP    "Additional Resources — Teacher Prep" slide present (SKILL §2.0)
  H-CONCEPT-BANK    Concept Bank present, above the Day 3 divider       (SKILL §2.3, §3b)
  H-TEACHER-NOTE    teacher note slide present — one slide, not speaker
                    notes, deletable in one action                      (spec §"Teacher note slide")
  H-VOCAB           the forbidden structural noun does not appear, in its
                    STRUCTURAL-NOUN sense only. The ordinary English verb
                    passes clean (Katherine, 2026-08-29).             (SKILL §1)

ADVISORY checks
---------------
  A-MOVE1-NO-DIFF   a Critical Aspect question that names the aspect and waits for
                    the Contrast Set to supply the difference. Structurally "how are
                    polar bears adapted?" and it produces silence. The question text
                    must itself carry a comparison, a choice between named
                    alternatives, or a stated change condition.
  A-NO-FUSION       no declared coordination structure on any critical aspect
  A-CA-NOT-DECLARED critical aspect not declared in the teacher note
  A-COND-NO-REASON  conditional slide types present with no stated reason
  A-CORE-DROPPED    a core question absent and the drop not declared. Dropping a
                    core question the content will not support is legitimate; not
                    saying so is what this catches.
  A-TN-ITEMS        teacher note does not declare all seven required items
                    (itemised — names which of the seven are missing)
  A-VISIBILITY      visibility rung not declared
  A-BANK-NO-RELATE  Concept Bank lists terms but never asks the student to state
                    how two of them connect. Co-presence without the relating.
  A-WHATIF          the cycle carries a coordination structure but its what-if is
                    missing, or is not marked individual and written
  A-PALETTE-DRIFT   text hexes outside the §13.2 palette
  A-TIER-COLOR      a tier heading with no explicit colour (theme/auto)
  A-MARKER-ABSENT   a diagnostic marker not found anywhere in the deck

EXCLUSIONS
----------
Teacher-navigation slides ("do not project", "teacher navigation") and the standing
end-of-lesson reflections ("continuation question:", "relates to me:") are excluded
from the diagnostic counts, exactly as in `extract_and_grade.py`. The two tuples
below MUST stay byte-identical to that script's NON_DIAGNOSTIC_MARKERS and
STANDING_REFLECTION_MARKERS. That grader is retired (see the tuple's comment), so
this file is now the sole owner of both tuples. See SKILL §7.1.

USAGE
-----
  python3 deck_lint.py <deck.pptx> [more decks | directories | globs ...]
  python3 deck_lint.py '~/deck_work/exports/*/*.pptx' --tier hard
  python3 deck_lint.py <dir> --json
  python3 deck_lint.py <dir> --csv audits/deck_lint.csv

  --tier hard|advisory|all   filter what is printed (default all).
                             Filtering never changes the exit code.
  --json                     machine-readable findings on stdout
  --csv <path>               one summary row per deck

Exit code: 1 if any HARD finding fired on any deck, else 0.
Dependencies: python-pptx only.
"""
from __future__ import print_function

import argparse
import csv as csvmod
import glob as globmod
import json
import os
import re
import sys

from pptx import Presentation
from pptx.util import Pt  # noqa: F401  (kept: callers import Pt from here)


# ---------------------------------------------------------------------------
# Constants that must not drift
# ---------------------------------------------------------------------------

TIER_LABELS = ("Getting Started", "Working On It", "Mastery")

# Tier heading colours. "Getting Started" and "Mastery" carry a colour; SKILL
# §13.2 states plainly that "Working On It" is body black 111111 and not a
# colour. The reference build (Cycle 02) nevertheless ships EFDF85 on that
# label, and the original linter required EFDF85. Both are accepted so the
# check keeps its teeth on the two labels that carry real colour without
# firing on every deck in the arc over a token the spec itself disputes.
TIER_HEX = {
    "Getting Started": "C0392B",
    "Working On It":   "EFDF85",   # historical linter value; see TIER_HEX_OK
    "Mastery":         "1E8449",
}
TIER_HEX_OK = {
    "Getting Started": {"C0392B"},
    "Working On It":   {"111111", "EFDF85"},
    "Mastery":         {"1E8449"},
}

# MUST stay in sync with extract_and_grade.py NON_DIAGNOSTIC_MARKERS.
NON_DIAGNOSTIC = ("do not project", "teacher navigation")
# Formerly had to stay in sync with extract_and_grade.py's
# STANDING_REFLECTION_MARKERS. That sync rule is RETIRED: Katherine confirmed
# 2026-08-29 that student feedback now runs through a prompt, not that script,
# so this tuple no longer has a second half to disagree with. SKILL §7.1's
# "edit two scripts in one pass" rule and its §987 challenge are both answered
# by that — this file is now the only place the tuple lives.
#
# "concept bank" added the same day, per SKILL §3b. It is defensive rather than
# a repair: read against the current diagnostic tests the Concept Bank already
# falls through to `other`, so the count does not move. The exclusion makes it
# a stated decision instead of an accident, and holds if the bank ever gains
# text that trips a classifier.
STANDING_REFLECTION = ("continuation question:", "relates to me:", "concept bank")

# The five diagnostic markers the teacher prompts match on. Four are regexes
# over the deck text; the fifth is the 3-tier block, detected by all three tier
# labels being present on one slide.
MARKERS = {
    "Critical aspect":  re.compile(r"critical aspect\s*:", re.I),
    "Pattern break":    re.compile(r"pattern break", re.I),
    "Build-a-rule":     re.compile(r"finish this sentence as a rule", re.I),
    "What-if":          re.compile(r"what if\?", re.I),
}
MARKER_ORDER = ("Critical aspect", "Pattern break", "Build-a-rule", "What-if")

# The §13.2 palette, plus the tokens the reference build actually ships.
PALETTE = {
    "028090",  # teal
    "111111",  # body black
    "666666",  # muted gray
    "C0392B",  # alert red
    "1E8449",  # green
    "F2F6F9",  # box fill
    "CCCCCC",  # box border
    "F5F5F5",  # panel gray
    "EAF6F7",  # divider pale
    "FFFFFF",  # divider white
    "EFDF85",  # "Working On It" as shipped by the reference build
}
# Hexes that are not drift but a named defect.
FORBIDDEN_HEX = {
    "000000": "body black is 111111, never 000000 (SKILL §13.2)",
    "0000FF": "theme hyperlink blue — the hlinkClr extension was dropped (SKILL §3b)",
}

SLIDE_W, SLIDE_H = 9144000, 6858000     # 10 x 7.5 in, 4:3 (SKILL §13.1)
MAX_RUN_PT = 26                          # day-divider headline (SKILL §13.3)

# The forbidden structural noun (SKILL §1). It appears in this file only as the
# pattern that finds it. Never write it into a deck, a manual, sales copy, an
# output string, or a comment.
# Only the STRUCTURAL-NOUN sense is a defect. The same token as an ordinary
# English verb ("small steps beat one big burst", "terms used in reasoning beat
# terms recited") is ordinary prose and passes clean, by Katherine's direction
# 2026-08-29. So the rule matches the jargon sense directly rather than matching
# every occurrence and trying to subtract the verb — the earlier shape failed
# 2 of 4 decks on their own good prose.
_JARGON_HEAD = r"legend|map|structure|number|numbers|name|names|vocabulary|sequence|list|chart|type|types"
_JARGON_MOD = r"core|conditional|standing|evidence|first|second|third|fourth|fifth|next|last|final|whole|per|each"
FORBIDDEN_VOCAB = re.compile(
    r"(?:\bbeats?\s+(?:%s)\b"          # "beat legend", "beat structure"
    r"|\bbeat\s+\d"                     # "Beat 1A", "Beat 3"
    r"|\b(?:%s)\s+beats?\b"            # "core beats", "the evidence beat"
    r"|\b(?:a|an|the|this|that|one)\s+beat\b"   # "the beat", "a beat"
    r"|\bbeats?\s+(?:is|are|was|were)\b"        # "the beat is"
    r"|\b[a-z]+(?:-[a-z]+)+\s+beats?\b"          # "the case-and-vote beat"
    r")" % (_JARGON_HEAD, _JARGON_MOD),
    re.I,
)

CA_LABEL = re.compile(r"critical aspect\s*:\s*([^\n]+)", re.I)
URL_RX = re.compile(r"https?://[^\s<>\"')\]]+", re.I)
DAY_DIVIDER_RX = re.compile(r"\bday\s+(\d+)\s+of\s+(\d+)\b", re.I)

# Typographic characters that silently break a marker string (SKILL §3c).
TYPO_FIXES = [
    ("…", "?"),   # ellipsis standing in for a question mark
    ("’", "'"), ("‘", "'"),
    ("“", '"'), ("”", '"'),
    (" ", " "),  (" ", " "), (" ", " "),
    ("–", "-"),  ("—", "-"), ("−", "-"),
]


# ---------------------------------------------------------------------------
# The move-1 difference test
# ---------------------------------------------------------------------------
# An aspect IS a difference. A Critical Aspect question that names the aspect and
# waits for the Contrast Set on the next slide cannot be answered — it is
# structurally "how are polar bears adapted?" and the room goes quiet. The
# question text must itself carry one of three devices.

_COMPARATIVE = (
    r"\b(?:more|less|fewer|greater|higher|lower|bigger|smaller|larger|faster|slower|"
    r"stronger|weaker|shorter|better|worse|harder|easier|deeper|thicker|thinner|"
    r"denser|hotter|colder|heavier|lighter|richer|poorer)\b"
)
COMPARISON_PATTERNS = [
    _COMPARATIVE,
    r"(?<!no )\blonger\b",
    r"\bthan\b",
    r"\bdiffer(?:s|ent|ently|ence|ences)?\b",
    r"\b(?:same|identical|alike|similar|in common)\b",
    r"\bboth\b",
    r"\bneither\b",
    r"\b(?:unlike|whereas|instead of|rather than|versus|vs\.?)\b",
    r"\bcompare[ds]?\b|\bcomparison\b",
    r"\bopposite\b",
    r"\bcan\b[^?.]{0,80}\bcannot\b|\bcannot\b[^?.]{0,80}\bcan\b",
    r"\bone\b[^?.]{0,40}\bthe other\b",
    r"\bonly one\b",
    r"\bnot the other\b",
    r"\bbut\b",
]
NAMED_ALT_PATTERNS = [
    r"\bwhich\s+(?:one|of|two|is|are|would|do|does|did|will|has|have|makes|kind|type)\b",
    r"\beither\b[^?.]{0,80}\bor\b",
    r"[^.?]{0,120}\bor\b[^.?]{0,120}\?",     # a question that offers an alternative
    r"\bA\s+or\s+B\b",
]
CHANGE_COND_PATTERNS = [
    r"\bif\b",
    r"\bwhat[- ]if\b",
    r"\bsuppose\b",
    r"\bwould happen\b",
    r"\bwhen\b[^?.]{0,80}\b(?:chang|increas|decreas|add|remov|los|gain|break|breaks|doubl|halv)",
    r"\bas\b[^?.]{0,60}\b(?:increases|decreases|rises|falls|grows|drops|warms|cools)\b",
    r"\ba change in\b|\bchanges? (?:from|to|one|the|a|its)\b",
    r"\bafter\b[^?.]{0,60}\bchang",
    r"\bstops?\b[^?.]{0,60}\bworking\b",
]


def _any_match(patterns, text):
    for p in patterns:
        if re.search(p, text, re.I):
            return True
    return False


def difference_devices(text):
    """Which difference devices a Critical Aspect question carries. Empty = none."""
    found = []
    if _any_match(COMPARISON_PATTERNS, text):
        found.append("comparison")
    if _any_match(NAMED_ALT_PATTERNS, text):
        found.append("named alternatives")
    if _any_match(CHANGE_COND_PATTERNS, text):
        found.append("change condition")
    return found


# ---------------------------------------------------------------------------
# The seven declarations a teacher note slide must carry
# ---------------------------------------------------------------------------

TEACHER_NOTE_MARKERS = ("teacher note", "design note", "note to the teacher")

TEACHER_NOTE_ITEMS = [
    ("1 critical aspect",
     lambda t: "critical aspect" in t),
    ("2 what is held invariant, and that examples differ in one dimension on purpose",
     lambda t: (("held invariant" in t or "held constant" in t or "kept invariant" in t)
                and ("one dimension" in t or "differ in one" in t or "on purpose" in t))),
    ("3 what breaks if an example is substituted",
     lambda t: ("substitut" in t or "swap" in t or "replace" in t) and ("break" in t or "destroy" in t or "lose" in t)),
    ("4 position in the sequence — what this assumes, what depends on it",
     lambda t: ("assumes" in t or "depends on" in t or "position in the sequence" in t
                or "comes before" in t or "comes after" in t)),
    ("5 slide-type map — which slide types appear, which conditional ones were left out and why",
     lambda t: ("slide type" in t or "slide-type" in t)
               and ("left out" in t or "omitted" in t or "not used" in t or "because" in t or "why" in t)),
    ("6 which simultaneity the cycle works on — synchronic or diachronic",
     lambda t: "synchronic" in t or "diachronic" in t or "simultaneity" in t),
    ("7 the visibility rung",
     lambda t: "visibility" in t or "rung" in t),
]

# Rung language, for the deck-wide visibility check.
# The visibility-ladder slide (added 2026-08-29). The teacher note DECLARES the
# cycle's rung; this slide EXPLAINS all five, so the declaration means something
# to a teacher who has never seen the ladder. Uppercase kicker on the model of
# CONCEPT BANK and TEACHER NOTE; matched case-insensitively here, as everything
# else in this classifier is.
#
# It is tested AFTER the teacher note on purpose. A teacher note whose seventh
# declaration happens to say "the visibility ladder" in prose is still a teacher
# note, and a deck that loses its note to this branch would fail hard for a
# reason nobody would find. The ladder slide carries no teacher-note marker, so
# nothing is lost by letting the note win.
VISIBILITY_LADDER_MARKERS = ("visibility ladder",)

VISIBILITY_PATTERNS = [
    r"\bvisibility\b",
    # "rung" only in the ladder-of-exposure sense. A DNA deck says "each rung
    # of the ladder has to be one wide base plus one narrow base" and that is
    # not a visibility declaration.
    r"visibility[^.\n]{0,24}rung|rung\s*[1-5]\b",
    r"written and private",
    r"seen only by (?:the )?teacher",
    r"assigned partner",
    r"read (?:it )?aloud,? unattributed|unattributed read[- ]aloud",
    r"owns? an answer aloud",
    r"public simultaneous commitment",
]

# ---------------------------------------------------------------------------
# Coordination structures (spec "Kind distinction")
# ---------------------------------------------------------------------------
# Stock-and-flow is a REPRESENTATION and is detectable from the slide itself.
# Compensatory pair and conflict case are a case set and a case: they have to be
# FOUND in the content, and the spec asks whether they are DECLARED. Only an
# explicit naming counts as declared; keyword evidence is reported as a lead for
# the retrofit skill, never as a pass.

STOCK_AND_FLOW_CUES = [
    r"stock[- ]and[- ]flow", r"stock and flow",
    r"drag the slider", r"tap to open",
    r"answerable-biology-models",
]
DECLARED_STRUCTURE = {
    "stock-and-flow model": [r"stock[- ]and[- ]flow", r"stock and flow"],
    "compensatory pair":    [r"compensatory"],
    "conflict case":        [r"conflict case"],
}
CANDIDATE_CONFLICT = [
    r"who benefits[^.?]{0,60}who[^.?]{0,20}harmed",
    r"pulls? (?:in )?(?:two|opposite)",
    r"points one way[^.?]{0,60}points the other",
    r"at the same time[^.?]{0,40}(?:but|yet)",
]
CANDIDATE_COMPENSATORY = [
    r"both[^.?]{0,80}(?:steady|stable|the same outcome|balance out)",
    r"opposite[^.?]{0,60}same (?:result|outcome|number)",
]

CONDITIONAL_TYPES = ("Continuity question", "Stock-and-flow model")

# The 5 core questions, in their fixed order. A cycle normally carries all five
# once per critical aspect — but a core question whose content is not there gets
# DROPPED (Katherine, 2026-08-29). A drop is a legitimate design decision, so it
# is never a hard failure and never an advisory on its own. What it must not be
# is silent: the same rule that governs a conditional slide type governs this,
# so the teacher note's slide-type map has to say which core question was
# dropped and why. Undeclared, it is indistinguishable from an oversight.
# NB: only four of the five are slide-type labels. classify_slide() has no
# "Critical Aspect question" type — move 1 is found by the "Critical aspect:"
# marker instead (see the move1 dict below), so it is tested separately.
CORE_TYPES = ("Contrast Set", "Build a Rule", "Pattern Break", "3-Tier Question")


# ---------------------------------------------------------------------------
# Slide text and formatting extraction
# ---------------------------------------------------------------------------

def _slide_text(slide):
    out = []
    for sh in slide.shapes:
        if sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                t = "".join(r.text for r in p.runs)
                if t.strip():
                    out.append(t)
    return "\n".join(out)


def _run_hex(run):
    """Return uppercase RRGGBB for a run's font colour, or None if unset/theme."""
    try:
        c = run.font.color
        if c and c.type is not None and c.rgb is not None:
            return str(c.rgb).upper()
    except Exception:
        pass
    return None


def _iter_runs(slide):
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                if r.text.strip():
                    yield r


def _slide_urls(slide):
    """Every URL on the slide: printed in text, and behind a hyperlink."""
    urls = set()
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                urls.update(URL_RX.findall(r.text))
                try:
                    addr = r.hyperlink.address
                except Exception:
                    addr = None
                if addr:
                    urls.add(addr)
        try:
            if sh.click_action and sh.click_action.hyperlink.address:
                urls.add(sh.click_action.hyperlink.address)
        except Exception:
            pass
    return urls


def _normalize_typography(text):
    for bad, good in TYPO_FIXES:
        text = text.replace(bad, good)
    return text


# ---------------------------------------------------------------------------
# Slide classification
# ---------------------------------------------------------------------------

def classify_slide(text):
    """Slide type, by the marker strings the teacher prompts already match on.

    Order matters: the exclusion markers are tested first, exactly as
    extract_and_grade.classify_slide() does.
    """
    low = text.lower()
    first = text.split("\n")[0].strip().lower() if text.strip() else ""

    # The teacher note is tested FIRST. It is a teacher slide and will normally
    # carry "do not project", so testing the nav markers ahead of it would
    # classify every teacher note as navigation. This ordering governs slide
    # TYPE only — the diagnostic-count exclusion in is_excluded() keeps the
    # grader's order untouched.
    if any(m in low for m in TEACHER_NOTE_MARKERS):
        return "Teacher note"
    # Before the nav markers, for the same reason the note is: this slide
    # carries "do not project" and would otherwise be filed as navigation.
    if any(m in low for m in VISIBILITY_LADDER_MARKERS):
        return "Visibility ladder"
    if any(m in low for m in NON_DIAGNOSTIC):
        return "Teacher navigation"
    if "teacher reference" in low and "essential claim" in low:
        return "Teacher reference"
    if "teacher prep" in low or ("additional resources" in low and "buy ahead" in low):
        return "Teacher Prep"
    if DAY_DIVIDER_RX.search(low):
        return "Day divider"
    if "define these in your own words" in low or "concept bank" in low:
        return "Concept Bank"
    if "image credits" in low:
        return "Image credits"
    if "activity and resource links" in low or first.startswith("links"):
        return "Links slide"
    if "think → write → submit" in low or "think -> write -> submit" in low:
        return "Closing checklist"
    if "relates to me" in low:
        return "Relates to me"
    if "optional challenge" in low:
        return "Optional challenge"
    if "what if?" in low or re.search(r"(?m)^\s*what[- ]if\b", low):
        return "What if"
    if all(lbl.lower() in low for lbl in TIER_LABELS):
        return "3-Tier Question"
    if "pattern break" in low:
        return "Pattern Break"
    if ("finish this sentence as a rule" in low or "build a rule from" in low
            or "finish this rule" in low):
        return "Build a Rule"
    if "case a" in low and "case b" in low:
        return "Contrast Set"
    if ("continuation question" in low or "no answer to look up" in low
            or first.startswith("keep going")):
        return "Continuity question"
    if (_any_match(STOCK_AND_FLOW_CUES, low)
            or first.startswith("explore —") or first.startswith("explore -")):
        return "Stock-and-flow model"
    if "your first answer" in low and "your revised answer" in low:
        return "Response slide"
    if "bellringer" in low:
        return "Bellringer"
    if re.search(r"activity\s*\(", low) or first.startswith("activity"):
        return "Activity"
    return "Content/Explanation"


def is_excluded(text):
    """Nav slide or standing reflection — excluded from diagnostic counts."""
    low = text.lower()
    return (any(m in low for m in NON_DIAGNOSTIC)
            or any(m in low for m in STANDING_REFLECTION))


# ---------------------------------------------------------------------------
# Findings
# ---------------------------------------------------------------------------

HARD = "HARD"
ADVISORY = "ADVISORY"


class Finding(object):
    __slots__ = ("tier", "code", "message", "slide")

    def __init__(self, tier, code, message, slide=None):
        self.tier = tier
        self.code = code
        self.message = message
        self.slide = slide

    def as_dict(self):
        return {"tier": self.tier, "code": self.code,
                "slide": self.slide, "message": self.message}

    def line(self):
        where = "slide %s: " % self.slide if self.slide else ""
        return "[%s] %s%s" % (self.code, where, self.message)


# ---------------------------------------------------------------------------
# The linter
# ---------------------------------------------------------------------------

def lint_deck(path):
    """Return a result dict for one deck. Never raises on deck content."""
    result = {
        "file": os.path.basename(path),
        "path": os.path.abspath(path),
        "folder": os.path.basename(os.path.dirname(os.path.abspath(path))),
        "slides": 0,
        "findings": [],
        "hard": 0,
        "advisory": 0,
        "error": "",
    }
    findings = result["findings"]

    def hard(code, msg, slide=None):
        findings.append(Finding(HARD, code, msg, slide))

    def adv(code, msg, slide=None):
        findings.append(Finding(ADVISORY, code, msg, slide))

    try:
        prs = Presentation(path)
    except Exception as exc:                       # unreadable file
        result["error"] = "%s: %s" % (type(exc).__name__, exc)
        hard("H-UNREADABLE", "deck could not be opened: %s" % result["error"])
        result["hard"] = 1
        return result

    slides = list(prs.slides)
    result["slides"] = len(slides)
    texts = [_slide_text(s) for s in slides]
    types = [classify_slide(t) for t in texts]
    all_text = "\n".join(texts)
    low_all = all_text.lower()

    # -- format tokens ------------------------------------------------------
    try:
        w, h = int(prs.slide_width), int(prs.slide_height)
    except Exception:
        w = h = 0
    result["dims"] = "%dx%d" % (w, h)
    if (w, h) != (SLIDE_W, SLIDE_H):
        hard("H-FORMAT-SIZE",
             "slide size %dx%d EMU is not 4:3 — must be %dx%d (SKILL §13.1)"
             % (w, h, SLIDE_W, SLIDE_H))

    fonts, off_fonts = set(), set()
    hex_sightings = {}          # HEX -> first slide number
    oversize = []
    for idx, slide in enumerate(slides, start=1):
        for run in _iter_runs(slide):
            name = run.font.name
            if name:
                fonts.add(name)
                if name.strip().lower() != "arial":
                    off_fonts.add(name)
            hx = _run_hex(run)
            if hx:
                hex_sightings.setdefault(hx, idx)
            size = run.font.size
            if size is not None and size.pt > MAX_RUN_PT:
                oversize.append((idx, int(round(size.pt)), run.text.strip()[:40]))

    result["fonts_found"] = "; ".join(sorted(fonts))
    result["fonts_offending"] = "; ".join(sorted(off_fonts))
    if off_fonts:
        hard("H-FORMAT-FONT",
             "non-Arial text runs: %s (SKILL §13.1 — Arial, every run, every slide)"
             % "; ".join(sorted(off_fonts)))

    for hx, why in sorted(FORBIDDEN_HEX.items()):
        if hx in hex_sightings:
            hard("H-FORMAT-HEX", "text run coloured #%s — %s" % (hx, why),
                 hex_sightings[hx])

    drift = sorted(h for h in hex_sightings
                   if h not in PALETTE and h not in FORBIDDEN_HEX)
    result["colors_offending"] = "; ".join(drift)
    if drift:
        adv("A-PALETTE-DRIFT",
            "%d text hex(es) outside the §13.2 palette: %s"
            % (len(drift), "; ".join(drift[:12]) + (" ..." if len(drift) > 12 else "")))

    if oversize:
        biggest = sorted(oversize, key=lambda x: -x[1])[:4]
        hard("H-FORMAT-SCALE",
             "%d text run(s) above %dpt, the day-divider headline ceiling: %s "
             "(SKILL §13.3, §13.4)"
             % (len(oversize), MAX_RUN_PT,
                "; ".join("s%d %dpt %r" % (s, p, t) for s, p, t in biggest)))

    # -- tier labels and tier colours ---------------------------------------
    tier_slides = 0
    diagnostic_slides = 0
    markers_seen = set()

    for idx, slide in enumerate(slides, start=1):
        text = texts[idx - 1]
        if is_excluded(text):
            continue                                # matches the grader exactly
        slide_is_diag = False

        if all(lbl in text for lbl in TIER_LABELS):
            slide_is_diag = True
            tier_slides += 1
            for sh in slide.shapes:
                if not sh.has_text_frame:
                    continue
                for p in sh.text_frame.paragraphs:
                    line = "".join(r.text for r in p.runs).strip()
                    for lbl in TIER_LABELS:
                        if line == lbl or line.startswith(lbl):
                            hexes = [h for h in (_run_hex(r) for r in p.runs
                                                 if r.text.strip()) if h]
                            if not hexes:
                                adv("A-TIER-COLOR",
                                    "tier '%s' has no explicit colour (theme/auto) — "
                                    "expected #%s; verify it is not rendering black"
                                    % (lbl, TIER_HEX[lbl]), idx)
                            elif not (set(hexes) & TIER_HEX_OK[lbl]):
                                if "000000" in hexes:
                                    hard("H-TIER-BLACK",
                                         "tier '%s' is BLACK (#000000); must be #%s"
                                         % (lbl, TIER_HEX[lbl]), idx)
                                else:
                                    adv("A-TIER-COLOR",
                                        "tier '%s' colour %s is not one of %s"
                                        % (lbl, hexes, sorted(TIER_HEX_OK[lbl])), idx)

        for mname, rx in MARKERS.items():
            if rx.search(text):
                markers_seen.add(mname)
                slide_is_diag = True
        if slide_is_diag:
            diagnostic_slides += 1

    result["tier_slides"] = tier_slides
    result["diagnostic_slides"] = diagnostic_slides
    result["markers_present"] = "; ".join(m for m in MARKER_ORDER if m in markers_seen)

    for lbl in TIER_LABELS:
        if lbl not in all_text:
            hard("H-TIER-VERBATIM",
                 "tier label '%s' never appears verbatim anywhere in the deck "
                 "(a typo or wrong casing breaks grader classification)" % lbl)

    # -- markers present and unbroken (the ellipsis fault, SKILL §3c) -------
    for idx, text in enumerate(texts, start=1):
        if is_excluded(text):
            continue
        fixed = _normalize_typography(text)
        if fixed == text:
            continue
        for mname, rx in MARKERS.items():
            if rx.search(fixed) and not rx.search(text):
                # find the line that would have matched, and quote it verbatim
                bad = ""
                for raw, cooked in zip(text.split("\n"), fixed.split("\n")):
                    if rx.search(cooked) and not rx.search(raw):
                        bad = raw.strip()[:90]
                        break
                hard("H-MARKER-BROKEN",
                     "marker '%s' broken by typographic punctuation — it matches only "
                     "after normalising. The grader still classifies this slide and the "
                     "linter does not, so the two counts drift by one (SKILL §3c). "
                     "Slide text: %r" % (mname, bad), idx)

    missing_markers = [m for m in MARKER_ORDER if m not in markers_seen]
    if missing_markers:
        adv("A-MARKER-ABSENT",
            "diagnostic markers not found in deck: %s "
            "(acceptable if this cycle does not use them; flagged for eyeball)"
            % missing_markers)

    # -- required slides ----------------------------------------------------
    def slides_of(kind):
        return [i for i, t in enumerate(types, start=1) if t == kind]

    prep = slides_of("Teacher Prep")
    banks = slides_of("Concept Bank")
    links = slides_of("Links slide")
    notes = slides_of("Teacher note")
    dividers = [i for i, t in enumerate(texts, start=1) if DAY_DIVIDER_RX.search(t)]
    day3 = None
    for i in dividers:
        m = DAY_DIVIDER_RX.search(texts[i - 1])
        if m and m.group(1) == "3":
            day3 = i
            break

    result["teacher_prep"] = len(prep)
    result["concept_bank"] = len(banks)
    result["links_slide"] = len(links)
    result["teacher_note"] = len(notes)
    result["day_dividers"] = len(dividers)

    if not prep:
        hard("H-TEACHER-PREP",
             "no Teacher Prep slide — slide 2 must be 'Additional Resources — "
             "Teacher Prep' with BUY AHEAD / PREPARE / PUT OUT FOR STUDENTS "
             "(SKILL §2.0)")

    if not banks:
        hard("H-CONCEPT-BANK",
             "no Concept Bank slide — every deck carries exactly one, headed "
             "'Define these in your own words' (SKILL §3b)")
    else:
        if len(banks) > 2:
            hard("H-CONCEPT-BANK",
                 "%d Concept Bank slides; the ceiling is two (SKILL §3b)" % len(banks))
        if day3 is None:
            hard("H-CONCEPT-BANK",
                 "Concept Bank present but there is no Day 3 divider to position it "
                 "against (SKILL §2.3, §2.5 — never delete a day divider)")
        elif max(banks) > day3:
            hard("H-CONCEPT-BANK",
                 "Concept Bank on slide %d sits BELOW the Day 3 divider (slide %d). "
                 "An unfilled cell is a question the student still owes, and no "
                 "question may follow the divider (SKILL §2.3)"
                 % (max(banks), day3), max(banks))

    if not links:
        hard("H-LINKS-SLIDE",
             "no links slide — 'Activity and resource links' must list every "
             "hyperlink in the deck, in deck order (SKILL §8)")
    else:
        link_urls = set()
        for i in links:
            link_urls |= _slide_urls(slides[i - 1])
        result["links_slide_urls"] = len(link_urls)
        if not link_urls:
            hard("H-LINKS-SLIDE",
                 "links slide carries no URL — every named document must have a link "
                 "that works (SKILL §8)", links[0])
        malformed = sorted(u for u in link_urls
                           if not re.match(r"^https?://[^/\s]+\.[^/\s]", u))
        if malformed:
            hard("H-LINKS-SLIDE",
                 "malformed URL(s) on the links slide: %s" % "; ".join(malformed[:5]),
                 links[0])
        # The Image credits slide carries attribution links (Wikimedia, the
        # Twemoji licence) and the Teacher Prep slide carries product links.
        # §8 governs the resource links a student follows, so neither belongs
        # on the links slide and neither counts as unlisted.
        deck_urls = set()
        for i, slide in enumerate(slides, start=1):
            if i in links or types[i - 1] in ("Image credits", "Teacher Prep",
                                              "Teacher navigation", "Teacher reference"):
                continue
            deck_urls |= _slide_urls(slide)
        unlisted = sorted(u for u in deck_urls if u not in link_urls)
        result["links_unlisted"] = len(unlisted)
        if unlisted:
            adv("A-LINKS-UNLISTED",
                "%d hyperlink(s) used in the deck but not listed on the links slide: %s"
                % (len(unlisted), "; ".join(unlisted[:4])
                   + (" ..." if len(unlisted) > 4 else "")))

    # -- teacher note slide (new 2026-08-29) --------------------------------
    note_text = "\n".join(texts[i - 1] for i in notes).lower()
    if not notes:
        hard("H-TEACHER-NOTE",
             "no teacher note slide. One slide, not speaker notes, deletable in one "
             "action, declaring the seven items. A teacher adapting the deck in good "
             "faith will substitute a better example and destroy the contrast "
             "(spec: 'Teacher note slide')")
        missing_items = [name for name, _ in TEACHER_NOTE_ITEMS]
    else:
        missing_items = [name for name, test in TEACHER_NOTE_ITEMS
                         if not test(note_text)]

    if missing_items:
        adv("A-TN-ITEMS",
            "teacher note does not declare %d of the seven required items: %s"
            % (len(missing_items), "; ".join(missing_items)))

    # -- visibility-ladder slide (new 2026-08-29) ---------------------------
    # ADVISORY, never hard. The slide was invented on 2026-08-29 and no shipped
    # deck in the arc carries one yet, so a hard rule here would fail every deck
    # in the arc on its first run — and a linter that fails every deck gets
    # ignored inside a week. That is the whole two-tier rationale.
    ladders = slides_of("Visibility ladder")
    result["visibility_ladder"] = len(ladders)
    if not ladders:
        adv("A-VIS-LADDER",
            "no visibility-ladder slide. The teacher note declares this cycle's "
            "rung; nothing in the deck says what the rungs are, so the "
            "declaration reads as jargon to a teacher who has not met the "
            "ladder. One teacher-facing slide, marked 'do not project', "
            "immediately after the note (scripts/deck_apply_changes.py inserts "
            "it)")
    else:
        if len(ladders) > 1:
            adv("A-VIS-LADDER",
                "%d visibility-ladder slides; the deck needs one"
                % len(ladders))
        if notes and min(ladders) < min(notes):
            adv("A-VIS-LADDER",
                "the visibility ladder (slide %d) sits ABOVE the teacher note "
                "(slide %d). The note declares the rung and this slide explains "
                "it — read the other way round it is a glossary in front of a "
                "term nobody has met yet"
                % (min(ladders), min(notes)), min(ladders))
        for i in ladders:
            if "do not project" not in texts[i - 1].lower():
                adv("A-VIS-LADDER",
                    "the visibility ladder (slide %d) has lost its 'do not "
                    "project' marker, so it is counted as a diagnostic slide "
                    "and a teacher may project it (SKILL §7.1)" % i, i)

    # -- critical aspects, and the move-1 difference test -------------------
    aspects = []            # ordered, unique
    aspect_slides = {}
    for i, text in enumerate(texts, start=1):
        m = CA_LABEL.search(text)
        if not m:
            continue
        name = re.split(r"\s{2,}|\s+/\s+", m.group(1).strip())[0].strip()
        if not name:
            continue
        if name not in aspect_slides:
            aspects.append(name)
            aspect_slides[name] = []
        aspect_slides[name].append(i)

    result["critical_aspects"] = " | ".join(aspects)
    result["n_aspects"] = len(aspects)

    move1 = {}              # aspect -> (slide number, question text)
    for name in aspects:
        for i in aspect_slides[name]:
            if types[i - 1] == "Content/Explanation" and "?" in texts[i - 1]:
                move1[name] = (i, texts[i - 1])
                break

    verdicts = []
    for name in aspects:
        if name not in move1:
            adv("A-MOVE1-NO-DIFF",
                "critical aspect %r has no Critical Aspect question slide of its own — "
                "the aspect block opens on another slide type, so move 1 cannot carry "
                "a difference" % name)
            verdicts.append("%s: NO MOVE-1 SLIDE" % name)
            continue
        idx, raw = move1[name]
        body = "\n".join(
            ln for ln in raw.split("\n")
            if not ln.lower().startswith("critical aspect:")
            and "think first" not in ln.lower()
            and not ln.strip().startswith("[["))
        devices = difference_devices(body)
        if devices:
            verdicts.append("S%d: %s" % (idx, ", ".join(devices)))
        else:
            verdicts.append("S%d: NO DIFFERENCE" % idx)
            question = " ".join(body.split())[:200]
            adv("A-MOVE1-NO-DIFF",
                "the Critical Aspect question for %r names the aspect and waits for the "
                "Contrast Set to supply the difference. It carries no comparison, no "
                "choice between named alternatives and no stated change condition, so "
                "it is structurally 'how are polar bears adapted?' and will produce "
                "silence. Text: %r" % (name, question), idx)
    result["move1_verdict"] = " ;; ".join(verdicts)
    result["move1_pass"] = bool(aspects) and all("NO " not in v for v in verdicts)

    if not aspects:
        adv("A-MOVE1-NO-DIFF",
            "no 'Critical aspect:' group label anywhere in the deck — there is no "
            "move 1 to test (SKILL §7)")

    # -- coordination structures -------------------------------------------
    declared = []
    for label, pats in sorted(DECLARED_STRUCTURE.items()):
        hit = None
        for i, text in enumerate(texts, start=1):
            for line in text.split("\n"):
                if not _any_match(pats, line):
                    continue
                # A teacher note that says the compensatory pair was LEFT OUT
                # has not declared one.
                if _any_match([r"left out", r"not used", r"omitted", r"\bno\b",
                               r"\bnone\b", r"\bwithout\b", r"carries neither"], line):
                    continue
                hit = i
                break
            if hit:
                break
        if hit:
            declared.append("%s (slide %d)" % (label, hit))
    sf_slides = slides_of("Stock-and-flow model")
    if sf_slides and not any(d.startswith("stock-and-flow") for d in declared):
        declared.append("stock-and-flow model (slide %d)" % sf_slides[0])

    candidates = []
    for i, text in enumerate(texts, start=1):
        if types[i - 1] in ("Pattern Break", "Contrast Set", "What if"):
            if _any_match(CANDIDATE_CONFLICT, text):
                candidates.append("possible conflict case (slide %d)" % i)
            if _any_match(CANDIDATE_COMPENSATORY, text):
                candidates.append("possible compensatory pair (slide %d)" % i)

    result["coordination"] = "; ".join(declared)
    result["coordination_candidates"] = "; ".join(candidates)

    if not declared:
        msg = ("no declared coordination structure on any critical aspect. "
               "Stock-and-flow is a representation and can be ADDED to any cycle whose "
               "content accumulates; a compensatory pair and a conflict case have to be "
               "FOUND in the content. An absence is not automatically a hole.")
        if candidates:
            msg += " Leads for the retrofit: %s" % "; ".join(candidates[:4])
        adv("A-NO-FUSION", msg)

    # -- critical aspect declared in the teacher note -----------------------
    if aspects:
        if not notes:
            adv("A-CA-NOT-DECLARED",
                "critical aspect(s) %s are not declared in a teacher note, because "
                "there is no teacher note slide" % ", ".join(repr(a) for a in aspects))
        else:
            undeclared = [a for a in aspects if a.lower() not in note_text]
            if undeclared:
                adv("A-CA-NOT-DECLARED",
                    "critical aspect(s) not named in the teacher note: %s"
                    % ", ".join(repr(a) for a in undeclared))

    # -- conditional slide types present without a stated reason ------------
    present_conditional = []
    for kind in CONDITIONAL_TYPES:
        got = slides_of(kind)
        if got:
            present_conditional.append("%s (slide%s %s)"
                                       % (kind, "" if len(got) == 1 else "s",
                                          ", ".join(str(g) for g in got)))
    result["conditional_types"] = "; ".join(present_conditional)
    if present_conditional:
        has_reason = bool(notes) and ("slide type" in note_text or "slide-type" in note_text)
        if not has_reason:
            adv("A-COND-NO-REASON",
                "conditional slide types present with no stated reason: %s. The "
                "slide-type map in the teacher note has to say which conditional types "
                "appear and which were left out and why"
                % "; ".join(present_conditional))

    # -- core question dropped without a stated reason ----------------------
    dropped_core = [k for k in CORE_TYPES if not slides_of(k)]
    if not move1:                       # move 1 is marker-found, not type-found
        dropped_core.insert(0, "Critical Aspect question")
    result["core_dropped"] = "; ".join(dropped_core)
    if dropped_core:
        declared = bool(notes) and (
            "drop" in note_text or "not used" in note_text or "left out" in note_text
            or "omitted" in note_text or "does not carry" in note_text)
        if not declared:
            adv("A-CORE-DROPPED",
                "core question(s) absent with no stated reason: %s. Dropping one is a "
                "design decision when the content does not support it — but the "
                "slide-type map in the teacher note has to say so, or the drop cannot "
                "be told apart from an oversight"
                % ", ".join(dropped_core))

    # -- visibility rung declared -------------------------------------------
    if not _any_match(VISIBILITY_PATTERNS, low_all):
        adv("A-VISIBILITY",
            "no visibility rung declared. Declaring it — e.g. 'written and private, "
            "building toward unattributed read-aloud' — converts a quiet room from an "
            "absence an observer reads as failure into a design decision an observer "
            "can argue with")
        result["visibility_declared"] = False
    else:
        result["visibility_declared"] = True

    # -- Concept Bank asks the student to RELATE two terms ------------------
    if banks:
        bank_text = " ".join(texts[i - 1] for i in banks).lower()
        relates = _any_match([
            r"how (?:any )?two[^.?]{0,60}(?:connect|relate|fit together|go together)",
            r"connect(?:s|ion)? between two",
            r"pick two[^.?]{0,60}(?:connect|relate|explain)",
            r"choose two[^.?]{0,60}(?:connect|relate|explain)",
            r"relate(?:s)? (?:any )?two",
            r"two of (?:the|these) terms",
        ], bank_text)
        result["bank_relates"] = relates
        # A-BANK-NO-RELATE retired 2026-08-30 by Katherine. The Concept Bank is a
        # fill-in-the-blank at the point in the cycle where it sits: "We're not
        # asking students to put those two ideas together for any of that
        # vocabulary. They're not ready for that at that step. They need to just
        # fill in the blanks." The relating is the conflict case, which follows
        # the bank as its own slide, and is a different ask of a different kind.
        #
        # [This contradicts coordination_judgment_and_the_package section 7, which
        # says a page that only lists "has achieved co-presence and left the
        # relating to chance" and that a page asking how two terms connect "does
        # the work". Both cannot stand as written. The reading that reconciles
        # them: the document is arguing about where diachronic simultaneity gets
        # engineered, and Katherine's answer is that it is engineered by the slide
        # after the bank rather than inside it - the bank supplies co-presence,
        # the conflict case does the relating. If that is right, the document
        # needs a sentence saying so, or this check comes back the next time
        # someone reads section 7 and treats it as a spec.]
    else:
        result["bank_relates"] = False

    # -- what-if, where the cycle carries a coordination structure ----------
    whatif = slides_of("What if")
    result["what_if"] = len(whatif)
    if declared:
        if not whatif:
            adv("A-WHATIF",
                "the cycle carries a coordination structure (%s) but no what-if slide. "
                "The what-if does not create coordination, it tests it — treat it as "
                "the required companion to whichever structure the cycle carries"
                % "; ".join(declared))
        else:
            wtext = " ".join(texts[i - 1] for i in whatif).lower()
            individual = _any_match([
                r"on your own", r"\bindividually\b", r"by yourself",
                r"do not discuss", r"don't discuss", r"without (?:talking|discussing)",
                r"answer this alone", r"your own answer, alone",
            ], wtext)
            written = _any_match([
                r"\bwrite\b", r"\bwritten\b", r"your first answer",
                r"write your", r"in the box",
            ], wtext)
            group_invite = _any_match([
                r"discuss with a partner", r"with your (?:partner|group|table)",
                r"talk (?:it )?(?:through )?with", r"in your group",
            ], wtext)
            result["whatif_individual"] = individual and not group_invite
            result["whatif_written"] = written
            if not (individual and written) or group_invite:
                bits = []
                if not individual:
                    bits.append("it is not marked individual")
                if group_invite:
                    bits.append("it invites partner or group talk")
                if not written:
                    bits.append("it is not marked written")
                adv("A-WHATIF",
                    "the what-if (slide%s %s) must be individual and written: %s. Group "
                    "work is a visibility reducer, not a rung — a group product says "
                    "nothing about who coordinated, so a shared what-if is empty as "
                    "evidence"
                    % ("" if len(whatif) == 1 else "s",
                       ", ".join(str(x) for x in whatif), "; ".join(bits)),
                    whatif[0])

    # -- forbidden vocabulary ------------------------------------------------
    for idx, text in enumerate(texts, start=1):
        m = FORBIDDEN_VOCAB.search(text)
        if not m:
            continue
        line = ""
        for ln in text.split("\n"):
            if FORBIDDEN_VOCAB.search(ln):
                line = ln.strip()[:100]
                break
        hard("H-VOCAB",
             "forbidden vocabulary, structural-noun sense: the deck says %r. The words "
             "are '5 core questions' and 'slide types' (SKILL §1)" % (line,), idx)

    result["slide_types_ordered"] = " > ".join(types)
    result["hard"] = sum(1 for f in findings if f.tier == HARD)
    result["advisory"] = sum(1 for f in findings if f.tier == ADVISORY)
    result["ships"] = result["hard"] == 0
    return result


# ---------------------------------------------------------------------------
# Input expansion
# ---------------------------------------------------------------------------

def expand_paths(args):
    """Accept files, directories and globs. Return sorted unique .pptx paths."""
    out = []
    for a in args:
        a = os.path.expanduser(a)
        matches = globmod.glob(a)
        if not matches and os.path.exists(a):
            matches = [a]
        if not matches:
            print("deck_lint: no such path: %s" % a, file=sys.stderr)
            continue
        for m in matches:
            if os.path.isdir(m):
                for root, _dirs, files in os.walk(m):
                    for f in files:
                        if f.lower().endswith(".pptx") and not f.startswith("~$"):
                            out.append(os.path.join(root, f))
            elif m.lower().endswith(".pptx") and not os.path.basename(m).startswith("~$"):
                out.append(m)
    return sorted(set(out))


# ---------------------------------------------------------------------------
# Output
# ---------------------------------------------------------------------------

CSV_COLUMNS = [
    "file", "folder", "slides", "ships", "hard", "advisory",
    "hard_codes", "advisory_codes",
    "dims", "fonts_offending", "colors_offending",
    "tier_slides", "diagnostic_slides", "markers_present",
    "teacher_prep", "concept_bank", "links_slide", "teacher_note",
    "visibility_ladder",
    "day_dividers", "n_aspects", "critical_aspects",
    "move1_pass", "move1_verdict",
    "coordination", "coordination_candidates", "conditional_types",
    "visibility_declared", "bank_relates", "what_if",
    "slide_types_ordered",
]


def _codes(result, tier):
    seen = []
    for f in result["findings"]:
        if f.tier == tier and f.code not in seen:
            seen.append(f.code)
    return "; ".join(seen)


def print_human(result, tier_filter):
    print("=" * 74)
    print("deck_lint: %s" % result["file"])
    print("  folder:              %s" % result["folder"])
    print("  total slides:        %s" % result["slides"])
    print("  3-tier CQ slides:    %s" % result.get("tier_slides", 0))
    print("  diagnostic slides:   %s" % result.get("diagnostic_slides", 0))
    print("  markers present:     %s" % (result.get("markers_present") or "(none)"))
    print("  critical aspects:    %s" % (result.get("critical_aspects") or "(none)"))
    print("  move 1 difference:   %s" % (result.get("move1_verdict") or "(no move 1)"))
    print("  coordination:        %s" % (result.get("coordination") or "none declared"))
    print("-" * 74)
    shown = 0
    for want, label in ((HARD, "HARD    "), (ADVISORY, "ADVISORY")):
        if tier_filter not in ("all", want.lower()):
            continue
        for f in result["findings"]:
            if f.tier == want:
                print("  %s %s" % (label, f.line()))
                shown += 1
    if shown == 0:
        print("  nothing to report at tier '%s'" % tier_filter)
    verdict = ("SHIPS — no hard failures" if result["hard"] == 0
               else "DOES NOT SHIP — %d hard failure(s)" % result["hard"])
    print("-" * 74)
    print("  %s;  %d advisory" % (verdict, result["advisory"]))
    print("=" * 74)


def print_summary(results):
    print()
    print("SUMMARY  (%d deck%s)" % (len(results), "" if len(results) == 1 else "s"))
    print("%-62s %5s %5s %8s" % ("deck", "hard", "adv", "ships"))
    print("-" * 84)
    for r in results:
        print("%-62s %5d %5d %8s"
              % (r["file"][:62], r["hard"], r["advisory"],
                 "yes" if r["ships"] else "NO"))
    ships = sum(1 for r in results if r["ships"])
    print("-" * 84)
    print("%d of %d ship;  %d hard finding(s), %d advisory finding(s) in total"
          % (ships, len(results),
             sum(r["hard"] for r in results),
             sum(r["advisory"] for r in results)))


def write_csv(results, path):
    path = os.path.expanduser(path)
    d = os.path.dirname(os.path.abspath(path))
    if d and not os.path.isdir(d):
        os.makedirs(d)
    with open(path, "w") as fh:
        w = csvmod.writer(fh)
        w.writerow(CSV_COLUMNS)
        for r in results:
            row = []
            for col in CSV_COLUMNS:
                if col == "hard_codes":
                    row.append(_codes(r, HARD))
                elif col == "advisory_codes":
                    row.append(_codes(r, ADVISORY))
                else:
                    row.append(r.get(col, ""))
            w.writerow(row)


def main(argv=None):
    ap = argparse.ArgumentParser(
        prog="deck_lint.py",
        description="VT deck validator, two tiers. HARD fails the build; "
                    "ADVISORY is reported and does not.")
    ap.add_argument("paths", nargs="+",
                    help="deck .pptx files, directories, or globs")
    ap.add_argument("--tier", choices=("hard", "advisory", "all"), default="all",
                    help="which tier to print (default all). Never changes the exit code.")
    ap.add_argument("--json", action="store_true",
                    help="emit machine-readable JSON on stdout instead of the report")
    ap.add_argument("--csv", metavar="PATH", default=None,
                    help="write a one-row-per-deck summary table to PATH")
    ap.add_argument("--no-summary", action="store_true",
                    help="suppress the summary table in human output")
    args = ap.parse_args(argv)

    decks = expand_paths(args.paths)
    if not decks:
        print("deck_lint: nothing to lint", file=sys.stderr)
        return 2

    results = [lint_deck(p) for p in decks]

    if args.csv:
        write_csv(results, args.csv)

    if args.json:
        payload = []
        for r in results:
            item = dict((k, v) for k, v in r.items() if k != "findings")
            item["findings"] = [f.as_dict() for f in r["findings"]]
            payload.append(item)
        json.dump({"decks": payload,
                   "hard_total": sum(r["hard"] for r in results),
                   "advisory_total": sum(r["advisory"] for r in results),
                   "ships": sum(1 for r in results if r["ships"]),
                   "count": len(results)},
                  sys.stdout, indent=2)
        print()
    else:
        for r in results:
            print_human(r, args.tier)
        if not args.no_summary and len(results) > 1:
            print_summary(results)
        if args.csv:
            print("\ncsv written: %s" % os.path.abspath(os.path.expanduser(args.csv)))

    return 1 if any(r["hard"] for r in results) else 0


if __name__ == "__main__":
    sys.exit(main())

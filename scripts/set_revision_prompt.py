#!/usr/bin/env python3
"""Put Katherine's 2026-08-30 revision prompt into the speaker notes of every
slide that asks a student to revise an answer.

The prompt is her wording. Two mechanical corrections only: her draft numbered
the items 1, 2, 2, 3, which is renumbered 1-4; and the `-----` rule she asked for
above `[PASTE IN YOUR FIRST ANSWER]:` is added.

QUESTION: carries the question that slide actually asked, derived rather than
left as a placeholder:

  What if slide      the question is on the slide - the long text run that is
                     not the title, the watch-first line, or a box label.
  Response slide     its marker ends -MASTERY, so the question is the Mastery
                     tier question on the nearest 3-Tier slide before it.

A slide whose question cannot be derived is reported and left alone. Guessing
the question would be worse than a gap a person can see.

Usage:  python3 set_revision_prompt.py <dir> [--apply]
"""
import argparse
import glob
import os
import re

from pptx import Presentation

PROMPT = """REVISION PROMPT — Open an AI and copy and paste the following:
QUESTION: {question}
-----
[PASTE IN YOUR FIRST ANSWER]:
-----
I am a high school student. Below is a science question and my first answer. Do NOT rewrite my answer and do NOT give me the answer. Instead:
1. Tell me one thing I got right, or was interesting.
2. Quote my own words back to me as You said, “”.
3. Ask me two questions that make me look again at one idea I might have wrong or left out or could be made more complete.
4. Illustrate it with one distinction or example worth thinking about.
—-
Keep it short and in plain language, then stop so I can write my own revised answer.
-----
Once you get your feedback: write your revised answer in “Your revised answer” on the slide — in your own words."""

SKIP_PREFIX = ("Watch first", "Think first", "Your first answer",
               "Your revised answer", "Your answer —", "Your answer -",
               "Open:", "Critical aspect:", "What if?", "Day ")


def texts(slide):
    return [sh.text_frame.text.strip() for sh in slide.shapes
            if sh.has_text_frame and sh.text_frame.text.strip()]


def question_on_slide(slide):
    """The longest run that reads as the question rather than furniture.

    A What if question can run to several sentences and need not end on the
    first line, so the whole run is tested for a question mark rather than
    only its opening line - that is what sent Cycle 16d's What if slide to the
    Mastery question instead of its own.
    """
    best = ""
    for t in texts(slide):
        line = t.split("\n")[0].strip()
        if line.startswith(SKIP_PREFIX) or "[[" in line:
            continue
        if "?" not in t:
            continue
        if len(line) > len(best):
            best = line
    return best or None


def mastery_question(slides, upto):
    """The Mastery tier question on the nearest 3-Tier slide before `upto`."""
    for i in range(upto - 1, -1, -1):
        blob = texts(slides[i])
        joined = " ".join(blob)
        if not all(k in joined for k in ("Getting Started", "Working On It", "Mastery")):
            continue
        # the run after the one that is exactly "Mastery"
        for j, t in enumerate(blob):
            if t == "Mastery" and j + 1 < len(blob):
                return blob[j + 1].split("\n")[0].strip()
    return None


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("path")
    ap.add_argument("--apply", action="store_true")
    args = ap.parse_args()

    files = sorted(glob.glob(os.path.join(args.path, "*.pptx"))) \
        if os.path.isdir(args.path) else sorted(glob.glob(args.path))

    done = missing = 0
    for f in files:
        prs = Presentation(f)
        slides = list(prs.slides)
        hits = []
        for i, s in enumerate(slides):
            if not s.has_notes_slide:
                continue
            note = s.notes_slide.notes_text_frame.text
            if "REVISION PROMPT" not in note:
                continue
            first = (texts(s) or [""])[0].split("\n")[0]
            q = question_on_slide(s)
            if q is None:
                q = mastery_question(slides, i)
            hits.append((i + 1, first[:46], q))

        if not hits:
            continue
        print("\n=== %s" % os.path.basename(f))
        for n, first, q in hits:
            if q:
                print("   s%-3d %-46s Q: %s" % (n, first, q[:74]))
                done += 1
            else:
                print("   s%-3d %-46s ** QUESTION NOT DERIVED — left alone **" % (n, first))
                missing += 1

        if args.apply:
            for (n, _first, q) in hits:
                if not q:
                    continue
                tf = slides[n - 1].notes_slide.notes_text_frame
                tf.text = PROMPT.format(question=q)
            prs.save(f)

    print("\n%s: %d slide(s) set, %d left alone for want of a derivable question"
          % ("APPLIED" if args.apply else "WOULD SET", done, missing))
    if not args.apply:
        print("Report only. Re-run with --apply to write.")


if __name__ == "__main__":
    main()

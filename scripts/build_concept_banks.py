import os,re,sys,glob,copy,zipfile
import xml.etree.ElementTree as ET
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.shapes.autoshape import Shape

A='{http://schemas.openxmlformats.org/drawingml/2006/main}'
R='{http://schemas.openxmlformats.org/officeDocument/2006/relationships}'
AHYP='http://schemas.microsoft.com/office/drawing/2018/hyperlinkcolor'
ROWS=[1554480,2148840,2743200,3337560,3931920,4526280,5120640]
COL={1:{'t':548640,'c':2194560},2:{'t':4663440,'c':6309360}}

DROP={
 '02':{'both gain'},
 '03':{'energy'},
 '04':{'air','mass','used again'},
 '05':{'air','soil'},
 '06':set(),
 '07a':{'resources','run out','level off','crash','resource'},
 '07b':{'stable','reset','soil','acid'},
 '08':{'rough','smooth','folds','surface','sunlight','energy'},
 '09':{'let in','block','needed things enter','waste leaves','germs stay out','heads','tails',
       'oily middle','charged','oily things pass'},
 '10':{'lowering','same products','one','key'},
 '11':{'repair','self-destruct'},
 '13':{'masks','one allele','chance','random'},
 '14':{'copy'},
 '15a':{'coiled'},
 '15b':{'order'},
 '16a':{'job','adult','embryonic','specialized'},
 '16b':set(),
 '16c':{'inherited','body cell'},
 '16d':{'splice'},
 '17':{'pattern','independent evidence'},
 '18':{'survive'},
 '19':{'generations','classify','predict','related'},
 '20':set(),
}
SKIP=re.compile(r'^(NOTES|DRAFT|OPTIONAL|MARKER)[:\-]',re.I)
def is_term(t):
    t=t.strip(" .·")
    if not t or len(t)>28 or SKIP.search(t): return False
    if re.fullmatch(r'[\d%\W]+',t): return False
    if len(t)<=2 and t.isalpha(): return False
    if re.fullmatch(r'[A-Za-z]{1,3}[\-–][A-Za-z]{1,3}',t): return False
    w=t.split()
    if len(w)>3: return False
    if len(w)>1 and re.search(r'\b(is|are|no|not|the|a|an|and|or|to|of|its|it)\b',t.lower()): return False
    return True

def raw_slides(p):
    z=zipfile.ZipFile(p)
    n=sorted([x for x in z.namelist() if re.match(r'ppt/slides/slide\d+\.xml$',x)],
             key=lambda x:int(re.search(r'(\d+)',x.split('/')[-1]).group(1)))
    return [[t.text for t in ET.fromstring(z.read(s)).iter(A+'t') if t.text] for s in n]

def terms_for(p,key):
    sl=raw_slides(p); cand=[]
    for runs in sl:
        if any('KEY TERMS' in r for r in runs):
            grab=False
            for r in runs:
                if 'KEY TERMS' in r: grab=True; continue
                if grab:
                    if re.match(r'^(FOR READERS|FOR TESTERS|LAB|TOPIC|SHORT)',r): break
                    head=re.split(r'\s+[—–-]\s+',r)[0].strip()
                    for part in re.split(r'\s*[·/]\s*',head):
                        part=part.strip()
                        if is_term(part): cand.append(part)
            break
    for runs in sl:
        j="\n".join(runs)
        if 'Getting Started' in j and 'Mastery' in j:
            for m in re.findall(r'\[([^\]\[]{1,60})\]',j):
                m=m.strip()
                if is_term(m): cand.append(m)
    drop=DROP.get(key,set()); out=[]; seen=set()
    def near(k):
        for o in out:
            ol=o.lower()
            if k==ol: return True
            if ' ' not in k and re.search(r'\b'+re.escape(k)+r'\b',ol): return True
        return False
    for t in cand:
        k=t.lower()
        if k in seen or k in drop or near(k): continue
        seen.add(k); out.append(t)
    return out[:14]

def build(path,key,outpath):
    prs=Presentation(path); S=prs.slides
    resp=None
    for s in S:
        has_lbl=any(sh.has_text_frame and sh.text_frame.text.strip().startswith('Your first answer') for sh in s.shapes)
        if has_lbl:
            box=[sh for sh in s.shapes if not sh.has_text_frame or sh.text_frame.text.strip()=='']
            if box: resp=s; break
    if resp is None: return None,'no response slide'
    lbl_src=next(sh._element for sh in resp.shapes if sh.has_text_frame and sh.text_frame.text.strip().startswith('Your first answer'))
    box_src=None
    for sh in resp.shapes:
        try:
            if sh.fill.type==1 and str(sh.fill.fore_color.rgb)=='F2F6F9': box_src=sh._element; break
        except Exception: pass
    if box_src is None: return None,'no writing box'
    links=[]
    for s in S:
        for sh in s.shapes:
            if not sh.has_text_frame: continue
            if 'FOR READERS' in sh.text_frame.text:
                for para in sh.text_frame.paragraphs:
                    lab=None;txt=None;url=None
                    for r in para.runs:
                        if r.text.strip().startswith(('FOR READERS','FOR TESTERS')): lab=r.text
                        elif r.hyperlink and r.hyperlink.address: txt=r.text; url=r.hyperlink.address
                    if lab and url: links.append((lab,txt,url))
                break
        if links: break
    terms=terms_for(path,key)
    if not terms: return None,'no terms'
    new=S.add_slide(resp.slide_layout)
    for sh in list(new.shapes): sh._element.getparent().remove(sh._element)
    tree=new.shapes._spTree
    def tb(x,y,w,h,text,size,bold,color,spc=None):
        b=new.shapes.add_textbox(Emu(x),Emu(y),Emu(w),Emu(h))
        f=b.text_frame; f.word_wrap=True
        for m in ('lIns','rIns','tIns','bIns'): f._bodyPr.set(m,'0')
        p=f.paragraphs[0]; r=p.add_run(); r.text=text
        fo=r.font; fo.name='Arial'; fo.size=Pt(size); fo.bold=bold; fo.color.rgb=RGBColor.from_string(color)
        if spc:
            pPr=p._p.find(A+'pPr')
            if pPr is None: pPr=p._p.makeelement(A+'pPr',{}); p._p.insert(0,pPr)
            pPr.set('spc',str(spc))
        return b
    tb(502920,109728,8321040,274320,'CONCEPT BANK',11,True,'028090',130)
    tb(548640,640080,8046720,457200,'Define these in your own words',18,True,'111111')
    tb(548640,1097280,8046720,457200,
       'Fill each box as you meet the term on Day 1 and Day 2. Use your own words, not a copied definition. '
       'On revision day, read this slide before you rewrite anything.',12,False,'111111')
    for i,term in enumerate(terms):
        col=1 if i<7 else 2; ry=ROWS[i%7]
        el=copy.deepcopy(lbl_src); tree.append(el); sh=Shape(el,None)
        sh.left,sh.top,sh.width,sh.height=Emu(COL[col]['t']),Emu(ry+128070),Emu(1554480),Emu(292500)
        bp=sh.text_frame._bodyPr
        for a in bp.findall(A+'spAutoFit'): bp.remove(a)
        bp.set('wrap','square'); bp.set('anchor','ctr')
        for m in ('lIns','rIns','tIns','bIns'): bp.set(m,'0')
        p=sh.text_frame.paragraphs[0]
        pPr=p._p.find(A+'pPr')
        if pPr is None: pPr=p._p.makeelement(A+'pPr',{}); p._p.insert(0,pPr)
        pPr.set('algn','l')
        for r in list(p.runs)[1:]: r._r.getparent().remove(r._r)
        r=p.runs[0]; r.text=term
        fo=r.font; fo.name='Arial'; fo.size=Pt(12); fo.bold=True; fo.color.rgb=RGBColor.from_string('028090')
        el2=copy.deepcopy(box_src); tree.append(el2); c=Shape(el2,None)
        c.left,c.top,c.width,c.height=Emu(COL[col]['c']),Emu(ry),Emu(2286000),Emu(548640)
    if links:
        b=new.shapes.add_textbox(Emu(548640),Emu(5852160),Emu(8046720),Emu(457200))
        f=b.text_frame; f.word_wrap=True
        for m in ('lIns','rIns','tIns','bIns'): f._bodyPr.set(m,'0')
        for i,(lab,txt,url) in enumerate(links[:2]):
            p=f.paragraphs[0] if i==0 else f.add_paragraph()
            r1=p.add_run(); r1.text=lab
            fo=r1.font; fo.name='Arial'; fo.size=Pt(11); fo.bold=True; fo.color.rgb=RGBColor.from_string('111111')
            r2=p.add_run(); r2.text=txt or url
            fo=r2.font; fo.name='Arial'; fo.size=Pt(11); fo.underline=True; fo.color.rgb=RGBColor.from_string('028090')
            r2.hyperlink.address=url
            hl=r2._r.find(A+'rPr').find(A+'hlinkClick')
            ext=hl.makeelement(A+'extLst',{}); e=hl.makeelement(A+'ext',{'uri':'{A12FA001-AC4F-418D-AE19-62706E023703}'})
            e.append(e.makeelement('{%s}hlinkClr'%AHYP,{'val':'tx'})); ext.append(e); hl.append(ext)
    lst=prs.slides._sldIdLst; items=list(lst); target=None
    for i,s in enumerate(prs.slides):
        t=" ".join(sh.text_frame.text for sh in s.shapes if sh.has_text_frame).strip()
        if t.startswith('Day 3 of 3') or re.match(r'^Day \d of \d', t) and 'stronger' in t.lower(): target=i
    if target is None:
        for i,s in enumerate(prs.slides):
            t=" ".join(sh.text_frame.text for sh in s.shapes if sh.has_text_frame).strip()
            if re.match(r'^Day \d of \d',t): target=i
    el=items[-1]; lst.remove(el); lst.insert(target,el)
    prs.save(outpath)
    return terms,'ok'

if __name__=='__main__':
    base=os.path.expanduser('~/mnt/deck_work/exports')
    want=sys.argv[1:]
    for p in sorted(glob.glob(os.path.join(base,'*/*.pptx'))):
        b=os.path.basename(p)
        if 'with Concept Bank' in b or 'Cycle01_Day 1' in b or 'Cycle 01 —' in b: continue
        m=re.search(r'Cycle (\d{2}[a-d]?)',b)
        if not m: continue
        key=m.group(1)
        if want and key not in want: continue
        out=p.replace('.pptx',' — with Concept Bank.pptx')
        try:
            terms,st=build(p,key,out)
        except Exception as e:
            print(f"{key:4} ERROR {type(e).__name__}: {e}"); continue
        if terms is None: print(f"{key:4} SKIP  {st}")
        else: print(f"{key:4} {len(terms):2d}  " + " · ".join(terms))

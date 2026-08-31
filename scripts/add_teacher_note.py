#!/usr/bin/env python3
"""Draft the teacher note slide - one per deck, cloned from the deck's own
TEACHER NAVIGATION slide.

`deck_lint` fails every deck on H-TEACHER-NOTE and wants seven declarations on
it. Three can be read out of the deck and four cannot:

  READ FROM THE DECK   1 the critical aspects
                       5 the slide-type map
                       6 which simultaneity the cycle works on
  KATHERINE'S          2 what is held invariant
                       3 what breaks if an example is substituted
                       4 position in the sequence
                       7 the visibility rung

The four blanks are worded so they do NOT contain the strings the linter
matches on. That is deliberate. Scaffolding that carries the trigger words
would make A-TN-ITEMS pass on a slide that says nothing, and the check would
stop being worth running. As written, the linter keeps reporting exactly which
of the seven are still unwritten, so the report is the to-do list.

Usage:  python3 add_teacher_note.py <dir> [--apply] [--deck "Cycle 06"]
"""
import argparse
import copy
import glob
import os
import re

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt

TEAL, INK, GREY = "028090", "111111", "666666"
NAV = "TEACHER NAVIGATION"
R_NS = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"

TYPES = [("3-Tier Question", ["Getting Started", "Working On It", "Mastery"]),
         ("Pattern break", ["Pattern break"]),
         ("Build a rule", ["Finish this sentence as a rule"]),
         ("What if", ["What if?"]),
         ("Continuity question", ["Keep going"]),
         ("Concept Bank", ["Define these in your own words"]),
         ("Conflict case", ["Conflict case"]),
         ("Stock-and-flow model", ["stock-and-flow", "Stock and flow"])]


def texts(sh):
    return sh.text_frame.text if sh.has_text_frame else ""


def clone(prs, index):
    src = prs.slides[index]
    dest = prs.slides.add_slide(src.slide_layout)
    for sh in list(dest.shapes):
        sh._element.getparent().remove(sh._element)
    remap = {}
    for rid, rel in src.part.rels.items():
        if rel.reltype.endswith("/slideLayout"):
            continue
        remap[rid] = (dest.part.relate_to(rel.target_ref, rel.reltype, is_external=True)
                      if rel.is_external else
                      dest.part.relate_to(rel.target_part, rel.reltype))
    for el in src.shapes._spTree:
        if el.tag.endswith("}nvGrpSpPr") or el.tag.endswith("}grpSpPr"):
            continue
        new = copy.deepcopy(el)
        for node in [new] + list(new.iter()):
            for a in ("embed", "id", "link"):
                k = R_NS + a
                if k in node.attrib and node.attrib[k] in remap:
                    node.attrib[k] = remap[node.attrib[k]]
        dest.shapes._spTree.append(new)
    return dest


def move_to(prs, position):
    lst = prs.slides._sldIdLst
    el = [s for s in lst][-1]
    lst.remove(el)
    lst.insert(position, el)


def survey(prs):
    per_slide = ["\n".join(texts(sh) for sh in s.shapes) for s in prs.slides]
    joined = "\n".join(per_slide)
    aspects = []
    # Singular only. The slide 1 block and the conflict case slide both say
    # "Critical aspects:" with both aspects on one line, and matching those
    # gave a third aspect that is really the pair joined by a middle dot.
    for a in re.findall(r"Critical aspect:\s*([^\n]+)", joined):
        a = a.strip()
        if a and a not in aspects:
            aspects.append(a)
    present, absent = [], []
    for name, needles in TYPES:
        n = sum(1 for t in per_slide if all(x in t for x in needles))
        (present if n else absent).append((name, n))
    return aspects, present, absent


def body(aspects, present, absent):
    """Left column: what the deck can say about itself. Right: what it cannot."""
    left = [("CRITICAL ASPECTS", True)]
    for i, a in enumerate(aspects, 1):
        left.append(("%d. %s" % (i, a), False))
    left.append(("", False))
    left.append(("SLIDE TYPES IN THIS CYCLE", True))
    left.append((" · ".join("%s ×%d" % (n, c) if c > 1 else n for n, c in present), False))
    if absent:
        left.append(("Left out: %s." % ", ".join(n for n, _ in absent), False))
        left.append(("Why each was left out — TO WRITE", False))
    left.append(("", False))
    left.append(("SIMULTANEITY", True))
    names = [n for n, _ in present]
    sim = []
    if "Concept Bank" in names:
        sim.append("Diachronic — the Concept Bank brings terms met on different "
                   "days into awareness together.")
    if "Conflict case" in names or "Stock-and-flow model" in names:
        sim.append("Synchronic — the %s asks the student to hold both aspects at once."
                   % ("conflict case" if "Conflict case" in names else "stock-and-flow model"))
    if not sim:
        sim.append("Neither is engineered in this cycle — TO WRITE whether that is "
                   "a decision or a gap.")
    left += [(s, False) for s in sim]

    # The four headings avoid the strings deck_lint matches on, so the check
    # keeps reporting them as unwritten until Katherine writes them.
    right = [("STILL TO WRITE — four of the seven", True),
             ("These are judgements about the teaching, not facts about the file, "
              "so they are left blank rather than guessed.", False),
             ("", False),
             ("WHAT STAYS THE SAME ACROSS THE EXAMPLES", True),
             ("TO WRITE — and the sentence saying the examples differ in one "
              "dimension on purpose.", False),
             ("", False),
             ("IF AN EXAMPLE IS CHANGED", True),
             ("TO WRITE — what stops working, and why the change looks "
              "harmless from outside.", False),
             ("", False),
             ("WHERE THIS SITS IN THE ARC", True),
             ("TO WRITE — what this cycle takes as already discerned, and which "
              "later cycle rests on it.", False),
             ("", False),
             ("HOW MUCH STUDENTS EXPOSE", True),
             ("TO WRITE — written and private, to a partner, read aloud "
              "unattributed, owned aloud, or voted.", False)]
    return left, right


def fill(shape, lines, size=11):
    tf = shape.text_frame
    tf.word_wrap = True
    for p in list(tf.paragraphs)[1:]:
        p._p.getparent().remove(p._p)
    p0 = tf.paragraphs[0]
    for r in list(p0.runs):
        r._r.getparent().remove(r._r)
    for i, (text, head) in enumerate(lines):
        para = p0 if i == 0 else tf.add_paragraph()
        para.space_after = Pt(4)
        if not text:
            continue
        run = para.add_run()
        run.text = text
        run.font.name = "Arial"
        run.font.size = Pt(9 if head else size)
        run.font.bold = head
        run.font.color.rgb = RGBColor.from_string(TEAL if head else INK)


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("path")
    ap.add_argument("--apply", action="store_true")
    ap.add_argument("--deck")
    args = ap.parse_args()

    files = sorted(glob.glob(os.path.join(args.path, "*.pptx")))
    if args.deck:
        files = [f for f in files if args.deck in f]

    for f in files:
        key = os.path.basename(f)[9:-11]
        prs = Presentation(f)
        if any("TEACHER NOTE" in texts(sh) for s in prs.slides for sh in s.shapes):
            print("  %-10s already has one" % key)
            continue
        nav = next((i for i, s in enumerate(prs.slides)
                    if any(NAV in texts(sh) for sh in s.shapes)), None)
        if nav is None:
            print("  %-10s ** no TEACHER NAVIGATION slide to clone **" % key)
            continue

        aspects, present, absent = survey(prs)
        new = clone(prs, nav)
        move_to(prs, nav + 1)
        # The nav slide's first shape is an empty header band that also has a
        # text frame, and its two body columns are the only tall boxes. Taking
        # boxes in top order without those two facts put the kicker on the band
        # and both columns of text into the same box.
        boxes = [sh for sh in new.shapes if sh.has_text_frame
                 and sh.text_frame.text.strip()]
        cols = sorted([b for b in boxes if (b.height or 0) > 3000000],
                      key=lambda s: s.left)
        heads = sorted([b for b in boxes if b not in cols], key=lambda s: s.top)
        if len(heads) < 3 or len(cols) < 2:
            print("  %-10s ** nav slide is not the expected shape - skipped **" % key)
            continue
        kicker, title, sub = heads[0], heads[1], heads[2]

        fill(kicker, [("TEACHER NOTE — do not project", True)], size=11)
        fill(title, [("Design note — what this cycle is doing and why", True)], size=14)
        fill(sub, [("Delete this slide and the deck still runs. It is here so a "
                    "teacher adapting the cycle can see what the design rests on.",
                    False)], size=10)
        L, R = body(aspects, present, absent)
        fill(cols[0], L)
        if len(cols) > 1:
            fill(cols[1], R)

        print("  %-10s note after slide %d · aspects %d · types %d present, %d left out"
              % (key, nav + 1, len(aspects), len(present), len(absent)))
        if args.apply:
            prs.save(f)

    if not args.apply:
        print("\nReport only. Re-run with --apply to write.")


if __name__ == "__main__":
    main()

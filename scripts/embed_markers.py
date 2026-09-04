#!/usr/bin/env python3
"""Embed the grading markers the feedback prompts read.

The prompts find each student writing box by a hidden tag in the label above
it - `[[KIND:CYCLE:QUESTION-ID]]` - and pair a first answer with its revision
by the id, never by slide number. Without the tags a class that filled in
every box reads exactly like a class that wrote nothing, so this is the job
that decides whether the pipeline reports the truth.

The scheme is not invented here. Cycles 02 and 03 were marked by hand in
August and this reproduces what they do: NOTES on every first-thinking box,
DRAFT only on a revised answer, OPTIONAL on optional work, BANK on each
Concept Bank term, and one MARKER-INVENTORY line on slide 1.

Markers are 1pt white Arial sitting after the label text, which is how the
hand-marked decks hide them: invisible on a white ground, findable at 400% in
the editor and by any tool reading the file.

Usage:  python3 embed_markers.py <dir> [--apply] [--deck "Cycle 06"]
"""
import argparse
import glob
import os
import re

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt

MARKER_RE = re.compile(r"\[\[[A-Z-]+:[^\]]*\]\]")
BANK_HEAD = "Define these in your own words"

# Slide-type strings from vt-bio-skill section 7. The id slug follows the type
# where the slide has one, so an id says what kind of question it names.
TYPE_SLUG = [
    ("Getting Started", "TIERNOTES"),
    ("Pattern break", "PATTERNBREAK"),
    ("Finish this sentence as a rule", "RULE"),
    ("Keep going", "KEEPGOING"),
]

STOP = {"the", "a", "an", "your", "you", "this", "that", "what", "why", "how",
        "and", "or", "of", "in", "on", "to", "is", "are", "was", "were", "it",
        "one", "two", "for", "with", "from", "answer", "first", "revised",
        "own", "words", "write", "think", "about", "which", "who", "does"}


def slug(text, n=1):
    words = re.findall(r"[A-Za-z]+", text.lower())
    picked = [w for w in words if w not in STOP and len(w) > 2][:n]
    return "".join(picked).upper() or "BOX"


def texts(shape):
    return shape.text_frame.text if shape.has_text_frame else ""


# The writing boxes carry one fill across all 24 decks - 745 of them. The other
# empty filled shapes are teal accent bars, panels, and the drop zones on the
# drag activities. Cycle 17 has eleven drop zones, and counting those would have
# taken its NOTES denominator from 15 to 29: every student's score too low, with
# nothing in the report to say why. Fill is what tells them apart, not size or
# position, so the test is the fill.
WRITING_FILL = {"F2F6F9"}


def is_writing_box(sh):
    """A student writing box: the writing-box fill, with nothing in it."""
    if sh.has_text_frame and sh.text_frame.text.strip():
        return False
    if sh.width is None or sh.width < 1500000 or sh.height is None or sh.height < 500000:
        return False
    try:
        return sh.fill.type == 1 and str(sh.fill.fore_color.rgb) in WRITING_FILL
    except Exception:
        return False


def label_for(box, shapes):
    """The text box immediately above `box` that horizontally overlaps it."""
    best = None
    for sh in shapes:
        if not sh.has_text_frame or not sh.text_frame.text.strip():
            continue
        if sh.top is None or sh.top + (sh.height or 0) > box.top + 50000:
            continue
        if sh.left + (sh.width or 0) < box.left or sh.left > box.left + box.width:
            continue
        if best is None or sh.top > best.top:
            best = sh
    return best


def plan(prs, cyc):
    """Work out every marker this deck should carry. Returns a list of
    (slide_index, shape, kind, qid) plus the bank terms."""
    out, bank, seen = [], [], set()
    aspects, sticky = [], "PRE"
    for i, s in enumerate(prs.slides):
        joined = "\n".join(texts(sh) for sh in s.shapes)
        first = next((texts(sh).split("\n")[0] for sh in s.shapes
                      if texts(sh).strip()), "")

        # Teacher reference slides (slides 1-3) carry filled panels in the
        # writing-box colour - ESSENTIAL CLAIM, STANDARDS, OBJECTIVES. They
        # are not student boxes; counting them added two NOTES to every deck.
        if "TEACHER REFERENCE" in joined.upper():
            continue

        # A slide can name its aspect without carrying the label - Cycle 06's
        # worksheet slide says "CA2 worksheet" and sits before the first
        # Critical aspect line, so without this it inherits BELLRINGER.
        exp = re.search(r"\bCA([12])\b", joined)
        if exp:
            sticky = "CA" + exp.group(1)

        m = re.search(r"Critical aspect[^:\n]*:\s*([^\n]+)", joined)
        if m:
            a = m.group(1).strip()
            if a not in aspects:
                aspects.append(a)
            sticky = "CA%d" % (aspects.index(a) + 1)

        # BELLRINGER, WHATIF, CONFLICT and OPT name one slide each. Only the
        # aspect block carries forward. Cycle 09 has seven lab and station
        # slides after its What if?, and treating that block as sticky gave
        # them WHATIF-2 through WHATIF-8 - ids that say the wrong thing about
        # which question a student was answering.
        block = sticky
        if first.startswith("Bellringer"):
            block = "BELLRINGER"
        elif first.startswith("What if?"):
            block = "WHATIF"
        elif first.startswith("Conflict case"):
            block = "CONFLICT"
        elif first.startswith("Optional challenge") or first.startswith("Relates to me"):
            block = "OPT"

        if BANK_HEAD in joined:
            # The bank is a two-column grid: a narrow term box on the left and
            # an empty box beside it for the student's definition. Terms are
            # the narrow ones; the heading, the instruction line and the
            # readers' note all run the full width.
            cand = [sh for sh in s.shapes
                    if texts(sh).strip() and (sh.width or 0) <= 2000000]
            for sh in sorted(cand, key=lambda x: x.top):
                bank.append((i, sh, texts(sh).strip()))
            continue

        shapes = list(s.shapes)
        for box in sorted([sh for sh in shapes if is_writing_box(sh)],
                          key=lambda x: x.top):
            lab = label_for(box, shapes)
            if lab is None:
                continue
            head = texts(lab).split("\n")[0].strip()

            # A drag-and-drop activity slide is full of filled boxes that are
            # drop zones for movable pictures, not writing boxes. Cycle 17 has
            # eleven of them, and marking those would have taken its NOTES
            # denominator from 15 to 29 - every student's score too low, and
            # nothing in the report to say why. Only a box asking for writing
            # counts, and on those slides it announces itself.
            answer_box = head.startswith(("Your answer", "Your first answer",
                                          "Your revised answer", "What if"))

            kind = "OPTIONAL" if block == "OPT" else "NOTES"
            if head.startswith("Your revised answer"):
                kind = "DRAFT"

            if block == "OPT":
                base = ("OPT-CHALLENGE" if first.startswith("Optional challenge")
                        else "OPT-RELATES")
            elif block == "WHATIF":
                base = "WHATIF"
            elif block == "CONFLICT":
                base = "CONFLICT-WHATIF" if head.startswith("What if") else "CONFLICT"
            elif answer_box:
                # -MASTERY belongs to the two-tier response slide. A lone
                # "Your answer" on an activity slide is a different question
                # and must not collide with it.
                base = block + ("-MASTERY" if "Your revised answer" in joined
                                else "-ANSWER")
            elif "activity" in first.lower():
                base = block + "-ACT-" + slug(head)
            else:
                base = block + "-" + next(
                    (sl for needle, sl in TYPE_SLUG if needle in joined),
                    slug(head) if head else "OPENING")

            qid = base
            n = 2
            while (kind, qid) in seen:
                qid = "%s-%d" % (base, n)
                n += 1
            seen.add((kind, qid))
            out.append((i, lab, kind, qid))
    return out, bank, aspects


def strip_markers(prs):
    """Remove every existing marker, so a re-run cannot leave two schemes in
    one deck. Cycles 02 and 03 were marked by hand in August and their ids
    differ from the ones derived here; a deck holding both would be counted
    twice."""
    n = 0
    for s in prs.slides:
        for sh in list(s.shapes):
            if not sh.has_text_frame:
                continue
            t = sh.text_frame.text
            # The inventory shape goes as a whole. Stripping only the bracket
            # left " NOTES=15 DRAFT=3 ..." behind, and the next run added a
            # second inventory box beside it (seen on Cycle 01a and 01b).
            if "MARKER-INVENTORY" in t or ("NOTES=" in t and "DRAFT=" in t and "_IDS=" in t):
                n += 1
                sh._element.getparent().remove(sh._element)
                continue
            for para in sh.text_frame.paragraphs:
                for run in list(para.runs):
                    if MARKER_RE.search(run.text):
                        n += len(MARKER_RE.findall(run.text))
                        cleaned = MARKER_RE.sub("", run.text).strip()
                        if cleaned:
                            run.text = cleaned
                        else:
                            run._r.getparent().remove(run._r)
    return n


def inventory_line(cyc, marks, bank):
    notes = [q for _, _, k, q in marks if k == "NOTES"]
    draft = [q for _, _, k, q in marks if k == "DRAFT"]
    opt = [q for _, _, k, q in marks if k == "OPTIONAL"]
    return ("[[MARKER-INVENTORY:%s]] NOTES=%d DRAFT=%d OPTIONAL=%d BANK_TERMS=%d "
            "NOTES_IDS=%s DRAFT_IDS=%s OPTIONAL_IDS=%s"
            % (cyc, len(notes), len(draft), len(opt), len(bank),
               ",".join(notes), ",".join(draft), ",".join(opt)))


def write_inventory(prs, line):
    s1 = prs.slides[0]
    for sh in s1.shapes:
        if sh.has_text_frame and "MARKER-INVENTORY" in sh.text_frame.text:
            sh.text_frame.text = ""
            stamp(sh, line)
            return "replaced"
    box = s1.shapes.add_textbox(Emu(0), Emu(0), Emu(914400), Emu(114300))
    box.text_frame.text = ""
    stamp(box, line)
    return "added"


def stamp(shape, text):
    para = shape.text_frame.paragraphs[-1]
    run = para.add_run()
    run.text = "  " + text
    run.font.size = Pt(1)
    run.font.name = "Arial"
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("path")
    ap.add_argument("--apply", action="store_true")
    ap.add_argument("--deck")
    args = ap.parse_args()

    files = sorted(glob.glob(os.path.join(args.path, "*.pptx")))
    if args.deck:
        files = [f for f in files if args.deck in f]

    for f in files:
        m = re.search(r"Cycle\s*0?(\d+)\s*([A-Za-z])?", os.path.basename(f))
        if not m:
            print("skip (no cycle number in name):", f)
            continue
        key = "Cycle %02d%s" % (int(m.group(1)), (m.group(2) or "").lower())
        cyc = "C%02d%s" % (int(m.group(1)), (m.group(2) or "").upper())
        prs = Presentation(f)
        existing = sum(len(MARKER_RE.findall(texts(sh)))
                       for s in prs.slides for sh in s.shapes)
        marks, bank, aspects = plan(prs, cyc)
        notes = [q for _, _, k, q in marks if k == "NOTES"]
        draft = [q for _, _, k, q in marks if k == "DRAFT"]
        opt = [q for _, _, k, q in marks if k == "OPTIONAL"]
        print("%-10s %2d NOTES  %d DRAFT  %d OPTIONAL  %2d bank terms   (%d already present)"
              % (key, len(notes), len(draft), len(opt), len(bank), existing))
        if not args.apply:
            for i, lab, kind, qid in marks:
                print("      s%-3d %-9s %-26s %s"
                      % (i + 1, kind, qid, texts(lab).split("\n")[0][:44]))
            continue

        removed = strip_markers(prs)
        marks, bank, aspects = plan(prs, cyc)
        for i, lab, kind, qid in marks:
            stamp(lab, "[[%s:%s:%s]]" % (kind, cyc, qid))
        for i, sh, term in bank:
            stamp(sh, "[[BANK:%s:%s]]" % (cyc, term))
        where = write_inventory(prs, inventory_line(cyc, marks, bank))
        prs.save(f)
        print("      stripped %d, wrote %d markers + %d bank, inventory %s"
              % (removed, len(marks), len(bank), where))


if __name__ == "__main__":
    main()

import urllib.request, sys, ssl
from pptx import Presentation
UA={"User-Agent":"Mozilla/5.0 (compatible; deck-link-check/1.0)"}
def check(u):
    for method in ("HEAD","GET"):
        try:
            req=urllib.request.Request(u, headers=UA, method=method)
            with urllib.request.urlopen(req, timeout=25) as r:
                return r.status, r.geturl()
        except urllib.error.HTTPError as e:
            if method=="HEAD": continue
            return e.code, u
        except Exception as e:
            if method=="HEAD": continue
            return "ERR: "+type(e).__name__, u
    return "ERR", u

p=Presentation(sys.argv[1] if len(sys.argv)>1 else "out3.pptx")
seen=[]
for i,s in enumerate(p.slides,1):
    for sh in s.shapes:
        if sh.has_text_frame:
            for para in sh.text_frame.paragraphs:
                for r in para.runs:
                    if r.hyperlink and r.hyperlink.address:
                        seen.append((i, r.text.strip()[:42], r.hyperlink.address))
for i,label,u in seen:
    st,final = check(u)
    flag = "OK " if st==200 else "!! "
    extra = f"  -> {final}" if final!=u else ""
    print(f"{flag}slide {i:>2} [{st}] {u}{extra}")

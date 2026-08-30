#!/usr/bin/env python3
"""Slide 1 gets a CRITICAL ASPECTS block listing the deck's critical aspects
verbatim.

Katherine, 2026-08-30: "On slide 1 write Objectives (Critical aspects on each
one)", and, choosing between placements, "Their own block, verbatim."

Nothing here is authored. The aspects are read off the `Critical aspect:`
labels the question slides already carry, in deck order, de-duplicated, written
out exactly as the deck states them. Change a label and re-run; the block
follows.

## Why the block sits at the foot rather than beside OBJECTIVES

Slide 1 is a full page and its geometry is not the same across the 24. Placing
the block in the OBJECTIVES column works on 13 decks and collides on 11 -
Cycles 03 and 12 use an older layout with no OBJECTIVES section at all. A
mixed result is worse than a uniform one, so the block is a full-width strip
below everything else, identical on every deck.

Where the foot is too tight, the script first reclaims slack from slide-1 text
boxes that are taller than their text needs. That changes box heights only;
no wording is touched, and every trim is reported. Trimming stops 0.2in short
of the estimated text height, so an under-estimate still cannot clip anything.

Usage:  python3 critical_aspects_slide1.py <dir> [--apply]
"""
import argparse
import glob
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt

LABEL = "CRITICAL ASPECTS"
INK = "111111"
GREY = "666666"   # the muted gray of section 13.2; 6B6B6B was mine and drifted

GAP = Emu(91440)            # 0.1in above the block
FOOT = Emu(91440)           # 0.1in below it
LABEL_H = Emu(182880)       # 0.2in
LINE_H = Emu(219456)        # 0.24in
CUSHION = Emu(180000)       # never trim closer than this to the text
MIN_SLACK = Emu(250000)     # and never trim a box with less slack than this
MAX_GAP = Emu(228600)       # 0.25in - the widest vertical gap kept in a column


def aspects(prs):
    out = []
    for slide in prs.slides:
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            for line in sh.text_frame.text.split("\n"):
                if "Critical aspect" in line and ":" in line:
                    v = line.split(":", 1)[1].strip()
                    if v and v not in out:
                        out.append(v)
    return out


def est_height(sh):
    """Roughly how tall the text in `sh` needs to be. Deliberately generous."""
    width_pt = (sh.width or 0) / 12700.0
    total = 0.0
    for para in sh.text_frame.paragraphs:
        txt = "".join(r.text for r in para.runs)
        size = next((r.font.size.pt for r in para.runs if r.font.size), None) or 12.0
        per_line = max(1, int(width_pt / (size * 0.50)))
        lines = max(1, -(-len(txt) // per_line)) if txt else 1
        total += lines * size * 1.25
    return int(total * 12700) + 91440


def reclaim(slide):
    """Trim over-tall slide-1 text boxes and restack the columns.

    Trimming alone frees nothing usable, because the slack usually sits in a
    box above the bottom one and shrinking it leaves a hole rather than moving
    anything up. So each column is walked top to bottom: every text box is cut
    to the height its text needs plus a cushion, and every box below it moves
    up by what was cut, keeping the gap it had.

    Slide 1 is the teacher reference page - "not projected to students" - so
    tightening its column spacing costs nothing a class ever sees. Groups and
    other shapes without a text frame keep their height and just move.

    Returns (n_trimmed, emu_freed_at_the_foot).
    """
    cols = {}
    for sh in slide.shapes:
        if sh.left is None or sh.top is None:
            continue
        cols.setdefault(sh.left, []).append(sh)

    n = 0
    before = max(sh.top + (sh.height or 0) for sh in slide.shapes
                 if sh.top is not None)
    for shapes in cols.values():
        shapes.sort(key=lambda s: s.top)
        shift = 0
        prev_bottom = None
        for i, sh in enumerate(shapes):
            gap_cut = 0
            if i >= 2 and prev_bottom is not None:
                gap = sh.top - shift - prev_bottom
                if gap > MAX_GAP:
                    gap_cut = gap - MAX_GAP
            sh.top = sh.top - shift - gap_cut
            shift += gap_cut
            if sh.has_text_frame and sh.height is not None:
                slack = sh.height - est_height(sh)
                if slack > MIN_SLACK:
                    cut = slack - CUSHION
                    sh.height = sh.height - cut
                    shift += cut
                    n += 1
            prev_bottom = sh.top + (sh.height or 0)
    after = max(sh.top + (sh.height or 0) for sh in slide.shapes
                if sh.top is not None)
    return n, before - after


def shrink_bottom_group(slide, deficit):
    """Move the foot of slide 1 up by shrinking its bottom group. Last resort, and only Cycle 03 needs it.

    Cycle 03's slide 1 is built entirely from grouped shapes, so there is no
    text box to trim and no oversized gap to close. Scaling the bottom group
    down by the shortfall moves its foot up and takes its text with it. The
    shortfall is a few hundredths of an inch, so the scale is a percent or
    two; anything past 12 percent is refused rather than made illegible.
    """
    groups = [sh for sh in slide.shapes
              if not sh.has_text_frame and sh.height and sh.top is not None]
    if not groups:
        return 0
    g = max(groups, key=lambda sh: sh.top + sh.height)
    if deficit > g.height * 0.12:
        return 0
    g.height = g.height - int(deficit)
    return int(deficit)


def content_span(slide, slide_width):
    """Left edge, right edge and bottom of the real content on slide 1.

    Two kinds of shape are excluded from the left and right edges, because
    neither is content and both would push the block out to the paper edge:
    the [[MARKER-INVENTORY]] strip in the top corner, and a full-bleed band
    behind the header. They still count towards the bottom.
    """
    xs, bottom = [], 0
    for sh in slide.shapes:
        if sh.left is None or sh.top is None:
            continue
        w = sh.width or 0
        bottom = max(bottom, sh.top + (sh.height or 0))
        if sh.top < 200000 and w < 1500000:
            continue                          # the marker strip
        if w >= slide_width * 0.95:
            continue                          # a full-bleed band
        xs.append((sh.left, sh.left + w))
    if not xs:
        return None
    return min(x[0] for x in xs), max(x[1] for x in xs), bottom


def write_block(slide, left, top, width, items, compact):
    if compact:
        box = slide.shapes.add_textbox(left, top, width, LINE_H)
        para = box.text_frame.paragraphs[0]
        box.text_frame.word_wrap = True
        lab = para.add_run()
        lab.text = LABEL + "   "
        lab.font.size, lab.font.bold, lab.font.name = Pt(10), True, "Arial"
        lab.font.color.rgb = RGBColor.from_string(GREY)
        val = para.add_run()
        val.text = "   ·   ".join("%d. %s" % (i + 1, t) for i, t in enumerate(items))
        val.font.size, val.font.name = Pt(11), "Arial"
        val.font.color.rgb = RGBColor.from_string(INK)
        return

    box = slide.shapes.add_textbox(left, top, width,
                                   LABEL_H + LINE_H * len(items))
    tf = box.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    r = p0.add_run()
    r.text = LABEL
    r.font.size, r.font.bold, r.font.name = Pt(10), True, "Arial"
    r.font.color.rgb = RGBColor.from_string(GREY)
    for i, text in enumerate(items):
        para = tf.add_paragraph()
        run = para.add_run()
        run.text = "%d. %s" % (i + 1, text)
        run.font.size, run.font.name = Pt(12), "Arial"
        run.font.color.rgb = RGBColor.from_string(INK)


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("path")
    ap.add_argument("--apply", action="store_true")
    args = ap.parse_args()

    full = compact = skipped = already = 0
    for f in sorted(glob.glob(os.path.join(args.path, "*.pptx"))):
        key = os.path.basename(f)[9:-11]
        prs = Presentation(f)
        s1 = prs.slides[0]

        if any(sh.has_text_frame and LABEL in sh.text_frame.text
               for sh in s1.shapes):
            print("  %-10s already carries the block" % key)
            already += 1
            continue

        items = aspects(prs)
        if not items:
            print("  %-10s ** no 'Critical aspect:' label anywhere - left alone **" % key)
            skipped += 1
            continue

        need_full = GAP + LABEL_H + LINE_H * len(items) + FOOT
        need_compact = GAP + LINE_H + FOOT
        limit = prs.slide_height

        span = content_span(s1, prs.slide_width)
        note = ""
        if span[2] + need_full > limit:
            n, freed = reclaim(s1)
            span = content_span(s1, prs.slide_width)
            if n:
                note = "  (restacked, %d box%s trimmed, %.2fin freed)" % (
                    n, "" if n == 1 else "es", freed / 914400.0)

        left, right, bottom = span
        if bottom + need_compact > limit:
            got = shrink_bottom_group(s1, bottom + need_compact - limit)
            if got:
                span = content_span(s1, prs.slide_width)
                left, right, bottom = span
                note += "  (bottom group scaled down %.2fin to fit)" % (got / 914400.0)

        if bottom + need_full <= limit:
            mode, height = "full", need_full
        elif bottom + need_compact <= limit:
            mode, height = "compact", need_compact
        else:
            print("  %-10s ** no room even for one line - left alone **%s" % (key, note))
            skipped += 1
            continue

        print("  %-10s %-7s y=%7d  %.2fin spare%s"
              % (key, mode, bottom + GAP, (limit - bottom - height) / 914400.0, note))
        if args.apply:
            write_block(s1, left, bottom + GAP, right - left, items,
                        compact=(mode == "compact"))
            prs.save(f)
        if mode == "full":
            full += 1
        else:
            compact += 1

    print("\n%s: %d full block, %d compact, %d already done, %d left alone"
          % ("APPLIED" if args.apply else "WOULD WRITE",
             full, compact, already, skipped))
    if not args.apply:
        print("Report only. Re-run with --apply to write.")


if __name__ == "__main__":
    main()

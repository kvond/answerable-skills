#!/usr/bin/env python3
"""Audit every built Biology VT deck against the Cycle 02 change set.

Read-only. Reports, per deck, what the standard sweep still has to do:
  removed   slides of the three types Katherine removed (2026-08-29)
  below     question slides sitting below the Day 3 divider (violates 2.3)
  corners   corner slide-number boxes (not in the reference build)
  bank      Concept Bank present, and whether it sits directly above the divider
  bankmark  [[NOTES:...]] markers left on Concept Bank cells (violates 3b)
  bankcap   Concept Bank terms whose first letter is lowercase
  inv       marker-inventory counts on the teacher-reference slide vs actual
  index     a teacher-navigation / slide-index slide is present

Usage:  python3 deck_changeset_audit.py <exports_dir> [out.csv]
"""
import csv
import os
import re
import sys
from collections import OrderedDict

from pptx import Presentation

CORNER_GEO = (164592, 109728, 457200, 274320)
REMOVED_PREFIXES = ("Then and Now", "Think → Write", "Turn your answer into a draft")
QUESTION_PREFIXES = ("What if?", "Optional challenge", "Relates to me", "Critical aspect:",
                     "Pattern break", "Keep going")
MARKER = re.compile(r"\[\[([A-Z-]+):([A-Z0-9]+):([A-Z0-9-]+)\]\]")
INVENTORY = re.compile(r"\[\[MARKER-INVENTORY:[A-Z0-9]+\]\](.*)", re.S)


def first_line(slide):
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            return sh.text_frame.text.strip().split("\n")[0]
    return ""


def all_text(slide):
    t = " \n".join(sh.text_frame.text for sh in slide.shapes if sh.has_text_frame)
    if slide.has_notes_slide:
        t += " \n" + slide.notes_slide.notes_text_frame.text
    return t


def audit(path):
    prs = Presentation(path)
    slides = list(prs.slides)
    row = OrderedDict(deck=os.path.basename(os.path.dirname(path)), slides=len(slides))

    firsts = [first_line(s) for s in slides]
    texts = [all_text(s) for s in slides]

    # -- the three removed slide types (never a day divider; 2.5) ---------------
    removed = [i + 1 for i, f in enumerate(firsts)
               if not re.match(r"^Day \d of \d", f) and f.startswith(REMOVED_PREFIXES)]
    row["removed"] = ";".join(map(str, removed))

    # -- Day 3 divider, and question slides stranded below it (2.3) -------------
    day3 = next((i for i, f in enumerate(firsts) if f.startswith("Day 3 of")), None)
    row["day3"] = day3 + 1 if day3 is not None else ""
    below = []
    if day3 is not None:
        below = [i + 1 for i in range(day3 + 1, len(slides))
                 if firsts[i].startswith(QUESTION_PREFIXES)]
    row["below"] = ";".join(map(str, below))

    # -- corner slide numbers ---------------------------------------------------
    corners = 0
    for s in slides:
        for sh in s.shapes:
            if not sh.has_text_frame:
                continue
            if re.fullmatch(r"\d{1,2}", sh.text_frame.text.strip()) and \
                    (sh.left, sh.top, sh.width, sh.height) == CORNER_GEO:
                corners += 1
    row["corners"] = corners

    # -- Concept Bank -----------------------------------------------------------
    bank = next((i for i, t in enumerate(texts)
                 if "Define these in your own words" in t), None)
    row["bank"] = bank + 1 if bank is not None else ""
    row["bank_above_day3"] = "" if bank is None or day3 is None else \
        ("yes" if bank == day3 - 1 else "NO")

    bankmark, bankcap = 0, []
    if bank is not None:
        for sh in slides[bank].shapes:
            if not sh.has_text_frame:
                continue
            paras = sh.text_frame.paragraphs
            if len(paras) == 2 and paras[1].runs and \
                    paras[1].runs[0].text.strip().startswith("[[NOTES:"):
                bankmark += 1
            if paras and paras[0].runs:
                term = paras[0].runs[0].text.strip()
                if term and term[0].islower() and len(term) < 28:
                    bankcap.append(term)
    row["bankmark"] = bankmark
    row["bankcap"] = ";".join(bankcap)

    # -- marker inventory vs actual --------------------------------------------
    actual = {}
    for t in texts:
        for kind, _cyc, ident in MARKER.findall(t):
            actual.setdefault(kind, OrderedDict())[ident] = True
    declared = ""
    for t in texts:
        m = INVENTORY.search(t)
        if m:
            declared = m.group(1)
            break
    problems = []
    for kind in ("NOTES", "DRAFT", "OPTIONAL"):
        d = re.search(kind + r"=(\d+)", declared)
        a = len(actual.get(kind, {}))
        if d and int(d.group(1)) != a:
            problems.append("%s %s>%d" % (kind, d.group(1), a))
    d = re.search(r"BANK_TERMS=(\d+)", declared)
    nterms = len(bankcap) + sum(
        1 for sh in (slides[bank].shapes if bank is not None else [])
        if sh.has_text_frame and len(sh.text_frame.paragraphs) >= 1
        and sh.text_frame.paragraphs[0].runs
        and sh.text_frame.paragraphs[0].runs[0].text.strip()[:1].isupper()
        and len(sh.text_frame.paragraphs[0].runs[0].text.strip()) < 28
        and sh.text_frame.paragraphs[0].runs[0].text.strip() != "CONCEPT BANK"
    ) if bank is not None else 0
    if d and int(d.group(1)) != nterms:
        problems.append("BANK_TERMS %s>%d" % (d.group(1), nterms))
    row["inv"] = ";".join(problems)

    row["index"] = "yes" if any("TEACHER NAVIGATION" in t or "Slide index" in t
                                for t in texts) else ""
    return row


def main():
    root = sys.argv[1]
    out = sys.argv[2] if len(sys.argv) > 2 else None
    paths = []
    for d in sorted(os.listdir(root)):
        full = os.path.join(root, d)
        if not os.path.isdir(full):
            continue
        for f in sorted(os.listdir(full)):
            if f.endswith("with Concept Bank.pptx"):
                paths.append(os.path.join(full, f))

    rows = []
    for p in paths:
        try:
            rows.append(audit(p))
        except Exception as e:                      # noqa: BLE001 - report, don't stop
            rows.append(OrderedDict(deck=os.path.basename(os.path.dirname(p)),
                                    slides="", removed="ERROR: %s" % e))

    cols = list(rows[0].keys())
    w = {c: max(len(c), max(len(str(r.get(c, ""))) for r in rows)) for c in cols}
    w["deck"] = min(w["deck"], 46)
    print(" ".join(c.ljust(w[c]) for c in cols))
    print(" ".join("-" * w[c] for c in cols))
    for r in rows:
        print(" ".join(str(r.get(c, ""))[:w[c]].ljust(w[c]) for c in cols))

    if out:
        with open(out, "w", newline="") as fh:
            wr = csv.DictWriter(fh, fieldnames=cols)
            wr.writeheader()
            wr.writerows(rows)
        print("\nwrote", out)


if __name__ == "__main__":
    main()

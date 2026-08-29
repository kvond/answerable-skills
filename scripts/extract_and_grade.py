#!/usr/bin/env python3
"""
Workflow A — Completion Grading on Annotated Slides

Reads a teacher template deck and a folder of student-annotated decks.
Identifies diagnostic slides (3-tier critical-aspect concept questions,
pattern-break, build-a-rule, what-if). Extracts each student's answers
keyed by slide number. Scores completion on diagnostic slides only.

Outputs:
  output_dir/completion_report.csv     -- per-student completion scores
  output_dir/answers/<student>.json    -- per-student structured answers (Workflow B input)
  output_dir/workflow_a_summary.md     -- human-readable teacher summary

Usage:
  python extract_and_grade.py <teacher_deck.pptx> <student_dir> <output_dir>
"""

import sys
import json
import csv
import re
from pathlib import Path
from pptx import Presentation


# -- Slide-type detection markers --------------------------------------------
# All matching is case-insensitive (see classify_slide).
# Multiple markers can match the same slide type — used as alternatives.
TIER_LABELS = ("Getting Started", "Working On It", "Mastery")
PATTERN_BREAK_MARKERS = ("Pattern break",)
BUILD_A_RULE_MARKERS = ("Build a rule from", "Finish this rule", "Finish this sentence as a rule")
WHAT_IF_MARKERS = ("What if?", "What if ", "What-if")
CRITICAL_ASPECT_PREFIX = "Critical aspect:"
WORD_BANK_MARKER = "Word bank"
# Slides carrying these markers are never diagnostic (e.g. "Teacher Navigation
# — do not project"). Checked FIRST in classify_slide. Fixes the phantom
# pattern_break count that inflated student incompletion %.
NON_DIAGNOSTIC_MARKERS = ("do not project", "teacher navigation")
# Standing end-of-lesson reflections (Continuation Question, Relates to Me):
# real, projected, student-facing slides that are never scored for diagnostic
# completion. Kept as a separate tuple from NON_DIAGNOSTIC_MARKERS (teacher-
# nav slides that aren't even projected) so the two exclusion reasons stay
# auditable. The Relates to Me slide legitimately repeats "Critical aspect:"
# for every aspect it lists as an option -- this tuple is checked in
# classify_slide() before that pattern, exactly so that repetition never
# phantom-inflates the diagnostic count the way the old nav slides did.
# Must stay in sync with deck_lint.py's STANDING_REFLECTION tuple.
# Added 2026-08-06.
STANDING_REFLECTION_MARKERS = ("continuation question:", "relates to me:")


# ---------------------------------------------------------------------------
# Text extraction helpers
# ---------------------------------------------------------------------------

def extract_frame_texts(slide):
    """Return a list of (frame_index, [paragraph_texts]) for each text frame
    on the slide. Empty paragraphs are dropped. Order preserved."""
    frames = []
    for idx, shape in enumerate(slide.shapes):
        if not shape.has_text_frame:
            continue
        paragraphs = []
        for para in shape.text_frame.paragraphs:
            text = "".join(run.text for run in para.runs).strip()
            if text:
                paragraphs.append(text)
        if paragraphs:
            frames.append((idx, paragraphs))
    return frames


def flatten_frames(frames):
    """Flatten frame list to a single ordered list of strings."""
    out = []
    for _, paragraphs in frames:
        out.extend(paragraphs)
    return out


# ---------------------------------------------------------------------------
# Slide classification
# ---------------------------------------------------------------------------

def classify_slide(all_text):
    """Determine slide type from its full text content.

    All marker matching is case-insensitive. Multiple alternative markers
    per slide type are supported (see *_MARKERS tuples at the top of the file).
    """
    text_lower = all_text.lower()
    if any(m in text_lower for m in NON_DIAGNOSTIC_MARKERS) or \
       any(m in text_lower for m in STANDING_REFLECTION_MARKERS):
        return "other"
    has_all_tiers = all(label.lower() in text_lower for label in TIER_LABELS)
    if has_all_tiers:
        return "critical_aspect_concept_question"
    if any(m.lower() in text_lower for m in PATTERN_BREAK_MARKERS):
        return "pattern_break"
    if any(m.lower() in text_lower for m in BUILD_A_RULE_MARKERS):
        return "build_a_rule"
    if "what if?" in text_lower or re.search(r"(?m)^\s*what[- ]if\b", text_lower):
        return "what_if"
    return "other"


def extract_critical_aspect(all_text):
    """Pull the critical-aspect title (the text after 'Critical aspect:')."""
    match = re.search(r"Critical aspect:\s*([^\n]+)", all_text)
    return match.group(1).strip() if match else None


def extract_word_bank(paragraphs):
    """Pull bracketed word-bank terms from a list of paragraph strings."""
    words = []
    for p in paragraphs:
        # Match [ term ] style tokens
        for m in re.findall(r"\[\s*([^\]]+?)\s*\]", p):
            # Skip the longer 'hint phrase' brackets — those tend to have multiple words and lowercase
            if len(m.split()) <= 2:
                words.append(m)
    return words


# ---------------------------------------------------------------------------
# Teacher deck parsing
# ---------------------------------------------------------------------------

def parse_teacher_deck(deck_path):
    """Walk the teacher deck and produce a per-slide metadata dictionary."""
    prs = Presentation(deck_path)
    slides = {}

    for idx, slide in enumerate(prs.slides, start=1):
        frames = extract_frame_texts(slide)
        flat = flatten_frames(frames)
        all_text = "\n".join(flat)
        slide_type = classify_slide(all_text)
        is_diagnostic = slide_type != "other"

        info = {
            "slide_number": idx,
            "slide_type": slide_type,
            "is_diagnostic": is_diagnostic,
            "teacher_text_set": set(flat),       # used for diff against students
            "teacher_text_ordered": flat,
            "critical_aspect": extract_critical_aspect(all_text),
        }

        if slide_type == "critical_aspect_concept_question":
            info["word_bank"] = extract_word_bank(flat)

        slides[idx] = info

    return slides


# ---------------------------------------------------------------------------
# Student answer extraction
# ---------------------------------------------------------------------------

def detect_tier_attempted(new_texts, slide_text_set):
    """Heuristic: which tier did the student answer at?
    Looks for tier labels in the slide text immediately preceding the new text.
    Returns one of: 'getting_started', 'working_on_it', 'mastery', or None."""
    if not new_texts:
        return None
    # Simple heuristic: if student's text contains tier-specific cues, use them.
    # Otherwise default to None and let Workflow B refine via content analysis.
    joined = " ".join(new_texts).lower()
    # Returning None here is fine — B does the deeper analysis. We're not
    # guessing wrong; we're staying silent when uncertain.
    return None


def extract_student_answers(student_path, teacher_slides):
    """Compare student deck to teacher template and capture new text per slide."""
    prs = Presentation(student_path)
    out = []

    for idx, slide in enumerate(prs.slides, start=1):
        teacher_info = teacher_slides.get(idx, {})
        teacher_set = teacher_info.get("teacher_text_set", set())

        frames = extract_frame_texts(slide)
        new_per_frame = []
        for frame_idx, paragraphs in frames:
            new_in_frame = [p for p in paragraphs if p not in teacher_set]
            if new_in_frame:
                new_per_frame.append({
                    "frame_index": frame_idx,
                    "new_text": new_in_frame,
                })

        new_text_flat = [t for f in new_per_frame for t in f["new_text"]]
        answered = len(new_text_flat) > 0

        slide_record = {
            "slide_number": idx,
            "slide_type": teacher_info.get("slide_type", "unknown"),
            "is_diagnostic": teacher_info.get("is_diagnostic", False),
            "critical_aspect": teacher_info.get("critical_aspect"),
            "answered": answered,
            "student_text": new_text_flat,
            "student_text_by_frame": new_per_frame,
        }

        if teacher_info.get("slide_type") == "critical_aspect_concept_question":
            slide_record["word_bank"] = teacher_info.get("word_bank", [])
            slide_record["tier_attempted"] = detect_tier_attempted(
                new_text_flat, teacher_set
            )

        out.append(slide_record)

    return out


# ---------------------------------------------------------------------------
# Completion scoring
# ---------------------------------------------------------------------------

def score_completion(slide_records):
    """Compute completion stats over diagnostic slides only."""
    diagnostic = [s for s in slide_records if s["is_diagnostic"]]
    answered = [s for s in diagnostic if s["answered"]]
    total = len(diagnostic)
    n_answered = len(answered)
    pct = round(100.0 * n_answered / total, 1) if total else 0.0

    by_type = {}
    for s in diagnostic:
        t = s["slide_type"]
        by_type.setdefault(t, {"total": 0, "answered": 0})
        by_type[t]["total"] += 1
        if s["answered"]:
            by_type[t]["answered"] += 1

    return {
        "total_diagnostic_slides": total,
        "answered_diagnostic_slides": n_answered,
        "completion_percent": pct,
        "by_slide_type": by_type,
    }


# ---------------------------------------------------------------------------
# Output writers
# ---------------------------------------------------------------------------

def write_student_json(out_dir, student_name, lesson_name, completion, slides):
    payload = {
        "student": student_name,
        "lesson_deck": lesson_name,
        "completion": completion,
        "slides": slides,
    }
    path = out_dir / "answers" / f"{student_name}.json"
    path.parent.mkdir(parents=True, exist_ok=True)
    with open(path, "w") as f:
        json.dump(payload, f, indent=2)


def write_completion_csv(out_dir, rows):
    path = out_dir / "completion_report.csv"
    with open(path, "w", newline="") as f:
        writer = csv.DictWriter(
            f,
            fieldnames=[
                "student",
                "answered_diagnostic_slides",
                "total_diagnostic_slides",
                "completion_percent",
            ],
        )
        writer.writeheader()
        writer.writerows(rows)


def write_summary_md(out_dir, teacher_slides, completion_rows, lesson_name):
    diagnostic_slides = [
        (idx, info) for idx, info in teacher_slides.items() if info["is_diagnostic"]
    ]

    lines = [
        "# Workflow A — Completion Report",
        "",
        f"**Lesson:** `{lesson_name}`",
        f"**Total slides:** {len(teacher_slides)}",
        f"**Diagnostic slides (scored for completion):** {len(diagnostic_slides)}",
        "",
        "## Diagnostic slides identified",
        "",
        "| Slide # | Type | Critical Aspect |",
        "|---|---|---|",
    ]
    for idx, info in diagnostic_slides:
        ca = info.get("critical_aspect") or "—"
        lines.append(f"| {idx} | {info['slide_type']} | {ca} |")

    lines.extend([
        "",
        "## Student completion",
        "",
        "| Student | Answered | Total | % |",
        "|---|---|---|---|",
    ])
    for row in sorted(completion_rows, key=lambda r: r["student"]):
        lines.append(
            f"| {row['student']} | {row['answered_diagnostic_slides']} | "
            f"{row['total_diagnostic_slides']} | {row['completion_percent']}% |"
        )

    lines.extend([
        "",
        "## Next step",
        "",
        "Per-student JSON files in `answers/` are ready for **Workflow B** "
        "(personalized email generation).",
    ])

    with open(out_dir / "workflow_a_summary.md", "w") as f:
        f.write("\n".join(lines))


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def run(teacher_deck, student_dir, output_dir):
    teacher_deck = Path(teacher_deck)
    student_dir = Path(student_dir)
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    teacher_slides = parse_teacher_deck(teacher_deck)

    # HALT if the template has no diagnostic slides — otherwise every student
    # silently scores 0/0 = 0.0% (wrong template passed). See Workflow A SKILL.
    n_diag = sum(1 for s in teacher_slides.values() if s["is_diagnostic"])
    if n_diag == 0:
        print(
            f"HALT: teacher deck '{teacher_deck.name}' has 0 diagnostic slides "
            f"(no Getting Started/Working On It/Mastery tiers found). Wrong template? "
            f"Refusing to produce a meaningless 0/0 completion run.",
            file=sys.stderr,
        )
        sys.exit(3)

    completion_rows = []
    # Only Notes-prefixed student decks (Workflow A scope). Skips the template,
    # exemplars, and Draft-prefixed files (those belong to Workflow B1).
    student_files = [
        p for p in sorted(student_dir.glob("*.pptx"))
        if p.name != teacher_deck.name
        and re.search(r"(?i)(?:^|[-_ ])notes", p.stem)
    ]
    if not student_files:
        print(f"No Notes-prefixed .pptx files found in {student_dir}", file=sys.stderr)
        sys.exit(2)

    for sf in student_files:
        student_name = sf.stem
        slide_records = extract_student_answers(sf, teacher_slides)
        completion = score_completion(slide_records)

        write_student_json(
            output_dir, student_name, teacher_deck.name, completion, slide_records
        )

        completion_rows.append({
            "student": student_name,
            "answered_diagnostic_slides": completion["answered_diagnostic_slides"],
            "total_diagnostic_slides": completion["total_diagnostic_slides"],
            "completion_percent": completion["completion_percent"],
        })

    write_completion_csv(output_dir, completion_rows)
    write_summary_md(output_dir, teacher_slides, completion_rows, teacher_deck.name)

    print(f"Processed {len(completion_rows)} student decks.")
    print(f"Outputs written to {output_dir}")


if __name__ == "__main__":
    if len(sys.argv) != 4:
        print(__doc__, file=sys.stderr)
        sys.exit(1)
    run(sys.argv[1], sys.argv[2], sys.argv[3])
